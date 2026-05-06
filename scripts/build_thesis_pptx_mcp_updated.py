# -*- coding: utf-8 -*-
"""
Build thesis presentation with an updated MCP workflow slide.

This wrapper first runs the existing PPTX generator, then appends a slide that
summarizes the latest five-tool MCP workflow:

- initialize_environment
- sample_point
- learn_impacts
- run_window_direct
- rank_actions

Output:
- outputs/papers/thesis_presentation_zh_mcp_updated.pptx
"""

from __future__ import annotations

from pathlib import Path
import shutil
import subprocess
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
ROOT = SCRIPT_DIR.parent
PAPERS = ROOT / "outputs" / "papers"
BASE_PPTX = PAPERS / "thesis_presentation_zh.pptx"
OUTPUT_PPTX = PAPERS / "thesis_presentation_zh_mcp_updated.pptx"

BACKGROUND_COLOR = RGBColor(238, 242, 247)
HEADER_FILL = RGBColor(23, 37, 61)
HEADER_TEXT = RGBColor(255, 255, 255)
HEADER_SUBTITLE = RGBColor(211, 225, 241)
TEXT_COLOR = RGBColor(25, 32, 40)
ACCENT_COLOR = RGBColor(0, 83, 130)
CARD_FILL = RGBColor(255, 255, 255)
CARD_LINE = RGBColor(124, 143, 165)
MUTED_COLOR = RGBColor(70, 80, 92)


def run_base_generator() -> None:
    subprocess.run([sys.executable, str(SCRIPT_DIR / "build_thesis_pptx.py")], cwd=ROOT, check=True)
    if not BASE_PPTX.exists():
        raise FileNotFoundError(f"Base presentation not found: {BASE_PPTX}")


def add_title(slide, title: str, subtitle: str = "") -> None:
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.16))
    header.fill.solid()
    header.fill.fore_color.rgb = HEADER_FILL
    header.line.color.rgb = HEADER_FILL

    title_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.17), Inches(12.2), Inches(0.62))
    frame = title_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "PingFang TC"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = HEADER_TEXT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(12.1), Inches(0.34))
        sub_frame = sub_box.text_frame
        sub_frame.clear()
        sub_p = sub_frame.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle
        sub_run.font.name = "PingFang TC"
        sub_run.font.size = Pt(11)
        sub_run.font.color.rgb = HEADER_SUBTITLE


def add_footer(slide, page: int) -> None:
    box = slide.shapes.add_textbox(Inches(12.25), Inches(7.0), Inches(0.6), Inches(0.25))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(page)
    run.font.name = "Arial"
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED_COLOR


def add_card(slide, left: float, top: float, width: float, height: float, title: str, body: str) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_FILL
    shape.line.color.rgb = CARD_LINE
    shape.line.width = Pt(1.1)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    title_run = p.add_run()
    title_run.text = title
    title_run.font.name = "Consolas"
    title_run.font.size = Pt(14)
    title_run.font.bold = True
    title_run.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(6)

    bp = frame.add_paragraph()
    bp.text = body
    bp.font.name = "PingFang TC"
    bp.font.size = Pt(11)
    bp.font.color.rgb = TEXT_COLOR


def add_mcp_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BACKGROUND_COLOR
    add_title(
        slide,
        "MCP-enabled Interactive Digital Twin",
        "新版 MCP 從 demo-oriented tools 收斂為 workflow-oriented tools",
    )

    cards = [
        (
            0.65,
            1.45,
            "initialize_environment",
            "建立房間、設備、家具、室內 baseline 與外部邊界條件。",
        ),
        (
            4.75,
            1.45,
            "sample_point",
            "查詢任意座標的 temperature / humidity / illuminance；支援 elapsed time 與 steady state。",
        ),
        (
            8.85,
            1.45,
            "learn_impacts",
            "由 before/after 感測資料學習非連網設備影響係數。",
        ),
        (
            2.7,
            3.45,
            "run_window_direct",
            "輸入外部溫度、濕度、日照與開窗比例，模擬窗戶邊界條件影響。",
        ),
        (
            6.8,
            3.45,
            "rank_actions",
            "針對指定點與目標值排序控制動作，回傳 penalty、improvement 與 resulting values。",
        ),
    ]
    for left, top, title, body in cards:
        add_card(slide, left, top, 3.8, 1.45, title, body)

    note = slide.shapes.add_textbox(Inches(0.9), Inches(5.55), Inches(11.6), Inches(0.85))
    frame = note.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = "重點：AI client 不只是執行預設 demo，而是能初始化環境、查詢點位、學習設備影響，並轉化為控制建議。"
    p.font.name = "PingFang TC"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = TEXT_COLOR

    add_footer(slide, len(prs.slides))


def main() -> None:
    PAPERS.mkdir(parents=True, exist_ok=True)
    run_base_generator()
    shutil.copyfile(BASE_PPTX, OUTPUT_PPTX)
    prs = Presentation(str(OUTPUT_PPTX))
    add_mcp_slide(prs)
    prs.save(str(OUTPUT_PPTX))
    print(f"Wrote {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
