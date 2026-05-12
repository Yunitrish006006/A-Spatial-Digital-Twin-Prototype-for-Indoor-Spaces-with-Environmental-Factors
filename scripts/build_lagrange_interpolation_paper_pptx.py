from __future__ import annotations

from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "outputs" / "papers"
DOCS_THESIS = ROOT / "docs" / "thesis"

PPTX_PATH = OUTPUT_DIR / "lagrange_interpolation_paper_report_zh.pptx"
OUTLINE_PATH = DOCS_THESIS / "lagrange_interpolation_paper_report_zh.md"
REPORT_PATH = DOCS_THESIS / "lagrange_interpolation_paper_report_full_zh.md"
SPEAKER_NOTES_PATH = DOCS_THESIS / "lagrange_interpolation_paper_speaker_notes_zh.md"
PDF_SOURCE = ROOT / "docs" / "papers" / "data source" / "s13660-022-02841-w.pdf"

SLIDE_W = 13.333
SLIDE_H = 7.5

BACKGROUND = RGBColor(244, 247, 251)
HEADER = RGBColor(18, 32, 51)
HEADER_SUBTITLE = RGBColor(211, 225, 241)
ACCENT = RGBColor(0, 97, 148)
ACCENT_LIGHT = RGBColor(45, 155, 200)
TEXT = RGBColor(30, 39, 51)
MUTED = RGBColor(86, 97, 112)
CARD = RGBColor(255, 255, 255)
CARD_LINE = RGBColor(209, 219, 231)
SOFT = RGBColor(226, 232, 240)
WARN = RGBColor(176, 77, 50)
GOOD = RGBColor(42, 125, 84)

BODY_FONT = "Noto Sans TC"
LATIN_FONT = "Arial"
FORMULA_FONT = "Cambria Math"


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND
    return slide


def add_shape(slide, shape_type, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line is not None else fill
    return shape


def configure_frame(frame, margin: float = 0.04) -> None:
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: int = 16,
    color: RGBColor = TEXT,
    bold: bool = False,
    font: str = BODY_FONT,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    configure_frame(frame)
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header(slide, title: str, subtitle: str, page: int) -> None:
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, 0.86, HEADER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0.86, SLIDE_W, 0.03, ACCENT_LIGHT)
    add_text(slide, 0.55, 0.12, 11.6, 0.34, title, size=20, color=RGBColor(255, 255, 255), bold=True)
    add_text(slide, 0.55, 0.51, 11.4, 0.24, subtitle, size=10, color=HEADER_SUBTITLE)
    add_text(slide, 12.25, 6.98, 0.55, 0.2, str(page), size=9, color=MUTED, align=PP_ALIGN.RIGHT, font=LATIN_FONT)


def add_card(slide, x: float, y: float, w: float, h: float, title: str | None = None, title_color: RGBColor = ACCENT):
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.03, y + 0.04, w, h, SOFT)
    card = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, CARD, CARD_LINE)
    if title:
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, 0.08, ACCENT_LIGHT)
        add_text(slide, x + 0.18, y + 0.18, w - 0.36, 0.26, title, size=15, color=title_color, bold=True)
    return card


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    items: Sequence[str],
    size: int = 14,
    color: RGBColor = TEXT,
    leading: float = 1.12,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    configure_frame(frame, margin=0.02)
    frame.clear()
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = f"• {item}"
        p.font.name = BODY_FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
        p.line_spacing = leading
    return box


def add_formula(slide, x: float, y: float, w: float, h: float, formula: str, size: int = 22, color: RGBColor = ACCENT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    configure_frame(frame, margin=0.05)
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = formula
    run.font.name = FORMULA_FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return box


def add_connector(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = ACCENT_LIGHT
    line.line.width = Pt(2.2)


def title_slide(prs: Presentation) -> tuple[str, list[str]]:
    slide = new_slide(prs)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, HEADER)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 5.95, SLIDE_W, 0.06, ACCENT_LIGHT)
    add_text(
        slide,
        0.8,
        1.05,
        11.8,
        1.25,
        "論文報告\nLinear Lagrange Interpolation 的誤差常數估計",
        size=32,
        color=RGBColor(255, 255, 255),
        bold=True,
    )
    add_text(
        slide,
        0.85,
        2.85,
        11.4,
        0.9,
        "Galindo, Ike, and Liu (2022)\nError-constant estimation under the maximum norm for linear Lagrange interpolation",
        size=20,
        color=HEADER_SUBTITLE,
    )
    add_text(
        slide,
        0.85,
        4.65,
        10.9,
        0.5,
        "報告重點：研究問題、數學背景、方法設計、定理與數值結果。",
        size=17,
        color=RGBColor(255, 255, 255),
    )
    add_text(slide, 0.85, 6.55, 11.6, 0.28, f"PDF: {PDF_SOURCE.relative_to(ROOT)}", size=10, color=HEADER_SUBTITLE, font=LATIN_FONT)
    return (
        "封面",
        [
            "這次先專門報告 Galindo、Ike、Liu 這篇論文本身。",
            "報告主軸是 interpolation error constant，不延伸到自己的 thesis 應用。",
        ],
    )


def add_step_flow(slide, steps: Sequence[tuple[str, str]], y: float) -> None:
    x = 0.7
    w = 2.35
    for idx, (title, body) in enumerate(steps):
        add_card(slide, x, y, w, 1.32, title)
        add_text(slide, x + 0.16, y + 0.58, w - 0.32, 0.52, body, size=11)
        if idx < len(steps) - 1:
            add_connector(slide, x + w + 0.08, y + 0.66, x + w + 0.48, y + 0.66)
        x += w + 0.55


def add_simple_table(
    slide,
    x: float,
    y: float,
    widths: Sequence[float],
    row_h: float,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    font_size: int = 12,
) -> None:
    total_w = sum(widths)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, total_w, row_h, ACCENT, ACCENT)
    cx = x
    for width, header in zip(widths, headers):
        add_text(slide, cx + 0.04, y + 0.12, width - 0.08, row_h - 0.14, header, size=font_size, color=RGBColor(255, 255, 255), bold=True, align=PP_ALIGN.CENTER)
        cx += width
    for ridx, row in enumerate(rows):
        ry = y + row_h * (ridx + 1)
        fill = CARD if ridx % 2 == 0 else RGBColor(248, 251, 254)
        cx = x
        for width, value in zip(widths, row):
            add_shape(slide, MSO_SHAPE.RECTANGLE, cx, ry, width, row_h, fill, CARD_LINE)
            add_text(slide, cx + 0.06, ry + 0.12, width - 0.12, row_h - 0.16, value, size=font_size, align=PP_ALIGN.CENTER)
            cx += width


def build_slides(prs: Presentation) -> list[tuple[str, list[str]]]:
    notes: list[tuple[str, list[str]]] = [title_slide(prs)]

    page = 2
    slide = new_slide(prs)
    add_header(slide, "論文基本資料", "Journal of Inequalities and Applications, 2022", page)
    add_card(slide, 0.72, 1.22, 12.0, 5.0, "Citation")
    add_text(
        slide,
        1.05,
        1.86,
        11.35,
        0.78,
        "Shirley Mae Galindo, Koichiro Ike, and Xuefeng Liu.\nError-constant estimation under the maximum norm for linear Lagrange interpolation.",
        size=19,
        bold=True,
    )
    add_bullets(
        slide,
        1.08,
        3.0,
        11.0,
        2.35,
        [
            "期刊：Journal of Inequalities and Applications",
            "年份：2022，文章編號 109",
            "DOI：10.1186/s13660-022-02841-w",
            "關鍵詞：Lagrange interpolation、finite-element method、Fujino-Morley interpolation、Bernstein polynomial",
        ],
        size=15,
    )
    add_text(slide, 0.85, 6.36, 11.7, 0.26, "一句話定位：這是一篇數值分析論文，目標是嚴格估計線性 Lagrange 補間在最大範數下的誤差常數。", size=12, color=MUTED)
    notes.append(("論文基本資料", ["先交代作者、題目、期刊、年份與 DOI。", "定位要明確：這篇是數值分析與 FEM 補間誤差常數論文。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "報告地圖：這篇論文怎麼讀？", "先理解補間問題，再看常數估計方法", page)
    add_card(slide, 0.65, 1.14, 12.05, 5.12, "Paper reading map")
    add_step_flow(
        slide,
        [
            ("Intro", "補間誤差背景與既有 constant estimates"),
            ("Sect. 2", "一般三角形 Cᴸ(K) 的 raw upper bound"),
            ("Sect. 3", "FEM + Bernstein 求具體三角形的 sharp bounds"),
            ("Sect. 4-5", "數值結果、interval arithmetic 與結論"),
        ],
        1.85,
    )
    add_bullets(
        slide,
        1.05,
        4.32,
        11.1,
        1.05,
        [
            "報告時先把「補間」與「最大誤差常數」講清楚，後面定理才不會抽象",
            "這篇的主貢獻不只是給公式，而是把常數估計變成可驗證的數值分析流程",
        ],
        size=13,
    )
    notes.append(("報告地圖", ["這頁先告訴聽眾整篇 paper 的閱讀順序。", "報告主線是：問題定義、一般上界、銳利估計演算法、數值驗證。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "研究背景：什麼是 Lagrange interpolation？", "用節點值建立通過節點的多項式近似", page)
    add_card(slide, 0.65, 1.18, 5.85, 5.05, "1D interval")
    add_formula(slide, 1.0, 1.86, 5.15, 0.58, "Πᴸu(0)=u(0),  Πᴸu(1)=u(1)", size=22)
    add_bullets(slide, 0.95, 2.82, 5.25, 2.85, ["在區間端點知道函數值", "用一條直線近似中間的函數", "補間函數在端點完全等於原函數", "中間誤差由函數彎曲程度決定"], size=14)
    add_card(slide, 6.85, 1.18, 5.85, 5.05, "2D triangle")
    add_formula(slide, 7.2, 1.86, 5.15, 0.58, "(u - Πᴸu)(pᵢ)=0,  i=1,2,3", size=22)
    add_bullets(slide, 7.15, 2.82, 5.25, 2.85, ["在三角形三個頂點知道函數值", "用一個線性平面近似三角形內部", "是 finite element method 中常見的 P1 元素", "本論文研究它的最大誤差常數"], size=14)
    notes.append(("Lagrange interpolation", ["先用一維端點連線說明，再推到三角形三頂點的線性平面。", "聽眾只要知道：補間在節點上完全正確，但內部會有誤差。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "研究問題：補間誤差可以怎麼界定？", "maximum norm 看的是整個區域內最壞的點誤差", page)
    add_card(slide, 0.65, 1.15, 12.05, 5.15, "Error estimate")
    add_formula(slide, 1.0, 1.82, 11.3, 0.65, "||u - Πᴸu||∞,K ≤ Cᴸ(K) |u|₂,K", size=27)
    add_bullets(
        slide,
        1.08,
        3.0,
        11.05,
        2.42,
        [
            "左邊：補間函數和原函數在三角形 K 內的最大點誤差",
            "右邊：二階 Sobolev seminorm 代表函數的二階變化量",
            "Cᴸ(K)：本篇要估計的 interpolation error constant",
            "核心問題：對不同形狀的三角形，Cᴸ(K) 到底多大？",
        ],
        size=15,
    )
    add_text(slide, 0.88, 6.35, 11.7, 0.26, "這類誤差界的用途是把「補間會錯多少」變成可計算、可比較、可驗證的常數估計問題。", size=12, color=MUTED)
    notes.append(("研究問題", ["這頁是整篇論文的主問題。", "它要找的不是某個函數的單次誤差，而是三角形元素 K 對所有足夠平滑函數的最壞誤差常數。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "符號表：先把公式讀懂", "這些符號會反覆出現在定理與演算法中", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "Notation")
    add_simple_table(
        slide,
        0.95,
        1.78,
        [1.7, 5.05, 4.45],
        0.52,
        ["符號", "意義", "報告時怎麼說"],
        [
            ["K", "三角形元素 domain", "補間發生的區域"],
            ["Πᴸu", "linear Lagrange interpolation", "通過頂點值的線性近似"],
            ["||·||∞,K", "K 內 maximum norm", "看最壞點誤差"],
            ["|u|₂,K", "H² seminorm", "衡量二階變化量"],
            ["Cᴸ(K)", "interpolation error constant", "本文要估的主角"],
            ["λ(K)", "Cᴸ(K) 對應的最小化量", "先估 λ，再換回 Cᴸ"],
        ],
        font_size=11,
    )
    add_text(slide, 0.98, 5.55, 11.2, 0.34, "核心公式可以讀成：最大點誤差 ≤ 幾何常數 × 函數二階變化量。", size=13, color=GOOD, bold=True)
    notes.append(("符號表", ["這頁讓聽眾先理解符號。", "可以用一句話讀公式：最大點誤差被一個幾何相關常數和函數二階變化量控制。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "已知的一維最佳誤差界", "作者在 introduction 中先回顧一般 Lagrange interpolation 文獻", page)
    add_card(slide, 0.65, 1.18, 12.05, 4.95, "1D well-known optimal estimate")
    add_formula(slide, 1.05, 1.88, 11.25, 0.7, "||u - Πᴸu||∞,I ≤ (1/8)||u⁽²⁾||∞,I    for I=(0,1)", size=23)
    add_formula(slide, 1.05, 2.92, 11.25, 0.7, "一般長度 h：  ||u - Πᴸu||∞ ≤ (h²/8) sup |u''|", size=23)
    add_bullets(
        slide,
        1.08,
        4.08,
        11.1,
        1.26,
        [
            "這個 1/8 常數是 optimal estimate，不是任意選的係數",
            "文章以 u(x)=x² 說明 L∞ 誤差界可達等號",
            "直觀上，線性補間只能捕捉直線，不能捕捉二階彎曲",
        ],
        size=14,
    )
    notes.append(("1D 結果", ["這頁是數學背景，不是這篇的主要新貢獻。", "但它幫聽眾建立直覺：補間跨度越長、二階導數越大，最大誤差越大。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "示範 1：一維 u(x)=x² 真的達到 1/8 常數", "用最簡單函數把 maximum error 算出來", page)
    add_card(slide, 0.65, 1.13, 12.05, 5.15, "I=(0,1),  u(x)=x²")
    add_formula(slide, 0.95, 1.78, 11.45, 0.55, "u(0)=0,  u(1)=1  ⇒  Πᴸu(x)=x", size=23)
    add_formula(slide, 0.95, 2.58, 11.45, 0.55, "error = |u-Πᴸu| = |x²-x| = x(1-x)", size=23)
    add_formula(slide, 0.95, 3.38, 11.45, 0.55, "max at x=0.5：  |0.25-0.5| = 0.25", size=23)
    add_formula(slide, 0.95, 4.18, 11.45, 0.55, "(1/8)||u''||∞ = (1/8)×2 = 0.25", size=23)
    add_bullets(
        slide,
        1.05,
        5.18,
        11.0,
        0.65,
        [
            "所以這個例子剛好達到等號，說明 1/8 是最佳常數，而不是經驗係數",
            "若區間長度改成 h，誤差會按 h² 放大：h=2 時最大誤差變成 1.0",
        ],
        size=12,
    )
    notes.append(("一維示範", ["用 u=x² 示範 1/8 常數怎麼出現。", "端點線性補間是 x，中間最大誤差在 x=0.5，剛好等於右邊的 1/8 乘上二階導數上界。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "示範 1b：區間長度 h 變大，誤差按 h² 放大", "同樣用 u(x)=x² 的直覺看 scaling", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "Scaling example")
    add_simple_table(
        slide,
        1.05,
        1.82,
        [2.0, 2.8, 3.0, 3.0],
        0.62,
        ["h", "max error", "(h²/8) sup|u''|", "意義"],
        [
            ["0.5", "0.0625", "(0.25/8)×2=0.0625", "跨度變短，誤差縮小"],
            ["1.0", "0.25", "(1/8)×2=0.25", "標準區間"],
            ["2.0", "1.0", "(4/8)×2=1.0", "跨度加倍，誤差變 4 倍"],
        ],
        font_size=12,
    )
    add_bullets(
        slide,
        1.05,
        4.55,
        11.1,
        1.0,
        [
            "這個 scaling 是後面看三角形大小與形狀時的重要直覺",
            "補間點距離越遠、函數越彎，中間位置的最壞誤差越大",
        ],
        size=13,
    )
    notes.append(("區間 scaling 示範", ["這頁補充 h² scaling。", "可以用 h=2 時誤差變 4 倍來讓聽眾記住：補間跨度很重要。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "既有文獻與缺口", "2D 三角形上的 maximum norm constant 仍需要更明確估計", page)
    add_card(slide, 0.65, 1.15, 3.85, 5.1, "已有結果")
    add_bullets(slide, 0.9, 1.82, 3.38, 3.85, ["1D interpolation 的 L²、H¹、L∞ optimal estimates", "三角形上的 L² / H¹ error constants", "Waldron 對 simplex 頂點線性補間的 L∞ sharp inequality", "D'Azevedo、Shewchuk、Cao 等討論三角形形狀與補間誤差"], size=12)
    add_card(slide, 4.75, 1.15, 3.85, 5.1, "本文切入點")
    add_bullets(slide, 5.0, 1.82, 3.38, 3.85, ["考慮三角形 K 上的 linear Lagrange interpolation", "目標範數是 L∞ maximum norm", "右側使用 H² seminorm", "希望得到 Cᴸ(K) 的明確上界與銳利數值估計"], size=12)
    add_card(slide, 8.85, 1.15, 3.85, 5.1, "難點")
    add_bullets(slide, 9.1, 1.82, 3.38, 3.85, ["L∞ constraint 不容易直接最佳化", "三角形形狀會影響 error constant", "退化三角形會讓常數變大", "需要嚴格 lower / upper bound，而不是只跑浮點近似"], size=12)
    notes.append(("文獻缺口", ["這頁說明這篇論文站在哪些既有研究之後。", "它的缺口是三角形元素上 maximum norm error constant 的明確估計。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Waldron 結果與本文差異", "這篇論文不是從零開始，而是在既有 L∞ interpolation theory 上推進", page)
    add_card(slide, 0.65, 1.14, 5.85, 5.1, "Waldron (1998)")
    add_formula(slide, 0.95, 1.82, 5.25, 0.55, "||u-Πᴸu||∞,K ≤ 1/2(R²-d²)||u⁽²⁾||∞,K", size=15)
    add_bullets(
        slide,
        0.95,
        2.72,
        5.25,
        2.8,
        [
            "使用二階 directional derivative 的 L∞ bound",
            "幾何量包含外接圓半徑 R 與距離 d",
            "是 simplex vertex interpolation 的 sharp inequality",
        ],
        size=13,
    )
    add_card(slide, 6.85, 1.14, 5.85, 5.1, "Galindo et al. (2022)")
    add_formula(slide, 7.15, 1.82, 5.25, 0.55, "||u-Πᴸu||∞,K ≤ Cᴸ(K)|u|₂,K", size=17)
    add_bullets(
        slide,
        7.15,
        2.72,
        5.25,
        2.8,
        [
            "右側使用 H² seminorm，而不是 W²,∞ 型條件",
            "目標是明確估計 Cᴸ(K)",
            "對具體三角形建立 FEM-based sharp bound algorithm",
        ],
        size=13,
    )
    notes.append(("Waldron 比較", ["這頁說明 introduction 中 Waldron 結果和本文主問題的差異。", "Galindo et al. 的式子用 H² seminorm，所以作者說它在某種意義上比 Waldron 的 W²,∞ 設定更 general。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "三角形參數化", "K_{α,θ,h} 用三個幾何參數描述一般三角形", page)
    add_card(slide, 0.65, 1.15, 12.05, 5.05, "Triangle configuration")
    add_formula(slide, 1.05, 1.85, 11.25, 0.65, "p₁=(0,0),  p₂=(h,0),  p₃=(αh cosθ, αh sinθ)", size=24)
    add_bullets(
        slide,
        1.08,
        3.0,
        11.05,
        2.45,
        [
            "h 是基準邊長，αh 是另一條邊長，θ 是夾角",
            "當 h=1 時，記為 K_{α,θ}",
            "Cᴸ(α,θ,h) 可由縮放性質化成 hCᴸ(α,θ,1)",
            "這個參數化讓作者能系統性討論不同形狀三角形的誤差常數",
        ],
        size=15,
    )
    notes.append(("三角形參數化", ["這頁交代作者如何把一般三角形整理成可分析的幾何參數。", "後面 Theorem 2.1 的上界會依賴 α 和 θ。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Theorem 2.1：一般三角形的上界", "由 affine transformation 從右等腰三角形推到一般三角形", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.1, "Upper bound for Cᴸ(α,θ)")
    add_formula(slide, 0.95, 1.82, 11.45, 0.7, "Cᴸ(α,θ) ≤ [v₊(α,θ)/(√2 α sinθ)] Cᴸ(1,π/2)", size=23)
    add_formula(slide, 0.95, 2.8, 11.45, 0.62, "v₊(α,θ)=1+α²+√(1+2α²cos2θ+α⁴)", size=21)
    add_bullets(
        slide,
        1.08,
        3.78,
        11.1,
        1.55,
        [
            "證明核心：把一般三角形用 affine transformation 轉回 reference triangle",
            "L∞ norm 在座標轉換下保持對應，H² seminorm 會受幾何形狀影響",
            "因此三角形越扁或角度越差，誤差常數上界越可能變大",
        ],
        size=14,
    )
    notes.append(("Theorem 2.1", ["這頁講這篇第一個主要理論結果。", "重點不是背公式，而是說明一般三角形的補間誤差常數可由 reference triangle 和幾何變形量控制。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Theorem 2.1 證明流程", "報告時不用逐行證明，但要知道每一步在做什麼", page)
    add_card(slide, 0.65, 1.13, 12.05, 5.16, "Proof map")
    add_simple_table(
        slide,
        0.95,
        1.78,
        [1.45, 4.0, 5.65],
        0.58,
        ["步驟", "做法", "意義"],
        [
            ["1", "定義 affine map", "把一般 K_{α,θ} 對應到 reference triangle K_{1,π/2}"],
            ["2", "比較 L∞ norm", "座標轉換後最大函數值保持對應"],
            ["3", "比較 H² seminorm", "二階導數會被幾何形變放大或縮小"],
            ["4", "套用 reference constant", "得到 Cᴸ(α,θ) 的 general upper bound"],
            ["5", "分析 shape regularity", "說明退化三角形會讓常數上界變差"],
        ],
        font_size=11,
    )
    add_text(slide, 1.0, 5.45, 11.1, 0.34, "一句話：Theorem 2.1 把「幾何形狀」轉成「常數上界中的 multiplier」。", size=13, color=GOOD, bold=True)
    notes.append(("Theorem 2.1 證明流程", ["這頁是定理證明的報告版。", "不用逐行推導偏導數，只要講清楚 affine transformation 如何把一般三角形的幾何影響帶進常數上界。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "示範 2：Theorem 2.1 代入角度看幾何影響", "同樣 α=1，只改變 θ，raw upper bound 會明顯變化", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "使用 Cᴸ(1,π/2) ≤ 0.41596 作為 reference")
    add_simple_table(
        slide,
        1.05,
        1.85,
        [1.55, 1.75, 2.65, 2.65, 2.15],
        0.56,
        ["θ", "sinθ", "multiplier", "raw bound", "直覺"],
        [
            ["60°", "0.866", "2.449", "≤ 1.019", "形狀較好"],
            ["90°", "1.000", "1.414", "≤ 0.588", "reference 附近"],
            ["150°", "0.500", "5.278", "≤ 2.195", "很鈍、較差"],
        ],
        font_size=12,
    )
    add_bullets(
        slide,
        1.05,
        4.55,
        11.15,
        1.05,
        [
            "這個 bound 是 general upper bound，所以偏保守；它的用途是說明幾何退化會讓常數上界變差",
            "後面 Section 3 的 FEM + Bernstein 方法，會對具體三角形給出更銳利的數值上下界",
        ],
        size=13,
    )
    notes.append(("Theorem 2.1 數字示範", ["這頁把 Theorem 2.1 代入三個角度。", "要強調 raw bound 偏保守，但它清楚展示三角形變鈍或變扁時，補間誤差常數上界會變大。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Shape regularity 與退化三角形", "幾何品質會直接影響 interpolation error constant", page)
    add_card(slide, 0.65, 1.18, 5.85, 5.05, "Shape-regular triangles")
    add_bullets(slide, 0.95, 1.85, 5.25, 3.75, ["若最小內角有正下界 δ", "三角形不會無限變扁", "Cᴸ(α,θ) 可保持 bounded", "這是 FEM mesh quality 中常見的穩定性概念"], size=14)
    add_card(slide, 6.85, 1.18, 5.85, 5.05, "Degenerate triangles")
    add_bullets(slide, 7.15, 1.85, 5.25, 3.75, ["若三角形趨近一條線段", "面積趨近 0，形狀條件惡化", "Cᴸ(K) 可能趨近無限大", "因此補間誤差不只由函數平滑度決定，也受幾何支配"], size=14)
    notes.append(("Shape regularity", ["這頁用比較直覺的方式說明幾何品質的重要性。", "對 FEM 或網格方法來說，退化三角形會造成補間和數值解的不穩定。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "最佳常數估計：轉成 optimization problem", "Section 3 進入具體三角形的銳利數值估計", page)
    add_card(slide, 0.65, 1.14, 12.05, 5.1, "Optimization view")
    add_formula(slide, 1.0, 1.82, 11.3, 0.68, "Cᴸ(K)= sup  ||u-Πᴸu||∞,K / |u|₂,K", size=25)
    add_formula(slide, 1.0, 2.78, 11.3, 0.68, "λ(K)= inf |u|²₂,K / ||u||²∞,K,    Cᴸ(K)=1/√λ(K)", size=23)
    add_bullets(
        slide,
        1.08,
        3.85,
        11.05,
        1.38,
        [
            "把誤差常數估計改寫為最小化問題",
            "要取得 Cᴸ(K) 的上界，就要取得 λ(K) 的下界",
            "困難點仍是 L∞ norm 出現在 constraint 中，不容易直接處理",
        ],
        size=14,
    )
    notes.append(("Optimization problem", ["這頁不用推導太細，但要講清楚：作者把常數估計轉成一個可計算的最佳化問題。", "L∞ constraint 是後面引入 Bernstein polynomial 的原因。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Optimization 轉換細節", "為什麼可以從 Cᴸ(K) 換成 λ(K)？", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "From supremum to infimum")
    add_formula(slide, 0.95, 1.75, 11.45, 0.55, "Cᴸ(K)= sup ||u||∞,K / |u|₂,K   with Πᴸu=0", size=22)
    add_formula(slide, 0.95, 2.52, 11.45, 0.55, "λ(K)= inf |u|²₂,K / ||u||²∞,K", size=22)
    add_formula(slide, 0.95, 3.29, 11.45, 0.55, "⇒  Cᴸ(K) = 1 / √λ(K)", size=24)
    add_bullets(
        slide,
        1.05,
        4.25,
        11.1,
        1.05,
        [
            "Πᴸu=0 的意思是只看在三角形三個頂點為 0 的誤差函數空間",
            "如果能證明 λ(K) 不小於某數，Cᴸ(K) 就不會大於對應的 1/√λ",
        ],
        size=13,
    )
    notes.append(("Optimization 轉換細節", ["這頁補上 Cᴸ 和 λ 的關係。", "報告時可說：最大化比值和最小化倒數平方是同一件事，這讓作者可以改做 eigenvalue/optimization 型問題。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Fujino-Morley FEM space 的角色", "用有限維空間逼近無限維 optimization problem", page)
    add_card(slide, 0.65, 1.16, 12.05, 5.08, "Discretization idea")
    add_bullets(
        slide,
        1.0,
        1.85,
        11.1,
        3.7,
        [
            "原本 optimization 是在 H²(K) 這種無限維函數空間上進行",
            "作者使用 Fujino-Morley finite element space 建立可計算的離散問題",
            "利用空間分解與正交性，把連續問題和離散問題的誤差連接起來",
            "Theorem 3.1 給出由離散解推回 λ(K) lower bound 的方法",
        ],
        size=15,
    )
    add_text(slide, 0.92, 6.28, 11.6, 0.28, "簡報講法：Fujino-Morley space 是把理論上的函數最佳化問題變成有限元素矩陣問題的橋梁。", size=12, color=MUTED)
    notes.append(("Fujino-Morley", ["這頁講方法層的第一個工具。", "不需要介紹完整 element 定義，重點是它把無限維問題離散化，並保留可證明的上下界關係。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Fujino-Morley 在證明中的功能", "不只是數值離散化，還用來建立上下界關係", page)
    add_card(slide, 0.65, 1.15, 12.05, 5.1, "Why this FEM space matters")
    add_bullets(
        slide,
        1.0,
        1.82,
        11.1,
        3.65,
        [
            "作者把連續函數空間拆成可處理的 discrete part 與 remainder part",
            "Fujino-Morley interpolation 提供對 remainder 的誤差控制",
            "這讓離散問題算出的 λ_h 可以轉回原始連續 λ(K) 的 lower bound",
            "換句話說，FEM 在這裡不是單純近似，而是嚴格 bound construction 的一部分",
        ],
        size=15,
    )
    add_text(slide, 0.98, 5.72, 11.2, 0.36, "報告簡化說法：Fujino-Morley 負責把「算得到的離散結果」和「原本連續問題的嚴格界」接起來。", size=12, color=GOOD, bold=True)
    notes.append(("Fujino-Morley 細節", ["這頁補強 Fujino-Morley 的角色。", "可強調它不是只為了算快，而是為了從離散計算回推嚴格的連續問題界線。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Bernstein polynomial 的角色", "用 convex-hull property 處理 L∞ 最大範數限制", page)
    add_card(slide, 0.65, 1.15, 12.05, 5.1, "Convex-hull property")
    add_formula(slide, 0.95, 1.82, 11.45, 0.65, "p(x)=Σ dᵢⱼₖ Jᵢⱼₖ(x)", size=24)
    add_bullets(
        slide,
        1.08,
        2.9,
        11.05,
        2.38,
        [
            "Bernstein polynomial 可由 control coefficients 表示",
            "convex-hull property：函數值會被 Bernstein coefficients 的範圍控制",
            "因此可把難處理的 maximum norm constraint 轉成係數層級的線性限制",
            "這是作者演算法能有效處理 L∞ constraint 的關鍵",
        ],
        size=15,
    )
    notes.append(("Bernstein polynomial", ["這頁是方法層第二個工具。", "可用一句話說：Bernstein representation 讓最大值約束可以被 control coefficients 包住。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Bernstein convex-hull property 的示範", "用 control coefficients 直觀看最大值限制", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "Simple interpretation")
    add_formula(slide, 0.95, 1.78, 11.45, 0.55, "p(x)=Σ dᵢ Bᵢ(x),   Bᵢ(x)≥0,   ΣBᵢ(x)=1", size=22)
    add_formula(slide, 0.95, 2.58, 11.45, 0.55, "if  -1 ≤ dᵢ ≤ 1  for all i,  then  -1 ≤ p(x) ≤ 1", size=22)
    add_bullets(
        slide,
        1.05,
        3.55,
        11.1,
        1.65,
        [
            "因為 p(x) 是 control coefficients 的 convex combination",
            "所以限制 coefficients 的範圍，就能保守地限制整個函數的 L∞ norm",
            "這正好解決原問題中最難處理的 maximum norm constraint",
        ],
        size=14,
    )
    add_text(slide, 1.0, 5.72, 11.2, 0.34, "報告口訣：Bernstein 把「到處都要檢查最大值」變成「檢查一組係數」。", size=13, color=GOOD, bold=True)
    notes.append(("Bernstein 示範", ["這頁用簡單 convex combination 解釋 Bernstein。", "聽眾不需要知道所有 Bernstein basis 細節，只要知道係數界可以控制函數最大值。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "演算法流程", "從 triangle mesh 到 rigorous bounds", page)
    add_step_flow(
        slide,
        [
            ("1. Mesh", "對三角形 K 做 triangulation"),
            ("2. FEM", "建立 Fujino-Morley space 與矩陣"),
            ("3. Bernstein", "轉成 Bernstein coefficients constraint"),
            ("4. Bounds", "計算 λ 的 lower bound 與 Cᴸ 的 upper bound"),
        ],
        1.75,
    )
    add_card(slide, 0.72, 4.2, 12.0, 1.58, "演算法目的")
    add_bullets(
        slide,
        1.05,
        4.76,
        11.25,
        0.72,
        [
            "不是只求一個浮點近似值，而是取得可驗證的 lower / upper bounds",
            "網格細化後，對具體三角形可取得相對誤差小於 1% 的估計",
        ],
        size=13,
    )
    notes.append(("演算法流程", ["這頁把 Section 3 的演算法轉成報告聽眾容易理解的流程。", "重點是 mesh、FEM、Bernstein、bounds 四個步驟。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "演算法細節：實際上算了什麼？", "從矩陣問題取得 λ_h,B，再回推出 Cᴸ(K)", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "Detailed computation path")
    add_simple_table(
        slide,
        0.95,
        1.78,
        [1.4, 4.35, 5.35],
        0.56,
        ["順序", "計算內容", "輸出或作用"],
        [
            ["a", "建立 triangulation T_h", "把 K 切成較小三角形"],
            ["b", "建立 Fujino-Morley basis 與 stiffness matrix", "把 |u|₂,K 寫成矩陣 quadratic form"],
            ["c", "轉成 Bernstein coefficients", "讓 L∞ constraint 可由係數限制處理"],
            ["d", "解離散 minimization problem", "得到 λ_h,B 或 λ_h 的 lower estimate"],
            ["e", "套用 Theorem 3.1 / Corollary 3.1", "得到 λ(K) lower bound 與 Cᴸ(K) upper bound"],
        ],
        font_size=10,
    )
    add_text(slide, 1.0, 5.55, 11.2, 0.34, "這也是本文題目中 error-constant estimation 的核心：不是只建模，而是把常數估計流程做成可驗證演算法。", size=12, color=MUTED)
    notes.append(("演算法細節", ["這頁補充 Section 3 的實作型流程。", "可以照順序講 a 到 e，讓聽眾知道作者到底算了什麼。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "數值結果一：unit right isosceles triangle", "具體三角形上的常數上下界", page)
    add_card(slide, 0.65, 1.14, 12.05, 5.1, "K = K_{1,π/2}")
    add_formula(slide, 1.0, 1.85, 11.3, 0.72, "0.40432 ≤ Cᴸ(1,π/2) ≤ 0.41596", size=30)
    add_bullets(
        slide,
        1.08,
        3.05,
        11.1,
        2.05,
        [
            "這是文中最容易報告的一個代表性結果",
            "上界由離散 optimization 與 Corollary 3.1 推得",
            "下界由高次多項式近似 minimizer 推得",
            "上下界接近，表示演算法對此三角形估計很銳利",
        ],
        size=15,
    )
    add_text(slide, 0.9, 6.22, 11.6, 0.28, "文中也用 interval arithmetic 檢查浮點 round-off error；例如 h=1/64 時，上界落在非常窄的區間內。", size=12, color=MUTED)
    notes.append(("數值結果一", ["這頁報告最具體的數字。", "可以說：直角等腰三角形的最佳常數大約是 0.41，作者給出上下界而不是單一近似值。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "示範 3：λ 的下界怎麼換成 Cᴸ 的上界？", "以 unit right isosceles triangle、N=64 為例", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "Corollary 3.1 的數字流程")
    add_formula(slide, 0.95, 1.75, 11.45, 0.55, "λₕ,B = 5.7812,   N=64", size=24)
    add_formula(slide, 0.95, 2.52, 11.45, 0.55, "λ(K) ≥ λₕ,B(1 - 1/N²)", size=24)
    add_formula(slide, 0.95, 3.29, 11.45, 0.55, "λ(K) ≥ 5.7812 × (1 - 1/4096) = 5.7798", size=24)
    add_formula(slide, 0.95, 4.06, 11.45, 0.55, "Cᴸ(K) = 1/√λ(K)  ⇒  Cᴸ(K) ≤ 1/√5.7798 = 0.41596", size=23)
    add_bullets(
        slide,
        1.05,
        5.15,
        11.1,
        0.72,
        [
            "這就是文中 h=1/64 時得到上界 0.41596 的計算邏輯",
            "報告時可說：作者不是直接猜 Cᴸ，而是先嚴格估 λ 的下界，再反推出 Cᴸ 的上界",
        ],
        size=12,
    )
    notes.append(("λ 到 C 的示範", ["這頁把論文中比較抽象的 λ 和 Cᴸ 關係帶數字算出來。", "重點是：λ 有下界，因為 Cᴸ=1/√λ，所以 Cᴸ 就有上界。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Table 1：λ lower bound 為什麼重要？", "λ 越大，Cᴸ=1/√λ 越小，補間誤差常數越好", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "Selected λ lower bounds at h=1/64")
    add_simple_table(
        slide,
        1.05,
        1.82,
        [1.8, 2.45, 3.0, 3.0],
        0.6,
        ["θ", "λ_h,B", "Corollary 3.1 λ lower", "對 Cᴸ 的影響"],
        [
            ["π/3", "15.457", "15.454", "Cᴸ 約 ≤ 0.254"],
            ["π/2", "5.7812", "5.7799", "Cᴸ 約 ≤ 0.416"],
            ["5π/6", "1.0212", "1.0210", "Cᴸ 約 ≤ 0.990"],
        ],
        font_size=12,
    )
    add_bullets(
        slide,
        1.05,
        4.62,
        11.1,
        0.9,
        [
            "Table 1 報告的是 λ 的 lower bound；Table 2 再把它轉成 Cᴸ 的 upper bound",
            "同樣 h=1/64，π/3 的 λ 遠大於 5π/6，因此 Cᴸ 較小、補間常數較好",
        ],
        size=13,
    )
    notes.append(("Table 1 解讀", ["這頁補充 Table 1 的意義。", "重點是 λ 和 Cᴸ 是反比平方根關係，所以 λ lower bound 越大，Cᴸ upper bound 越小。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "數值結果二：不同形狀的三角形", "角度與形狀會改變 Cᴸ(K)", page)
    add_card(slide, 0.65, 1.15, 5.85, 5.1, "Table 2 的觀察")
    add_bullets(slide, 0.95, 1.82, 5.25, 3.85, ["作者比較 θ=π/6、π/4、π/3、π/2、2π/3、3π/4、5π/6", "不同角度下 Cᴸ 的上下界不同", "θ 趨近很大或很小時，三角形品質變差，常數可能上升", "結果支持前面 shape regularity 的理論直覺"], size=13)
    add_card(slide, 6.85, 1.15, 5.85, 5.1, "Figure 11 的觀察")
    add_bullets(slide, 7.15, 1.82, 5.25, 3.85, ["隨著 mesh refinement，上界與下界逐漸收斂", "不同形狀三角形的收斂行為不完全相同", "作者指出收斂效率的理論分析留待後續研究", "數值實驗驗證了方法的可行性"], size=13)
    notes.append(("數值結果二", ["這頁從單一三角形推到多種形狀。", "報告時可強調：補間誤差常數不是固定常數，它和元素幾何形狀強相關。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "示範 4：Table 2 的數字怎麼解讀？", "同樣 h=1/64，三角形越差，Cᴸ 通常越大", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "Selected rows from Table 2")
    add_simple_table(
        slide,
        1.0,
        1.82,
        [1.6, 2.9, 2.05, 2.05, 2.05],
        0.58,
        ["θ", "形狀直覺", "lower", "upper", "gap"],
        [
            ["π/3", "接近正三角形", "0.25209", "0.25439", "約 0.9%"],
            ["π/2", "直角等腰", "0.40419", "0.41596", "約 2.9%"],
            ["5π/6", "非常鈍角", "0.92830", "0.98968", "約 6.4%"],
        ],
        font_size=12,
    )
    add_bullets(
        slide,
        1.05,
        4.5,
        11.1,
        1.05,
        [
            "lower / upper 很接近時，代表常數被夾得很窄，估計更銳利",
            "π/3 的常數比 π/2 小很多；5π/6 明顯變大，對應三角形品質變差",
        ],
        size=13,
    )
    notes.append(("Table 2 數字示範", ["這頁把 Table 2 抽出三個代表角度。", "用這張表說明：Cᴸ 不是固定值，它會隨三角形形狀改變，而且上下界差距可作為估計銳利程度。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "Interval arithmetic：為什麼結果比較嚴格？", "作者處理浮點 round-off error，而不是只相信一般小數計算", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "Rigorous computation")
    add_formula(slide, 0.95, 1.85, 11.45, 0.62, "Cᴸ_ub(1,π/2) ∈ [0.4159516728, 0.4159516793]", size=23)
    add_bullets(
        slide,
        1.05,
        3.05,
        11.1,
        2.1,
        [
            "一般浮點計算會有 round-off error，尤其矩陣很大時可能累積",
            "作者用 interval arithmetic 組裝矩陣與評估上界",
            "結果區間非常窄，表示 round-off error 沒有明顯破壞常數估計",
            "這強化了本文「rigorous estimation」的可信度",
        ],
        size=15,
    )
    notes.append(("Interval arithmetic", ["這頁補上嚴格數值計算的特色。", "可說作者不只報小數，也用 interval arithmetic 證明這個小數範圍在 round-off error 下仍可靠。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "論文主要貢獻", "作者在 conclusion 中總結的兩個層次", page)
    add_card(slide, 0.65, 1.15, 5.85, 5.1, "理論貢獻")
    add_bullets(slide, 0.95, 1.82, 5.25, 3.85, ["給出三角形元素上 L∞ interpolation error constant 的明確估計框架", "Theorem 2.1 提供 arbitrary triangle 的上界", "說明 shape regularity 與 degenerate triangle 對常數的影響", "把幾何品質和 interpolation error 連起來"], size=13)
    add_card(slide, 6.85, 1.15, 5.85, 5.1, "計算貢獻")
    add_bullets(slide, 7.15, 1.82, 5.25, 3.85, ["提出 FEM-based algorithm 估計具體三角形的最佳常數", "用 Fujino-Morley space 連接連續與離散問題", "用 Bernstein convex-hull property 處理 maximum norm constraint", "用數值上下界與 interval arithmetic 支持嚴格性"], size=13)
    notes.append(("主要貢獻", ["這頁是報告總結這篇 paper 的核心。", "它的貢獻可分成理論上界與可計算的嚴格估計演算法。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "限制與未來工作", "作者沒有宣稱已完整解決所有 maximum-norm optimization 問題", page)
    add_card(slide, 0.65, 1.15, 12.05, 5.1, "Scope")
    add_bullets(
        slide,
        1.0,
        1.85,
        11.1,
        3.6,
        [
            "主要處理 linear Lagrange interpolation over triangular elements",
            "Theorem 2.1 的 general bound 可用但較 raw；銳利估計依賴具體三角形上的數值程序",
            "作者指出 optimization approach 的收斂性與效率仍需後續系統性研究",
            "這篇的價值在嚴格估計 interpolation error constant，不是提出新的應用模型",
        ],
        size=15,
    )
    add_text(slide, 1.0, 5.78, 11.1, 0.32, "報告時可說：這篇論文的強項是嚴格數學與可驗證計算；限制是問題範圍相當專門。", size=12, color=WARN, bold=True)
    notes.append(("限制與未來工作", ["這頁幫你避免把文獻講過頭。", "它不是萬用補間理論，也不是應用實驗論文，而是針對一類有限元素補間常數做嚴格估計。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "可能被問的問題", "準備幾個短答，報告時比較穩", page)
    add_card(slide, 0.65, 1.12, 12.05, 5.18, "Q&A")
    add_simple_table(
        slide,
        0.95,
        1.72,
        [4.0, 7.0],
        0.68,
        ["問題", "建議回答"],
        [
            ["為什麼不用平均誤差？", "本文關心 maximum norm，因為它能控制最壞點誤差。"],
            ["為什麼三角形形狀重要？", "affine transformation 會改變 H² seminorm；退化元素會放大常數。"],
            ["Bernstein 的用處是什麼？", "用 convex-hull property 把 L∞ 約束轉成係數約束。"],
            ["這篇最主要的貢獻？", "一般三角形上界 + 具體三角形的嚴格數值估計流程。"],
            ["限制在哪裡？", "問題範圍專門，且 optimization convergence 的完整理論仍是未來工作。"],
        ],
        font_size=10,
    )
    notes.append(("Q&A", ["這頁可作為備用，不一定完整講。", "如果老師提問，可以用這些短答快速回到本文主線。"]))

    page += 1
    slide = new_slide(prs)
    add_header(slide, "報告結論", "用三句話收尾", page)
    add_card(slide, 0.65, 1.15, 12.05, 5.1, "Takeaways")
    add_bullets(
        slide,
        1.0,
        1.85,
        11.1,
        3.6,
        [
            "這篇論文研究的是 linear Lagrange interpolation 在 maximum norm 下的 error constant estimation",
            "它把三角形元素的幾何形狀、H² seminorm 與最大點誤差界連接起來",
            "主要方法是 FEM 離散化、Fujino-Morley space 與 Bernstein convex-hull property 的組合",
            "數值結果顯示，對具體三角形可取得接近的上下界，並能看出三角形形狀對常數的影響",
        ],
        size=15,
    )
    add_text(slide, 1.0, 5.78, 11.1, 0.32, "簡短結語：它是一篇把「線性補間最壞誤差」做成可嚴格估計問題的數值分析論文。", size=12, color=GOOD, bold=True)
    notes.append(("結論", ["最後用三句話收束：研究問題、方法、結果。", "如果老師問為什麼選這篇，可以回答：因為它提供 maximum-norm interpolation error bound 的嚴格背景。"]))

    return notes


def write_outline(notes: Iterable[tuple[str, list[str]]]) -> None:
    lines = [
        "# Galindo et al. (2022) 論文專題報告簡報大綱",
        "",
        f"- PDF source: `{PDF_SOURCE.relative_to(ROOT)}`",
        "- Generated PPTX: `outputs/papers/lagrange_interpolation_paper_report_zh.pptx`",
        "- Citation: Galindo, S. M., Ike, K., & Liu, X. (2022). Error-constant estimation under the maximum norm for linear Lagrange interpolation. Journal of Inequalities and Applications, 2022, 109. https://doi.org/10.1186/s13660-022-02841-w",
        "",
        "## 核心定位",
        "",
        "這是一篇數值分析與有限元素補間誤差常數論文。它研究 linear Lagrange interpolation over triangular elements 在 maximum norm 下的 error constant estimation，重點不在應用系統，而在如何對 `||u-Π^L u||_∞ ≤ C^L(K)|u|_{2,K}` 中的 `C^L(K)` 給出明確上界與可驗證數值估計。",
        "",
        "## Slide-by-slide 講稿",
        "",
    ]
    for idx, (title, bullets) in enumerate(notes, start=1):
        lines.append(f"### Slide {idx}: {title}")
        for bullet in bullets:
            lines.append(f"- {bullet}")
        lines.append("")
    lines.extend(
        [
            "## 60 秒摘要",
            "",
            "Galindo、Ike 與 Liu 這篇 2022 年論文研究 linear Lagrange interpolation 在三角形元素上的 maximum-norm error constant。作者先回顧一維線性補間的最佳誤差界，再把問題推到二維三角形元素，目標是估計 `||u-Π^L u||_{∞,K} ≤ C^L(K)|u|_{2,K}` 中的常數。理論上，Theorem 2.1 透過 affine transformation 給出一般三角形的上界，也說明三角形形狀越退化，誤差常數可能越大。計算上，作者把最佳常數估計轉為 optimization problem，使用 Fujino-Morley finite element space 離散化，再用 Bernstein polynomial 的 convex-hull property 處理 L∞ constraint。數值結果對多種三角形給出上下界，例如 unit right isosceles triangle 有 `0.40432 ≤ C^L(1,π/2) ≤ 0.41596`。因此，這篇論文的主要價值是把線性補間的最壞點誤差常數變成可理論分析、可數值驗證的問題。",
            "",
            "## 報告提醒",
            "",
            "本版本簡報專門報告該論文本身，不主動延伸到 thesis 公式 22。若報告後被問到和自己研究的關係，再補充它可作為 interpolation error bound 的數學背景即可。",
        ]
    )
    OUTLINE_PATH.write_text("\n".join(lines), encoding="utf-8")


def write_full_report() -> None:
    lines = [
        "# Galindo et al. (2022) 論文報告",
        "",
        "## 一、論文基本資料",
        "",
        "- 作者：Shirley Mae Galindo、Koichiro Ike、Xuefeng Liu",
        "- 題名：Error-constant estimation under the maximum norm for linear Lagrange interpolation",
        "- 期刊：Journal of Inequalities and Applications",
        "- 年份與文章編號：2022, 109",
        "- DOI：10.1186/s13660-022-02841-w",
        "- PDF：`docs/papers/data source/s13660-022-02841-w.pdf`",
        "",
        "這篇論文屬於數值分析與有限元素方法的補間誤差研究。它不是應用系統論文，而是處理一個很基礎但重要的問題：當我們用 linear Lagrange interpolation 在三角形元素上近似一個函數時，最大點誤差可以如何用一個明確的常數控制。",
        "",
        "## 二、研究動機與問題",
        "",
        "Lagrange interpolation 的基本想法是：已知節點上的函數值，建立一個通過這些節點的多項式近似。在一維區間中，linear Lagrange interpolation 就是用端點連成直線；在二維三角形中，則是用三個頂點值決定一個線性平面。這種補間在 finite element method 中非常常見，因為三角形元素與線性基底函數能把連續偏微分方程問題轉成可計算的離散問題。",
        "",
        "問題是，補間雖然在節點上完全正確，但在元素內部會有誤差。本文關心的不是平均誤差，而是 maximum norm，也就是整個三角形內最壞位置的點誤差。作者研究的核心不等式可以寫成：",
        "",
        "`||u - Π^L u||_{∞,K} ≤ C^L(K)|u|_{2,K}`",
        "",
        "其中 `K` 是三角形元素，`Π^L u` 是 linear Lagrange interpolation，`||·||_{∞,K}` 是最大範數，`|u|_{2,K}` 是二階 Sobolev seminorm，而 `C^L(K)` 就是本文要估計的 interpolation error constant。這個常數會受到三角形形狀影響，因此不是一個和幾何無關的固定數字。",
        "",
        "## 三、一維補間誤差的背景",
        "",
        "作者在 introduction 中先回顧一維 linear Lagrange interpolation 的標準最佳誤差界。若區間為 `I=(0,1)`，則有：",
        "",
        "`||u - Π^L u||_{∞,I} ≤ (1/8)||u''||_{∞,I}`",
        "",
        "這個 `1/8` 常數是 optimal estimate。以 `u(x)=x^2` 為例，端點為 `u(0)=0`、`u(1)=1`，線性補間為 `Π^L u(x)=x`。此時誤差為 `|x^2-x|=x(1-x)`，最大值出現在 `x=0.5`，最大誤差為 `0.25`。另一方面，`u''=2`，所以右側 `(1/8)||u''||∞=(1/8)×2=0.25`，剛好達到等號。這個例子說明，線性補間無法捕捉二階彎曲，而最大誤差正是由補間跨度與二階導數控制。",
        "",
        "若區間長度改為 `h`，誤差界會變成 `(h^2/8)sup|u''|`。例如同樣的二階導數上界下，`h=2` 時誤差會是 `h=1` 時的四倍。這個 scaling 也提供後面理解三角形元素幾何影響的直覺：元素越大或形狀越差，補間誤差常數通常越不利。",
        "",
        "## 四、既有文獻與本文切入點",
        "",
        "在二維三角形元素上，既有研究已經處理過不同範數下的 interpolation error constants。例如 L2 norm 與 H1 seminorm 的誤差常數估計已有相關結果；Waldron 也曾對 simplex 頂點上的線性補間提供 L∞ sharp inequality。Waldron 的結果使用幾何量如外接圓半徑 `R` 與距離 `d`，並以二階 directional derivative 的 L∞ 型條件控制誤差。",
        "",
        "Galindo et al. 的切入點不同：本文考慮的是三角形 `K` 上的 maximum norm interpolation error，但右側用的是 H2 seminorm，也就是 `|u|_{2,K}`。作者的目標是對 `C^L(K)` 給出明確上界，並對具體三角形提供更銳利的 lower/upper bounds。困難在於 maximum norm constraint 不容易直接最佳化，而且三角形形狀會明顯影響誤差常數。",
        "",
        "## 五、三角形參數化與 Theorem 2.1",
        "",
        "作者用三個幾何參數描述一般三角形：",
        "",
        "`p1=(0,0), p2=(h,0), p3=(αh cosθ, αh sinθ)`",
        "",
        "其中 `h` 是基準邊長，`αh` 是另一條邊長，`θ` 是夾角。當 `h=1` 時，三角形記為 `K_{α,θ}`。由 scaling 性質可知，`C^L(α,θ,h)=hC^L(α,θ,1)`，因此作者可以先分析單位尺度下不同形狀三角形的誤差常數。",
        "",
        "Theorem 2.1 給出一般三角形上的 raw upper bound：",
        "",
        "`C^L(α,θ) ≤ [v_+(α,θ)/(√2 α sinθ)] C^L(1,π/2)`",
        "",
        "其中 `v_+(α,θ)=1+α^2+√(1+2α^2 cos2θ+α^4)`。證明概念是使用 affine transformation，把一般三角形映射回 reference triangle，也就是 unit right isosceles triangle。L∞ norm 在對應座標下保持最大值關係，但 H2 seminorm 會受到幾何形變影響，因此誤差常數會出現幾何 multiplier。",
        "",
        "這個定理的意義是：三角形形狀會直接進入誤差常數。若三角形 shape-regular，也就是最小內角有正下界，常數可以保持 bounded；若三角形趨近退化線段，常數可能趨近無限大。這也呼應有限元素方法中 mesh quality 的重要性。",
        "",
        "## 六、Theorem 2.1 的數字示範",
        "",
        "若取 `α=1`，並使用文中 `C^L(1,π/2)≤0.41596` 作為 reference，可以代入不同角度觀察 raw bound 的變化。當 `θ=60°` 時，multiplier 約為 `2.449`，raw bound 約為 `≤1.019`；當 `θ=90°` 時，multiplier 約為 `1.414`，raw bound 約為 `≤0.588`；當 `θ=150°` 時，multiplier 約為 `5.278`，raw bound 約為 `≤2.195`。這些數字顯示，鈍角或形狀較差的三角形會讓一般上界明顯變大。",
        "",
        "需要注意的是，Theorem 2.1 的上界偏保守，它的主要用途是說明任意三角形都可以得到可用的 upper bound，以及幾何退化會讓 bound 變差。若要得到具體三角形上的銳利常數，則需要本文 Section 3 的 FEM-based optimization algorithm。",
        "",
        "## 七、從 C^L(K) 轉成 λ(K)",
        "",
        "Section 3 的重點是把最佳常數估計轉成 optimization problem。作者定義與 `C^L(K)` 對應的量 `λ(K)`，使得：",
        "",
        "`C^L(K)=1/√λ(K)`",
        "",
        "直觀上，`C^L(K)` 是一個 supremum ratio，而 `λ(K)` 是對應倒數平方形式的 infimum problem。因此，如果可以證明 `λ(K)` 不小於某個值，就能推出 `C^L(K)` 不大於 `1/√λ(K)`。這個轉換很重要，因為它讓問題可以被寫成有限元素矩陣與最佳化問題。",
        "",
        "## 八、Fujino-Morley 與 Bernstein polynomial 的角色",
        "",
        "原本的 optimization problem 是在無限維函數空間 `H^2(K)` 上進行，不可能直接計算。作者使用 Fujino-Morley finite element space 進行離散化，並利用空間分解與正交性把離散問題的估計連回原本的連續問題。因此，Fujino-Morley 在本文中不只是計算工具，也是建立嚴格上下界關係的橋梁。",
        "",
        "另一個困難是 L∞ constraint。最大範數表示要控制整個三角形上所有點的最大值，這在 optimization 中不容易直接處理。作者利用 Bernstein polynomial 的 convex-hull property。若一個 polynomial 可寫成 Bernstein basis 的加權和，且 Bernstein basis 非負、總和為 1，則 polynomial 的值會落在 control coefficients 的 convex hull 中。簡單說，如果所有 Bernstein coefficients 都在 `[-1,1]` 內，函數值也會被限制在對應範圍內。這讓作者能把原本難處理的 maximum norm constraint 轉成係數層級的限制。",
        "",
        "## 九、演算法流程",
        "",
        "本文 Section 3 的實際計算流程可以整理為五步。第一，對三角形 `K` 做 triangulation。第二，建立 Fujino-Morley basis 與矩陣，把 `|u|_{2,K}` 寫成矩陣 quadratic form。第三，把有限元素函數轉成 Bernstein coefficients，使 maximum norm constraint 可被係數限制處理。第四，解離散 minimization problem，得到 `λ_h,B` 或相關的離散下界。第五，套用 Theorem 3.1 或 Corollary 3.1，將離散結果轉成 `λ(K)` 的 lower bound，再反推出 `C^L(K)` 的 upper bound。",
        "",
        "這個流程的重點不是單純求一個浮點近似，而是取得可驗證的 lower/upper bounds。這也是本文題目中 error-constant estimation 的核心。",
        "",
        "## 十、數值結果",
        "",
        "文中最容易報告的代表性結果是 unit right isosceles triangle，也就是 `K_{1,π/2}`。作者得到：",
        "",
        "`0.40432 ≤ C^L(1,π/2) ≤ 0.41596`",
        "",
        "這表示最佳常數大約落在 0.41 附近。上界由離散 optimization 與 Corollary 3.1 推得，下界則由高次多項式近似 minimizer 得到。上下界相近，代表演算法對這個三角形的估計相當銳利。",
        "",
        "以 `N=64` 為例，文中有 `λ_h,B=5.7812`。根據 Corollary 3.1，`λ(K) ≥ λ_h,B(1-1/N^2)`，也就是 `λ(K)≥5.7812×(1-1/4096)=5.7798`。由於 `C^L(K)=1/√λ(K)`，因此 `C^L(K)≤1/√5.7798=0.41596`。這就是上界數字的計算邏輯。",
        "",
        "Table 1 報告的是不同三角形角度下的 `λ` lower bound，Table 2 則把它轉成 `C^L` 的 lower/upper bounds。例如在 `h=1/64` 時，`θ=π/3` 的常數約落在 `0.25209` 到 `0.25439`，`θ=π/2` 約落在 `0.40419` 到 `0.41596`，而 `θ=5π/6` 則約落在 `0.92830` 到 `0.98968`。這清楚顯示不同形狀三角形的補間誤差常數差異很大，鈍角三角形的常數明顯較差。",
        "",
        "## 十一、嚴格數值計算",
        "",
        "作者也注意到浮點運算可能有 round-off error，因此使用 interval arithmetic 組裝矩陣並評估上界。例如在 unit right isosceles triangle、mesh size 為 `1/64` 時，嚴格上界落在非常窄的區間：",
        "",
        "`C^L_ub(1,π/2) ∈ [0.4159516728, 0.4159516793]`",
        "",
        "這表示 round-off error 沒有明顯破壞結果，也強化了本文「rigorous estimation」的可信度。",
        "",
        "## 十二、主要貢獻、限制與結論",
        "",
        "這篇論文的貢獻可以分成兩層。理論上，它給出三角形元素上 L∞ interpolation error constant 的明確估計框架，並透過 Theorem 2.1 說明任意三角形的上界與幾何形狀之間的關係。計算上，它提出 FEM-based algorithm，結合 Fujino-Morley space 與 Bernstein convex-hull property，對具體三角形取得嚴格且相當銳利的上下界。",
        "",
        "限制方面，本文主要處理的是 triangular elements 上的 linear Lagrange interpolation，問題範圍相當專門。Theorem 2.1 的一般上界可用但偏保守，真正銳利的結果仍依賴具體三角形上的數值程序。另外，作者也提到 optimization approach 的收斂性與效率仍需未來更系統性的研究。",
        "",
        "總結來說，Galindo et al. (2022) 是一篇把「線性補間最壞點誤差」做成可理論分析、可數值驗證問題的數值分析論文。它的價值不只在於給出某個常數，而是提供一套從幾何分析、最佳化轉換、有限元素離散化到嚴格數值驗證的完整估計流程。",
    ]
    REPORT_PATH.write_text("\n".join(lines), encoding="utf-8")


def write_speaker_notes(notes: Iterable[tuple[str, list[str]]]) -> None:
    scripts = {
        "封面": "今天我要報告的是 Galindo、Ike 和 Liu 在 2022 年發表的論文，主題是 linear Lagrange interpolation 在 maximum norm 下的誤差常數估計。這篇不是應用系統論文，而是數值分析論文，所以報告重點會放在它的數學問題、方法和數值結果。",
        "論文基本資料": "這頁先交代論文基本資料。它發表在 Journal of Inequalities and Applications，文章編號是 2022 年第 109 篇。關鍵詞包含 Lagrange interpolation、finite-element method、Fujino-Morley interpolation 和 Bernstein polynomial，從這些關鍵詞就可以看出，它主要處理有限元素補間誤差常數的嚴格估計。",
        "報告地圖": "這篇論文可以照四個部分讀。第一是 introduction 的補間誤差背景，第二是 Section 2 對一般三角形建立 raw upper bound，第三是 Section 3 用 FEM 和 Bernstein 多項式處理具體三角形的 sharp bounds，最後是數值結果和結論。",
        "Lagrange interpolation": "Lagrange interpolation 的核心想法是用節點上的函數值建立通過節點的多項式。在一維就是用端點連成直線；在二維三角形中，就是用三個頂點值決定一個線性平面。補間在節點上是完全正確的，但元素內部會有誤差，本文就是要估計這個最壞誤差。",
        "研究問題": "本文的核心不等式是最大點誤差小於等於一個常數乘上函數的二階 seminorm。左邊是補間誤差在三角形內最糟的點，右邊的 Cᴸ(K) 是本文要估計的常數。這個常數和三角形形狀有關，所以研究問題就是：不同三角形下這個常數到底多大。",
        "符號表": "這頁把後面常出現的符號先整理起來。K 是三角形元素，Πᴸu 是 linear Lagrange interpolation，maximum norm 是看區域內最壞點誤差，H² seminorm 衡量函數二階變化，Cᴸ(K) 是本文主角，λ(K) 則是為了計算 Cᴸ(K) 轉換出來的最小化量。",
        "1D 結果": "作者先回顧一維結果。若區間長度是 1，線性補間的 maximum norm 誤差可以由 1/8 乘上二階導數最大值控制。這個結果是最佳估計，意思是常數 1/8 不能再任意縮小。",
        "一維示範": "這裡用 u=x² 示範。端點值是 0 和 1，所以線性補間是 x。誤差是 x²-x 的絕對值，也就是 x(1-x)，最大值在 x=0.5，等於 0.25。右邊 1/8 乘上 u'' 的最大值，也就是 1/8 乘以 2，同樣等於 0.25，所以剛好達到等號。",
        "區間 scaling 示範": "如果區間長度改變，誤差界會按 h² 縮放。h=0.5 時誤差是 0.0625，h=1 時是 0.25，h=2 時變成 1.0。這個例子讓我們看到，補間跨度變大，最壞誤差會平方級放大。",
        "文獻缺口": "既有文獻已經有一維補間誤差、三角形上的 L² 和 H¹ 誤差常數，也有 Waldron 等人處理 simplex 上的 L∞ 誤差。但本文的目標是用 H² seminorm 控制 maximum norm 下的 linear Lagrange interpolation error constant，並且對具體三角形做明確估計。",
        "Waldron 比較": "Waldron 的結果也是 L∞ 誤差估計，但它使用外接圓半徑等幾何量，以及 W²,∞ 型的二階導數上界。Galindo et al. 則把右側寫成 H² seminorm，並且關心如何估計 Cᴸ(K)。所以本文不是否定 Waldron，而是在另一個函數空間設定下推進 error constant estimation。",
        "三角形參數化": "作者用 h、α 和 θ 描述一般三角形。p1 在原點，p2 在 x 軸上，p3 由 αh 和 θ 決定。這樣做的好處是，任意三角形的形狀可以系統地用幾何參數分析，也方便研究三角形變扁或角度變差時常數怎麼變。",
        "Theorem 2.1": "Theorem 2.1 給出一般三角形的上界。公式的重點是 Cᴸ(α,θ) 可以被 reference triangle 的 Cᴸ(1,π/2) 乘上一個幾何 multiplier 控制。這個 multiplier 包含 α 和 θ，所以它直接反映三角形形狀對誤差常數的影響。",
        "Theorem 2.1 證明流程": "證明可以分成幾個步驟。先用 affine transformation 把一般三角形對應到 reference triangle。接著比較 L∞ norm 和 H² seminorm 在座標轉換下的變化。最後套用 reference triangle 的常數，就得到一般三角形的上界。",
        "Theorem 2.1 數字示範": "這頁把 Theorem 2.1 代入數字。固定 α=1，只改 θ。當 θ=60 度時，上界約為 1.019；θ=90 度時約為 0.588；θ=150 度時變成約 2.195。這說明角度很鈍時，一般上界會明顯變差。",
        "Shape regularity": "這頁說明 shape regularity。若三角形最小內角有正下界，就不會無限變扁，誤差常數也可以保持 bounded。相反地，如果三角形趨近一條線段，面積趨近零，誤差常數可能變得很大。",
        "Optimization problem": "Section 3 開始把常數估計轉成最佳化問題。Cᴸ(K) 原本是一個 supremum ratio。作者定義 λ(K)，讓 Cᴸ(K)=1/√λ(K)。因此只要能找到 λ(K) 的下界，就能得到 Cᴸ(K) 的上界。",
        "Optimization 轉換細節": "這頁補充為什麼可以這樣轉。最大化一個比值，等價於最小化它倒數平方形式。Πᴸu=0 表示我們其實在看節點上為零的誤差函數空間。這個轉換讓問題更適合用矩陣和有限元素方法求解。",
        "Fujino-Morley": "原本問題在 H²(K) 這種無限維空間上，不能直接算。作者使用 Fujino-Morley finite element space 建立離散問題，並利用空間分解與正交性把離散問題連回連續問題。",
        "Fujino-Morley 細節": "Fujino-Morley 在這篇論文中不是單純的數值近似工具。它負責讓離散計算結果可以變成原本連續問題的嚴格 bound。也就是說，作者不只是算出一個近似值，而是建立可以證明的界線。",
        "Bernstein polynomial": "另一個困難是 maximum norm constraint 很難處理。作者使用 Bernstein polynomial 的 convex-hull property。直觀上，Bernstein 表示法可以用一組 control coefficients 控制整個 polynomial 的值域。",
        "Bernstein 示範": "如果 p(x) 是 Bernstein basis 的非負加權和，而且這些 basis 加起來等於 1，那麼 p(x) 就是 coefficients 的 convex combination。因此，只要所有 coefficients 都在 -1 到 1 之間，p(x) 也會被控制在這個範圍內。這讓 L∞ 約束變得可計算。",
        "演算法流程": "這頁整理演算法流程。先切 mesh，再建立 Fujino-Morley space 與矩陣，接著轉成 Bernstein coefficients，最後計算 λ 的 lower bound 和 Cᴸ 的 upper bound。重點是它輸出的是可驗證的上下界。",
        "演算法細節": "更細地看，作者先建立 triangulation，接著把 H² seminorm 寫成矩陣 quadratic form，再用 Bernstein coefficients 處理最大範數約束。解出離散問題後，再套用定理或推論，把離散結果轉成連續問題的 bound。",
        "數值結果一": "最代表性的結果是 unit right isosceles triangle。作者得到 Cᴸ 在 0.40432 和 0.41596 之間。上下界很接近，所以可以說最佳常數大約是 0.41。",
        "λ 到 C 的示範": "這頁示範上界怎麼算。當 λ_h,B=5.7812 且 N=64 時，套用 Corollary 3.1 得到 λ(K) 至少是 5.7798。因為 Cᴸ(K)=1/√λ(K)，所以 Cᴸ(K) 至多是 0.41596。",
        "Table 1 解讀": "Table 1 是 λ lower bound。λ 越大，Cᴸ 的上界越小。從表中可以看到 π/3 的 λ 比 5π/6 大很多，因此 π/3 的補間誤差常數比較好。",
        "數值結果二": "Table 2 比較不同角度三角形。作者列出 π/6 到 5π/6 等不同 θ。主要觀察是三角形形狀會強烈影響 Cᴸ，而且上下界會隨 mesh refinement 逐漸收斂。",
        "Table 2 數字示範": "這裡抽出三個代表角度。π/3 接近正三角形，Cᴸ 約 0.25；π/2 直角等腰，Cᴸ 約 0.41；5π/6 非常鈍角，Cᴸ 接近 1。這清楚顯示三角形品質變差時常數會變大。",
        "Interval arithmetic": "作者也用 interval arithmetic 檢查浮點誤差。對直角等腰三角形，Cᴸ 上界落在非常窄的區間內，這代表 round-off error 沒有明顯影響結果，也強化嚴格數值估計的可信度。",
        "主要貢獻": "本文貢獻可以分成理論和計算兩層。理論上，它建立一般三角形的誤差常數上界。計算上，它提出 FEM-based algorithm，結合 Fujino-Morley 和 Bernstein convex-hull property，取得具體三角形上的嚴格估計。",
        "限制與未來工作": "限制是本文範圍很專門，主要處理 triangular elements 上的 linear Lagrange interpolation。一般上界較保守，銳利估計依賴具體數值程序。作者也指出，optimization approach 的收斂性和效率還需要後續研究。",
        "Q&A": "這頁可以當備用。如果被問為什麼不用平均誤差，就回答 maximum norm 控制最壞點誤差。如果被問三角形形狀為什麼重要，就說 affine transformation 會改變 H² seminorm，退化元素會放大常數。",
        "結論": "最後總結，這篇論文研究的是 linear Lagrange interpolation 在 maximum norm 下的 error constant estimation。它把三角形幾何、H² seminorm 和最大點誤差界連起來，並用 FEM、Fujino-Morley 和 Bernstein polynomial 建立可驗證的估計流程。",
    }
    lines = [
        "# Galindo et al. (2022) 論文報告逐頁講稿",
        "",
        f"- 對應簡報：`{PPTX_PATH.relative_to(ROOT)}`",
        f"- 對應大綱：`{OUTLINE_PATH.relative_to(ROOT)}`",
        f"- 完整報告：`{REPORT_PATH.relative_to(ROOT)}`",
        "",
        "## 使用方式",
        "",
        "每頁先照「講稿」講主線，再視時間補充「提示」。如果時間只有 10 分鐘，可略過 Q&A 頁與部分示範頁；如果時間是 15 到 20 分鐘，可完整講完。",
        "",
    ]
    for idx, (title, bullets) in enumerate(notes, start=1):
        lines.append(f"## Slide {idx}: {title}")
        lines.append("")
        lines.append(scripts.get(title, "這頁依照投影片內容說明即可，重點是把它和前後頁的論文主線連起來。"))
        lines.append("")
        lines.append("提示：")
        for bullet in bullets:
            lines.append(f"- {bullet}")
        lines.append("")
    SPEAKER_NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    DOCS_THESIS.mkdir(parents=True, exist_ok=True)
    prs = new_presentation()
    notes = build_slides(prs)
    prs.save(PPTX_PATH)
    write_outline(notes)
    write_full_report()
    write_speaker_notes(notes)
    print(f"Wrote {PPTX_PATH.relative_to(ROOT)}")
    print(f"Wrote {OUTLINE_PATH.relative_to(ROOT)}")
    print(f"Wrote {REPORT_PATH.relative_to(ROOT)}")
    print(f"Wrote {SPEAKER_NOTES_PATH.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
