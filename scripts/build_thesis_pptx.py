from __future__ import annotations

import json
from pathlib import Path
from typing import Iterable, List, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from build_thesis_docx import ensure_image_asset


ROOT = Path(__file__).resolve().parents[1]
OUTPUTS = ROOT / "outputs"
PAPERS = OUTPUTS / "papers"
DATA = OUTPUTS / "data"
FIGURES = OUTPUTS / "figures"
ARCHITECTURE = FIGURES / "architecture"
PUBLIC_BENCHMARK_FIGURES = FIGURES / "public_benchmarks"
THESIS_PAPERS = ROOT / "docs" / "papers" / "thesis"
PRESENTATION_PATH = PAPERS / "thesis_presentation_zh.pptx"
STORED_PRESENTATION_PATH = THESIS_PAPERS / "thesis_presentation_zh.pptx"
OUTLINE_PATH = ROOT / "docs" / "thesis" / "presentation_outline_zh.md"
LONG_PRESENTATION_PATH = PAPERS / "thesis_presentation_zh_30min.pptx"
STORED_LONG_PRESENTATION_PATH = THESIS_PAPERS / "thesis_presentation_zh_30min.pptx"
LONG_OUTLINE_PATH = ROOT / "docs" / "thesis" / "presentation_outline_zh_30min.md"
LONG_SPEAKER_NOTES_PATH = ROOT / "docs" / "thesis" / "presentation_speaker_notes_zh_30min.md"

BACKGROUND_COLOR = RGBColor(238, 242, 247)
HEADER_FILL = RGBColor(23, 37, 61)
HEADER_TEXT = RGBColor(255, 255, 255)
HEADER_SUBTITLE = RGBColor(211, 225, 241)
TITLE_COLOR = HEADER_TEXT
TEXT_COLOR = RGBColor(25, 32, 40)
ACCENT_COLOR = RGBColor(0, 83, 130)
MUTED_COLOR = RGBColor(70, 80, 92)
CARD_FILL = RGBColor(255, 255, 255)
CARD_LINE = RGBColor(124, 143, 165)
BODY_FONT = "Noto Sans TC"
LATIN_FONT = "Arial"
FORMULA_FONT = "Cambria Math"


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def average_field_mae(summary: dict) -> dict:
    scenarios = summary.get("scenarios", [])
    metrics = ("temperature", "humidity", "illuminance")
    return {
        metric: round(sum(item["field_mae"][metric] for item in scenarios) / max(len(scenarios), 1), 4)
        for metric in metrics
    }


def best_recommendations(summary: dict) -> List[Tuple[str, str]]:
    output = []
    for item in summary.get("scenarios", []):
        recommendations = item.get("recommendations", [])
        best = recommendations[0]["name"] if recommendations else "n/a"
        output.append((item["name"], best))
    return output[:5]


def scenario_map(summary: dict) -> dict:
    return {item["name"]: item for item in summary.get("scenarios", [])}


def init_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def style_slide(slide) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = BACKGROUND_COLOR


def new_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide)
    return slide


def add_title(slide, text: str, subtitle: str = "") -> None:
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.16))
    header.fill.solid()
    header.fill.fore_color.rgb = HEADER_FILL
    header.line.color.rgb = HEADER_FILL
    title_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.17), Inches(12.2), Inches(0.62))
    frame = title_box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(22 if len(text) > 24 else 24)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(12.1), Inches(0.34))
        sub_frame = sub_box.text_frame
        sub_frame.clear()
        sub_p = sub_frame.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle
        sub_run.font.name = BODY_FONT
        sub_run.font.size = Pt(11)
        sub_run.font.color.rgb = HEADER_SUBTITLE


def add_footer(slide, page: int) -> None:
    box = slide.shapes.add_textbox(Inches(12.25), Inches(7.0), Inches(0.6), Inches(0.25))
    box.name = "footer_page_number"
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(page)
    run.font.name = LATIN_FONT
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED_COLOR


def renumber_footers(prs: Presentation) -> None:
    for index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.name != "footer_page_number" or not shape.has_text_frame:
                continue
            frame = shape.text_frame
            frame.clear()
            paragraph = frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.RIGHT
            run = paragraph.add_run()
            run.text = str(index)
            run.font.name = LATIN_FONT
            run.font.size = Pt(9)
            run.font.color.rgb = MUTED_COLOR
            break


def add_styled_run(
    paragraph,
    text: str,
    size: int,
    color: RGBColor,
    bold: bool = False,
    font_name: str = BODY_FONT,
) -> None:
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def metric_triplet(values: dict, decimals: int = 4) -> str:
    return (
        f"T={values['temperature']:.{decimals}f}, "
        f"H={values['humidity']:.{decimals}f}, "
        f"L={values['illuminance']:.{decimals}f}"
    )


def is_formula_line(text: str) -> bool:
    return text.startswith(
        (
            "T(",
            "H(",
            "L(",
            "p =",
            "v ∈",
            "Fᵥ",
            "F̂",
            "b₀",
            "T₀",
            "H₀",
            "L₀",
            "ζ",
            "A_ac",
            "A_win",
            "A_light",
            "N_T",
            "N_H",
            "N_L",
            "B_T",
            "S_T",
            "B_H",
            "S_H",
            "B_ac",
            "S_ac",
            "B_win",
            "S_win",
            "B_light",
            "S_light",
            "L_win",
            "L_light",
            "A_j",
            "Aⱼ",
            "E_j",
            "Eⱼ",
            "R_j",
            "Rⱼ",
            "Cᵥ",
            "C_v",
            "C(p)",
            "rᵛ",
            "R_v",
            "Rᵥ",
            "Penalty",
            "F_hybrid",
            "Fᵛ",
            "F_true",
            "Loss",
            "ℒ",
            "φᵢ",
            "yᵢ",
            "I^",
            "Iʳ",
            "𝒱",
            "|Rᵥ",
            "MAE",
            "RMSE",
            "IDW",
            "w_s",
            "Score",
            "×",
            "+",
            "+ Σ",
        )
    )


def add_bullets(slide, left: float, top: float, width: float, height: float, items: Sequence[str], level0_size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = BODY_FONT
        p.font.size = Pt(level0_size)
        p.font.color.rgb = TEXT_COLOR
        add_styled_run(p, item, level0_size, TEXT_COLOR)


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body_lines: Sequence[str],
    title_size: int = 16,
    body_size: int = 12,
    formula_size: int = 14,
) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_FILL
    shape.line.color.rgb = CARD_LINE
    shape.line.width = Pt(1.2)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.font.name = BODY_FONT
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(8)
    add_styled_run(p, title, title_size, ACCENT_COLOR, bold=True)
    for line in body_lines:
        para = frame.add_paragraph()
        is_formula = is_formula_line(line)
        para.font.name = FORMULA_FONT if is_formula else BODY_FONT
        para.font.size = Pt(formula_size if is_formula else body_size)
        para.font.color.rgb = TEXT_COLOR
        para.space_after = Pt(3)
        add_styled_run(
            para,
            line,
            formula_size if is_formula else body_size,
            TEXT_COLOR,
            font_name=FORMULA_FONT if is_formula else BODY_FONT,
        )


def add_picture(slide, source: Path, left: float, top: float, width: float, height: float) -> None:
    png = ensure_image_asset({"path": str(source.relative_to(ROOT)), "asset_name": source.stem})
    slide.shapes.add_picture(str(png), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_two_column_title_body(slide, title: str, left_items: Sequence[str], right_image: Path, subtitle: str = "") -> None:
    add_title(slide, title, subtitle)
    add_bullets(slide, 0.8, 1.55, 5.0, 4.9, left_items, level0_size=18)
    add_picture(slide, right_image, 6.3, 1.45, 6.2, 4.9)


FORMULA_WALKTHROUGH = [
    (
        "公式說明 1：三因子場與查詢點",
        "場的定義",
        [
            "T(p,t)：位置 p、時間 t 的溫度",
            "H(p,t)：位置 p、時間 t 的相對濕度",
            "L(p,t)：位置 p、時間 t 的照度",
            "p = (x,y,z)，單位為公尺",
            "t 可代表啟動後時間、情境時間或 demo 時間軸",
        ],
        "主張邊界",
        [
            "本研究不是只估單一平均值",
            "輸出是三維空間中任意點的三個環境量",
            "8 顆角落感測器只提供稀疏觀測",
            "其他採樣點是由模型與校正場推估出來",
            "所以報告時要用「估計／推估」，不要說「量測到」",
        ],
    ),
    (
        "公式說明 2：總估計式",
        "主公式",
        [
            "F̂ᵥ(p,t) = Nᵥ(p,t) + Cᵥ(p,t)",
            "v ∈ {T,H,L}",
            "T：temperature；H：relative humidity；L：illuminance",
            "Nᵥ：nominal model，描述主要物理趨勢",
            "Cᵥ：8 顆角落感測器 residual 形成的校正場",
        ],
        "為什麼這樣拆",
        [
            "Nᵥ 讓模型先有設備、邊界與空間結構",
            "Cᵥ 讓模型對齊真實角落感測器",
            "三個變數共用此估計框架",
            "但 N_T、N_H、N_L 的物理項目分開設計",
            "這能回應「不能把同一套 bulk/local 硬套三種物理量」",
        ],
    ),
    (
        "公式說明 3：Indoor baseline",
        "baseline 定義",
        [
            "b₀ = (T₀, H₀, L₀)",
            "T₀：設備作用前的室內基準溫度",
            "H₀：設備作用前的室內基準相對濕度",
            "L₀：設備作用前的室內基準照度",
            "若有啟動前觀測，可由 8 顆感測器平均取得",
        ],
        "跟 baseline 比較法的差別",
        [
            "這裡的 baseline 是模型起始狀態",
            "不是第 5 章的 IDW baseline",
            "也不是公開資料集的 persistence 或 linear regression",
            "Web demo 左側 Indoor Baseline 就是在設定 T₀/H₀/L₀",
            "所有設備影響都是在 b₀ 上加減偏移",
        ],
    ),
    (
        "公式說明 4：baseline 的取得方式",
        "有啟動前觀測時",
        [
            "T₀ = (1/|S|)Σ_{s∈S} O_T(p_s,t_ref)",
            "H₀ = (1/|S|)Σ_{s∈S} O_H(p_s,t_ref)",
            "L₀ = (1/|S|)Σ_{s∈S} O_L(p_s,t_ref)",
            "S 是 8 顆角落感測器集合",
            "t_ref 是設備尚未加入作用的參考時間",
        ],
        "沒有啟動前觀測時",
        [
            "改由房間設計檔、情境設定或 demo 輸入提供",
            "例如標準房間預設 T₀=29°C、H₀=67%、L₀=90 lux",
            "因此 baseline 不是模型學出來的黑盒值",
            "它是後續設備影響與 residual correction 的共同起點",
            "報告時要明確說出資料來源是哪一種",
        ],
    ),
    (
        "公式說明 5：高度正規化",
        "垂直座標",
        [
            "ζ = z / Hᵣ - 1/2",
            "Hᵣ：房間高度",
            "z：查詢點高度",
            "ζ 約落在 -0.5 到 0.5",
            "ζ > 0 表示偏上方；ζ < 0 表示偏下方",
        ],
        "為什麼需要",
        [
            "室內溫度與濕度可能存在垂直分層",
            "冷空氣、熱源與混合程度會讓上下層不同",
            "ζ 提供低成本的高度修正項",
            "照度主要由光源幾何與遮蔽處理",
            "所以高度不是三種變數完全相同地使用",
        ],
    ),
    (
        "公式說明 6：設備 activation",
        "時間響應",
        [
            "Aⱼ(t) = aⱼ(1 - exp(-t/τⱼ))",
            "j 代表某個設備，例如冷氣、窗戶或燈具",
            "aⱼ：設備影響的穩態比例或強度尺度",
            "τⱼ：接近穩態所需的時間常數",
            "t 越大，Aⱼ(t) 越接近 aⱼ",
        ],
        "使用原因",
        [
            "設備不會在啟動瞬間把全室改變到穩態",
            "冷氣降溫、除濕與窗戶交換都需要時間",
            "這是一階收斂近似，計算成本低且可解釋",
            "不是完整 HVAC transient simulation",
            "目標是表達主要時間趨勢",
        ],
    ),
    (
        "公式說明 7：influence envelope",
        "空間作用範圍",
        [
            "Eⱼ(p,t) = Aⱼ(t) Rⱼ(p) Dⱼ(p,t) Vⱼ(p)",
            "Aⱼ(t)：設備目前啟動強度",
            "Rⱼ(p)：距離衰減",
            "Dⱼ(p,t)：方向性，例如冷氣出風或窗戶日照方向",
            "Vⱼ(p)：可見性或遮蔽程度",
        ],
        "距離衰減",
        [
            "Rⱼ(p) = exp(-||p - pⱼ|| / rⱼ)",
            "pⱼ：設備位置",
            "rⱼ：設備作用半徑或衰減尺度",
            "距離越遠，局部影響越小",
            "但全室平均項仍由各變數自己的 B 項處理",
        ],
    ),
    (
        "公式說明 8：溫度場主式",
        "溫度 nominal model",
        [
            "N_T(p,t) = T₀ + B_T(t) + S_T(p,t) + γ_T M(t) ζ",
            "T₀：室內基準溫度",
            "B_T(t)：全室平均熱響應",
            "S_T(p,t)：局部空間熱響應",
            "γ_T M(t)ζ：垂直溫度分層",
        ],
        "使用原因",
        [
            "溫度受熱交換、熱源與空氣混合影響",
            "只做局部衰減會低估冷氣的全室降溫",
            "所以保留 B_T 表示整體室溫移動",
            "S_T 再表達出風口、窗邊或燈具附近差異",
            "這是控制導向 reduced-order 熱場近似",
        ],
    ),
    (
        "公式說明 9：溫度的全室與局部項",
        "分解式",
        [
            "B_T(t) = B_ac,T(t) + B_win,T(t) + B_light,T(t)",
            "S_T(p,t) = S_ac,T(p,t) + S_win,T(p,t) + S_light,T(p,t)",
            "B_T 負責全室平均狀態改變",
            "S_T 負責某些位置比較強的局部差異",
            "兩者合起來避免 local-only 或 well-mixed-only 的偏誤",
        ],
        "三類來源",
        [
            "冷氣：依模式與設定溫差讓室內趨冷或趨暖",
            "窗戶：依 T_out - T₀ 表示外氣熱交換方向",
            "燈具：在溫度路徑中視為小型熱源",
            "注意：燈具在溫度是熱源，在照度才是光源",
            "這就是變數專屬公式的意義",
        ],
    ),
    (
        "公式說明 10：冷氣溫度項",
        "冷氣全室項",
        [
            "B_ac,T(t) = s_m k_ac,Tᵍ d_T P_ac A_ac(t)",
            "s_m：冷房或暖房模式符號",
            "k_ac,Tᵍ：冷氣對全室溫度的增益係數",
            "d_T：設定溫度與室內基準的需求差",
            "P_ac：校正後冷氣 power scale",
        ],
        "冷氣局部項",
        [
            "S_ac,T(p,t) = s_m k_ac,Tˢ d_T P_ac E_ac(p,t)",
            "k_ac,Tˢ：冷氣局部空間增益",
            "E_ac(p,t)：出風口附近、方向與遮蔽造成的空間權重",
            "解讀方式：B 是全室趨勢，S 是出風口附近差異",
            "兩者不是任意疊加，而是修正 local-only 的缺陷",
        ],
    ),
    (
        "公式說明 11：窗戶與燈具溫度項",
        "窗戶熱交換",
        [
            "B_win,T(t) = k_win,Tᵍ (T_out - T₀) P_win A_win(t)",
            "S_win,T(p,t) = k_win,Tˢ (T_out - T₀) P_win E_win(p,t)",
            "T_out > T₀ 時偏升溫",
            "T_out < T₀ 時偏降溫",
            "P_win 表示窗戶開啟比例或校正後影響尺度",
        ],
        "燈具熱源",
        [
            "B_light,T(t) = k_light,Tᵍ P_light A_light(t)",
            "S_light,T(p,t) = k_light,Tˢ P_light E_light(p,t)",
            "燈具在溫度模型裡只代表發熱",
            "照明造成的 lux 變化由照度模型另外處理",
            "這可避免把光學效果誤解成熱場效果",
        ],
    ),
    (
        "公式說明 12：濕度場主式",
        "濕度 nominal model",
        [
            "N_H(p,t) = clip[0,100]{H₀ + B_H(t) + S_H(p,t) - γ_H M(t) ζ}",
            "H₀：室內基準相對濕度",
            "B_H(t)：全室平均水氣或除濕響應",
            "S_H(p,t)：局部水氣交換或除濕差異",
            "clip[0,100]：相對濕度限制在 0% 到 100%",
        ],
        "使用原因",
        [
            "相對濕度有明確物理範圍",
            "冷氣常見效果是除濕，所以符號方向不同於溫度",
            "窗戶則由室外濕度與室內基準濕度差決定",
            "本研究不主張完整 psychrometric model",
            "而是用低階水氣交換近似再由 sensor residual 校正",
        ],
    ),
    (
        "公式說明 13：濕度來源項",
        "全室濕度項",
        [
            "B_H(t) = -k_ac,Hᵍ d_H P_ac A_ac(t)",
            "+ k_win,Hᵍ (H_out - H₀) P_win A_win(t)",
            "冷氣項為負：表示除濕",
            "窗戶項正負由 H_out - H₀ 決定",
            "外面較濕時開窗提高濕度，較乾時降低濕度",
        ],
        "局部濕度項",
        [
            "S_H(p,t) = -k_ac,Hˢ d_H P_ac E_ac(p,t)",
            "+ k_win,Hˢ (H_out - H₀) P_win E_win(p,t)",
            "E_ac 讓除濕效果在冷氣影響區附近更強",
            "E_win 讓窗邊水氣交換較強",
            "濕度沒有使用燈具照度那套光學公式",
        ],
    ),
    (
        "公式說明 14：照度場主式",
        "照度 nominal model",
        [
            "N_L(p,t) = max{0, L₀ + L_winᵈⁱʳ(p,t)",
            "+ L_lightᵈⁱʳ(p,t) + L_winᵃᵐᵇ(p,t) + Iʳᵉᶠˡ(p,t)}",
            "L₀：室內基準照度",
            "max{0,...}：照度不可為負",
            "照度由光源、窗戶、遮蔽與反射決定",
        ],
        "為什麼不同於溫濕度",
        [
            "照度不是空氣混合或水氣交換問題",
            "它更接近光線幾何與可視性問題",
            "燈具與窗戶可造成局部高照度峰值",
            "所以保留 direct source 與 obstruction",
            "再用一次漫反射補足間接光",
        ],
    ),
    (
        "公式說明 15：直射光與環境光",
        "窗戶直射光",
        [
            "L_winᵈⁱʳ(p,t) = S_out d_f k_sol P_win E_win(p,t)",
            "S_out：室外日照強度",
            "d_f：與時間、季節或日照方向相關的折減",
            "k_sol：窗戶日照轉換係數",
            "E_win：窗戶到室內點的距離、方向與遮蔽權重",
        ],
        "燈具與環境光",
        [
            "L_lightᵈⁱʳ(p,t) = G_light P_light E_light(p,t)",
            "L_winᵃᵐᵇ(p,t)：窗戶帶來的擴散環境光",
            "G_light：燈具光通量尺度",
            "P_light：燈具啟動或校正後影響尺度",
            "直射與環境光分開，可描述窗邊與全室背景亮度",
        ],
    ),
    (
        "公式說明 16：一次漫反射",
        "反射公式",
        [
            "Iʳᵉᶠˡ(p,t) = Σ_s ρ_s Ī_s A_sʳᵉˡ exp(-||p-c_s||/ℓ_s)",
            "× max(0, n_s·r̂_s→p) V_s(p)",
            "s：牆、地板、天花板或家具表面",
            "ρ_s：表面反射率",
            "Ī_s：表面接收到的平均照度",
        ],
        "限制與說法",
        [
            "一次漫反射用來補足非直射區域的回填亮度",
            "它不是完整 ray tracing 或 radiosity",
            "只計算一次反射，因此成本較低",
            "可主張改善照度主要趨勢",
            "不能主張達到精密光學模擬等級",
        ],
    ),
    (
        "公式說明 17：8 參數校正多項式",
        "三線性形式",
        [
            "C(p) = c₀ + c₁X + c₂Y + c₃Z",
            "+ c₄XY + c₅XZ + c₆YZ + c₇XYZ",
            "X,Y,Z 是正規化房間座標",
            "8 個係數對應 8 個角落感測器約束",
            "比 affine 多了交互項，仍保持低階可解釋",
        ],
        "為什麼剛好 8 點",
        [
            "房間有地面四角與天花板四角",
            "每個變數在同一時間有 8 個 residual",
            "三線性校正場也有 8 個自由度",
            "因此可由 8 個角點唯一決定此低階校正場",
            "但不是唯一決定任意複雜真實場",
        ],
    ),
    (
        "公式說明 18：角點 residual",
        "residual 定義",
        [
            "rᵛ_{abc}(t) = Oᵥ(p_{abc},t) - Nᵥ(p_{abc},t)",
            "a,b,c ∈ {0,1}",
            "p_{abc}：其中一個房間角點",
            "Oᵥ：該角點感測器觀測值",
            "Nᵥ：nominal model 在同一角點的預測值",
        ],
        "直覺意義",
        [
            "rᵛ 是主模型在感測點的誤差",
            "如果 rᵛ 為正，代表模型低估該角點",
            "如果 rᵛ 為負，代表模型高估該角點",
            "校正場 Cᵥ 的任務就是把這些角點誤差平滑帶入室內",
            "這一步讓模型與真實稀疏感測資料對齊",
        ],
    ),
    (
        "公式說明 19：三線性校正式",
        "校正公式",
        [
            "Cᵥ(X,Y,Z,t) = Σ_{a,b,c∈{0,1}} rᵛ_{abc}(t)",
            "× ℓ_a(X) ℓ_b(Y) ℓ_c(Z)",
            "ℓ₀(u)=1-u，ℓ₁(u)=u",
            "X/Y/Z 皆在 0 到 1 之間",
            "每個內部點都是 8 個角落 residual 的加權和",
        ],
        "重要性質",
        [
            "所有權重非負且總和為 1",
            "所以這是房間內部補間，不是無限制外插",
            "在任一角點上，對應權重為 1，其餘為 0",
            "因此校正後感測器位置會與觀測一致",
            "這是本研究 8 點推估最核心的數學基礎",
        ],
    ),
    (
        "公式說明 20：校正後估計值",
        "回到主公式",
        [
            "F̂ᵥ(p,t) = Nᵥ(p,t) + Cᵥ(p,t)",
            "在角點：Cᵥ 等於觀測 residual",
            "所以 F̂ᵥ(p_{abc},t) = Oᵥ(p_{abc},t)",
            "在非角點：Cᵥ 是 8 個 residual 的三線性補間",
            "Nᵥ 則保留設備與物理結構的空間趨勢",
        ],
        "主張邊界",
        [
            "8 顆感測器不能直接量到所有點",
            "我們證明的是三線性 residual correction 的一致性與可表示範圍",
            "其他點是 nominal model 加上低階 residual 補間的估計",
            "所以主張是接近主要空間趨勢",
            "不是宣稱無條件還原任意真實室內場",
        ],
    ),
    (
        "公式說明 21：可完全表示的 residual 空間",
        "函數空間",
        [
            "𝒱 = span{1, X, Y, Z, XY, XZ, YZ, XYZ}",
            "這個空間的維度是 8",
            "三線性函數可由 8 個角點取值唯一決定",
            "如果真實 residual 屬於 𝒱",
            "則 8 個角點 residual 可完全重建整個 residual 場",
        ],
        "嚴謹主張",
        [
            "這不是說所有室內場都一定是三線性",
            "而是說在三線性 residual 假設下可完全重建",
            "對平滑但非三線性的 residual，則用誤差界描述接近程度",
            "對突發局部熱源、光斑或遮蔽尖峰，8 點不保證捕捉",
            "這樣的 claim boundary 比較嚴謹",
        ],
    ),
    (
        "公式說明 22：平滑 residual 的誤差界",
        "誤差上界",
        [
            "|Rᵥ(p,t) - Cᵥ(p,t)| ≤ W²M_xx/8 + L²M_yy/8",
            "+ H²M_zz/8",
            "Rᵥ：真實 residual",
            "Cᵥ：三線性校正 residual",
            "W/L/H：房間寬、長、高",
        ],
        "如何解釋「接近」",
        [
            "M_xx/M_yy/M_zz 是 residual 二階曲率上界",
            "曲率越小，代表 residual 越平滑",
            "平滑時，8 點三線性補間誤差可被上界限制",
            "所以可說「在平滑 residual 假設下接近真實情況」",
            "不能說高頻局部變化也一定準確",
        ],
    ),
    (
        "公式說明 23：非連網裝置影響學習",
        "特徵向量",
        [
            "φᵢ = [xᵢ,yᵢ,zᵢ,tᵢ, baseline, outdoor, F_temp,",
            "F_hum, F_illum, activations, powers, envelopes]",
            "φᵢ 是模型已知的情境特徵",
            "包含位置、時間、室內外條件與設備作用強度",
            "目標是從觀測差異學出非連網設備的影響係數",
        ],
        "標籤定義",
        [
            "yᵢᵛ = F_trueᵛ(pᵢ,tᵢ) - Fᵥ(pᵢ,tᵢ)",
            "yᵢᵛ 是主模型尚未吸收的 residual target",
            "before/after 觀測提供設備啟用造成的差異訊號",
            "least-squares 適合小資料、可解釋且可快速重算",
            "若多個未知事件重疊，學到的是混合效果",
        ],
    ),
    (
        "公式說明 24：Hybrid residual",
        "第二層修正",
        [
            "F_hybridᵛ(p,t) = Fᵥ(p,t) + Rᵥ(p,t; θᵥ)",
            "Fᵥ：前面可解釋的 base estimator",
            "Rᵥ：小型 neural network 預測的 residual",
            "θᵥ：該變數 residual model 的參數",
            "每個變數可有自己的 residual 修正器",
        ],
        "定位",
        [
            "hybrid 不是取代物理模型",
            "它只修正主模型剩下的系統性誤差",
            "這樣可保留可解釋結構",
            "也能用資料修正不易手刻的偏差",
            "LOO 結果代表標準情境 family 內 residual 可學習",
        ],
    ),
    (
        "公式說明 25：Hybrid 訓練目標",
        "residual label",
        [
            "Rᵥ*(p,t) = F_trueᵛ(p,t) - Fᵥ(p,t)",
            "F_trueᵛ：訓練或合成 truth 場",
            "Fᵥ：base estimator 輸出",
            "Rᵥ*：希望 neural network 學到的剩餘誤差",
            "訓練時不是直接預測整個場，而是預測 residual",
        ],
        "損失函數",
        [
            "ℒ(θᵥ) = (1/N)Σᵢ ||Rᵥ* - Rᵥ(pᵢ,tᵢ;θᵥ)||²",
            "+ λ||θᵥ||²",
            "第一項是 residual 預測誤差",
            "第二項是正則化，降低過擬合",
            "溫濕度響應較平滑可低通；照度因光源/遮蔽需保留快速跳變",
        ],
    ),
    (
        "公式說明 26：MAE、RMSE 與 Correlation",
        "誤差指標",
        [
            "MAE = (1/n) Σᵢ |ŷᵢ - yᵢ|",
            "RMSE = √[(1/n) Σᵢ (ŷᵢ - yᵢ)²]",
            "ŷᵢ：模型估計值",
            "yᵢ：truth 或觀測參考值",
            "n：比較樣本數",
        ],
        "使用原因",
        [
            "MAE 代表平均偏差，最直觀",
            "RMSE 會放大尖峰或離群誤差",
            "Correlation 用於公開資料時序任務，檢查趨勢是否同向",
            "三者搭配可避免只看單一指標",
            "照度量級較大，圖表常用 log-scale 避免遮蔽溫濕度差異",
        ],
    ),
    (
        "公式說明 27：IDW baseline",
        "IDW 插值",
        [
            "IDW(p) = Σ_s w_s O_s / Σ_s w_s",
            "w_s = 1 / (dist(p,s) + ε)^q",
            "s：感測器索引",
            "O_s：感測器觀測值",
            "q：距離權重指數",
        ],
        "為什麼拿它比較",
        [
            "IDW 是無設備物理先驗的幾何插值 baseline",
            "它只知道感測器位置與距離",
            "不知道冷氣出風、窗戶日照或燈具位置",
            "若本研究優於 IDW，表示設備與空間結構有提供額外資訊",
            "但 IDW 在無光源、平坦場時可能表現不差",
        ],
    ),
    (
        "公式說明 28：推薦排序與驗證",
        "推薦分數",
        [
            "Score(a) = P_before - P_after",
            "Penalty 越大代表越偏離舒適目標",
            "P_before：採取動作前的目標區域懲罰",
            "P_after：反事實模擬動作後的懲罰",
            "Score 越高，預期改善越大",
        ],
        "必須說清楚的限制",
        [
            "目前推薦排序是 counterfactual simulation",
            "它是模型預測下的改善，不等於已證明因果效果",
            "真正驗證需做 before/after intervention",
            "比較實際 comfort penalty reduction 與預測改善是否一致",
            "demo 可以呈現比較流程，但不能替代實測介入驗證",
        ],
    ),
]
def add_formula_walkthrough(prs: Presentation, start_page: int, compact: bool = False) -> int:
    page = start_page
    for title, left_title, left_lines, right_title, right_lines in FORMULA_WALKTHROUGH:
        slide = new_slide(prs)
        add_title(slide, title, "公式、符號意義、限制與可主張範圍")
        card_height = 5.25 if compact else 5.35
        add_card(slide, 0.65, 1.35, 6.05, card_height, left_title, left_lines)
        add_card(slide, 6.95, 1.35, 5.75, card_height, right_title, right_lines)
        add_footer(slide, page)
        page += 1
    return page


def build_presentation() -> Presentation:
    prs = init_presentation()
    validation_summary = read_json(DATA / "validation_summary.json")
    submission_summary = read_json(DATA / "submission_readiness_summary.json")
    window_summary = read_json(DATA / "window_matrix_summary.json")
    bedroom_summary = read_json(DATA / "bedroom_01_weekly" / "weekly_simulation_summary.json")
    avg_mae = average_field_mae(validation_summary)
    bedroom_aggregate = bedroom_summary["aggregate"]

    # Slide 1
    slide = new_slide(prs)
    add_title(
        slide,
        "單房間非連網家電環境影響學習之稀疏感測空間數位孿生原型",
        "A Sparse-Sensing Spatial Digital Twin for Learning Environmental Impacts of Non-Networked Appliances in a Single Room",
    )
    add_bullets(
        slide,
        0.9,
        1.7,
        5.4,
        3.8,
        [
            "研究生：林昀佑",
            "指導教授：易昶霈教授、沈慧宇副教授",
            "系所：國立彰化師範大學資訊工程學系碩士班",
            "主題：以 8 顆角落感測器重建單房間溫度、濕度、照度場",
        ],
        level0_size=20,
    )
    add_picture(slide, ARCHITECTURE / "整體分層架構.svg", 6.7, 1.5, 5.6, 4.8)
    add_footer(slide, 1)

    # Slide 2
    slide = new_slide(prs)
    add_title(slide, "研究問題與動機")
    add_card(
        slide,
        0.7,
        1.5,
        3.8,
        3.9,
        "問題背景",
        [
            "一般房間內仍存在大量非連網裝置",
            "冷氣、窗戶、照明會改變環境",
            "但系統通常無法直接讀到其狀態",
        ],
    )
    add_card(
        slide,
        4.75,
        1.5,
        3.8,
        3.9,
        "核心挑戰",
        [
            "只有少量感測器，無法直接知道全室分布",
            "裝置可能新增、移動，家具也會阻擋傳遞",
            "早期純插值與 local-only 模型都出現不合理結果",
        ],
    )
    add_card(
        slide,
        8.8,
        1.5,
        3.8,
        3.9,
        "研究目標",
        [
            "重建三因子空間場",
            "學習非連網裝置對環境的影響",
            "以 Web 與工具介面提供查詢與決策能力",
        ],
    )
    add_footer(slide, 2)

    # Slide 3
    slide = new_slide(prs)
    add_title(slide, "系統架構")
    add_picture(slide, ARCHITECTURE / "整體分層架構.svg", 0.8, 1.4, 6.0, 5.2)
    add_bullets(
        slide,
        7.1,
        1.6,
        5.2,
        5.0,
        [
            "入口分成使用者互動層與 AI 工具呼叫層",
            "所有請求都先進入服務編排層",
            "後端主體為環境數位孿生核心與校正學習層",
            "hybrid residual 只作為 optional 第二層修正器",
        ],
        level0_size=17,
    )
    add_footer(slide, 3)

    # Slide 4
    slide = new_slide(prs)
    add_title(slide, "房間拓樸、感測器與目標區域")
    add_picture(slide, ARCHITECTURE / "房間感測器與目標區域配置.svg", 0.8, 1.4, 6.0, 5.3)
    add_card(
        slide,
        7.1,
        1.55,
        5.0,
        4.6,
        "固定設定",
        [
            "房間尺寸：6 m × 4 m × 3 m",
            "感測器：地面四角 + 天花板四角，共 8 顆",
            "區域：window_zone / center_zone / door_side_zone",
            "設備：ac_main / window_main / light_main",
            "採樣網格：16 × 12 × 6",
        ],
    )
    add_footer(slide, 4)

    # Slide 5
    slide = new_slide(prs)
    add_title(slide, "數學模型：變數專屬場模型")
    add_card(
        slide,
        0.7,
        1.45,
        5.9,
        4.8,
        "核心表示式",
        [
            "F̂ᵥ(p,t) = Nᵥ(p,t) + Cᵥ(p,t)",
            "Nᵥ：變數專屬 nominal model",
            "Cᵥ：由角落感測器殘差擬合的 trilinear 校正場",
            "溫度：熱交換、熱源、垂直分層",
            "濕度：除濕、水氣交換、物理範圍限制",
            "照度：直射、環境光、一次漫反射",
        ],
    )
    add_card(
        slide,
        6.85,
        1.45,
        5.8,
        4.8,
        "模型特點",
        [
            "冷氣、窗戶、照明都是模組化裝置",
            "家具會造成 obstacle-aware attenuation",
            "動態響應採一階收斂近似",
            "三個變數共用座標、設備與校正框架",
            "但 nominal model 不共用同一套物理公式",
            "推薦必須先有 sample / cluster 與三因子目標",
        ],
    )
    add_footer(slide, 5)

    # Slide 6
    slide = new_slide(prs)
    add_title(slide, "模型學習、推論與推薦資料流")
    add_picture(slide, ARCHITECTURE / "模型學習推論與推薦資料流.svg", 0.45, 1.12, 12.45, 5.72)
    add_footer(slide, 6)

    # Slide 7
    slide = new_slide(prs)
    add_title(slide, "系統實作與介面")
    add_card(
        slide,
        0.7,
        1.4,
        3.8,
        4.9,
        "MCP：工具化介面",
        [
            "MCP 不做預測；只提供 AI 可呼叫 runtime tools",
            "initialize：設定 scenario、baseline、外部邊界、設備/家具、時間與 estimator",
            "sample point：查指定座標在特定時間/穩定態的三因子",
            "learn impacts：建立 before/after 觀測紀錄再學係數",
            "window direct：直接輸入外部溫濕度、日照與開窗比例",
            "rank actions：指定座標 sample + T/H/L 目標後才排序註冊設備操作",
        ],
    )
    add_card(
        slide,
        4.75,
        1.4,
        3.8,
        4.9,
        "Web Demo",
        [
            "可旋轉 3D 預覽",
            "時間軸、播放、point sample",
            "裝置與家具都可模組化設定",
            "支援 hybrid estimator toggle",
        ],
    )
    add_card(
        slide,
        8.8,
        1.4,
        3.8,
        4.9,
        "輸入模式",
        [
            "標準情境 8 組",
            "窗戶矩陣 48 組",
            "窗戶 direct input",
            "自訂家具與自訂裝置",
        ],
    )
    add_footer(slide, 7)

    # Slide 8
    slide = new_slide(prs)
    add_title(slide, "驗證流程與比較原則")
    add_picture(slide, ARCHITECTURE / "驗證與實驗流程圖.svg", 0.8, 1.35, 5.8, 5.1)
    add_bullets(
        slide,
        6.95,
        1.5,
        5.25,
        5.2,
        [
            "E1-E3：synthetic full-field、IDW baseline、ablation",
            "E4：非連網裝置影響學習與推薦排序",
            f"E5：window matrix {window_summary.get('count', 0)} 組外部邊界條件",
            "E6：hybrid no-Fourier 對照與 LOO cross-validation",
            f"E7：bedroom_01 {bedroom_summary['snapshot_count']} 筆，pillow hold-out",
            "E8 protocol、E9 public benchmark；demo 不是量化實驗",
        ],
        level0_size=17,
    )
    add_footer(slide, 8)

    # Slide 9
    slide = new_slide(prs)
    add_title(slide, "主要結果：場重建與 baseline 比較")
    add_card(
        slide,
        0.7,
        1.45,
        3.7,
        2.0,
        "平均 Field MAE",
        [
            f"Temperature: {avg_mae['temperature']}",
            f"Humidity: {avg_mae['humidity']}",
            f"Illuminance: {avg_mae['illuminance']}",
        ],
    )
    add_card(
        slide,
        0.7,
        3.75,
        3.7,
        2.3,
        "標準情境",
        [
            f"驗證情境數：{len(validation_summary.get('scenarios', []))}",
            "比較項目：field / sensors / zones / IDW / recommendations",
        ],
    )
    add_picture(slide, FIGURES / "all_active_temperature_3d.svg", 4.7, 1.45, 3.8, 4.7)
    add_picture(slide, FIGURES / "window_only_illuminance_3d.svg", 8.7, 1.45, 3.8, 4.7)
    add_footer(slide, 9)

    # Slide 10
    slide = new_slide(prs)
    add_title(slide, "Hybrid Residual Neural Network 結果")
    default_hybrid = submission_summary["default_holdout_hybrid"]
    no_fourier = submission_summary["no_fourier_holdout_hybrid"]
    loo = submission_summary["leave_one_scenario_out"]
    add_card(
        slide,
        0.7,
        1.45,
        5.0,
        4.9,
        "Held-out + LOO",
        [
            f"Default samples: {default_hybrid['dataset']['train_samples']} / {default_hybrid['dataset']['test_samples']}",
            f"Default hybrid MAE: {metric_triplet(default_hybrid['hybrid_test_field_mae'])}",
            f"No-Fourier hybrid MAE: {metric_triplet(no_fourier['hybrid_test_field_mae'])}",
            f"LOO avg hybrid MAE: {metric_triplet(loo['average_hybrid_field_mae'])}",
            f"LOO reduction: T {loo['average_field_mae_reduction_percent']['temperature']:.2f}%, H {loo['average_field_mae_reduction_percent']['humidity']:.2f}%, L {loo['average_field_mae_reduction_percent']['illuminance']:.2f}%",
        ],
    )
    add_picture(slide, FIGURES / "submission" / "field_mae_comparison.png", 6.0, 1.35, 6.2, 3.1)
    add_bullets(
        slide,
        6.1,
        4.65,
        5.9,
        1.5,
        [
            "hybrid residual 是第二層修正器，不取代主模型",
            "LOO 結果證明標準情境 family 內殘差可學習，不代表任意房間泛化",
            "另以 bedroom_01 真實快照檢查 sparse calibration 對 pillow 點的改善",
        ],
        level0_size=15,
    )
    add_footer(slide, 10)

    # Slide 11
    slide = new_slide(prs)
    add_title(slide, "公開資料任務拆解：SML2010 / CU-BEMS")
    add_picture(slide, PUBLIC_BENCHMARK_FIGURES / "sml2010_task_breakdown.svg", 0.6, 1.35, 6.0, 3.35)
    add_picture(slide, PUBLIC_BENCHMARK_FIGURES / "cu_bems_task_breakdown.svg", 6.75, 1.35, 6.0, 3.35)
    add_bullets(
        slide,
        0.9,
        5.05,
        11.6,
        1.45,
        [
            "S3 是主要優勢：事件/邊界 delta 需要變化方向，structured prior 比 persistence 與 linear regression 更有用",
            "S1 與 C2 是主要劣勢：短視窗照度高度自相關，且公開資料缺實際幾何、遮蔽與多燈具資訊",
            "CU-BEMS C1/C3 可勝過 linear regression，但商辦 zone-level persistence 太強，不能宣稱全面勝出",
        ],
        level0_size=15,
    )
    add_footer(slide, 11)

    # Slide 12
    slide = new_slide(prs)
    add_title(slide, "研究貢獻與資料策略")
    add_bullets(
        slide,
        0.9,
        1.55,
        11.4,
        5.3,
        [
            "提出以單房間、8 顆角落感測器為前提的三因子空間數位孿生原型",
            "以變數專屬 nominal model、power calibration 與 trilinear correction 建立可解釋估測流程",
            "以 least-squares 學習非連網裝置影響，並用 hybrid residual 做第二層修正",
            "明確拆分 synthetic full-field、real sparse calibration、public task-aligned 與 intervention validation",
            "完整 3D 場比較以 canonical synthetic benchmark 為主",
            f"真實臥室快照校正後 pillow MAE: {metric_triplet(bedroom_aggregate['estimated_pillow_mae'])}",
            "公開資料集則採 task-aligned benchmark：CU-BEMS / SML2010 / ASHRAE 各比相容子任務",
        ],
        level0_size=18,
    )
    add_footer(slide, 12)

    # Slide 13
    slide = new_slide(prs)
    add_title(slide, "結論與未來工作")
    add_card(
        slide,
        0.8,
        1.6,
        5.7,
        4.8,
        "結論",
        [
            "有限角落感測器下仍可用分層式模型重建單房間三因子分布",
            "非連網裝置可透過環境變化進行影響學習與校正",
            "bedroom_01 7 天快照顯示校正後可改善未參與 fitting 的 pillow 點",
            "各資料來源支援的 claim boundary 已拆開說明",
            "模型已能輸出區域估計、反事實推薦排序與 AI 可查詢工具",
        ],
    )
    add_card(
        slide,
        6.8,
        1.6,
        5.7,
        4.8,
        "未來工作",
        [
            "擴大 ESP32 長期真實資料",
            "擴充 CO2 / PM2.5 等因子",
            "改進 multi-zone / partition 模型",
            "補足 dense real-room ground truth",
            "執行推薦動作 before/after 介入驗證",
            "研究遠端 MCP 與閉環控制",
        ],
    )
    add_footer(slide, 13)
    add_formula_walkthrough(prs, 14, compact=True)
    renumber_footers(prs)
    return prs


def build_presentation_30min() -> Presentation:
    prs = init_presentation()
    validation_summary = read_json(DATA / "validation_summary.json")
    submission_summary = read_json(DATA / "submission_readiness_summary.json")
    window_summary = read_json(DATA / "window_matrix_summary.json")
    bedroom_summary = read_json(DATA / "bedroom_01_weekly" / "weekly_simulation_summary.json")
    avg_mae = average_field_mae(validation_summary)
    scenarios = scenario_map(validation_summary)
    bedroom_aggregate = bedroom_summary["aggregate"]

    # 1 cover
    slide = new_slide(prs)
    add_title(
        slide,
        "單房間非連網家電環境影響學習之稀疏感測空間數位孿生原型",
        "論文口試簡報",
    )
    add_bullets(
        slide,
        0.9,
        1.7,
        5.4,
        4.0,
        [
            "研究生：林昀佑",
            "指導教授：易昶霈教授、沈慧宇副教授",
            "國立彰化師範大學資訊工程學系碩士班",
            "主題：單房間三因子數位孿生、非連網裝置影響學習、工具化服務介面",
        ],
        level0_size=20,
    )
    add_picture(slide, ARCHITECTURE / "整體分層架構.svg", 6.7, 1.5, 5.6, 4.8)
    add_footer(slide, 1)

    # 2 roadmap
    slide = new_slide(prs)
    add_title(slide, "報告流程")
    add_bullets(
        slide,
        1.0,
        1.55,
        11.0,
        5.4,
        [
            "1. 問題背景與研究動機",
            "2. 文獻定位與研究缺口",
            "3. 系統架構與數學模型",
            "4. 感測器校正、影響學習與 hybrid residual",
            "5. 實作系統、驗證設計與實驗結果",
            "6. 結論、限制、未來工作與公式整理",
        ],
        level0_size=21,
    )
    add_footer(slide, 2)

    # research core
    slide = new_slide(prs)
    add_title(slide, "研究主軸與輸入輸出")
    add_card(
        slide,
        0.85,
        1.45,
        11.6,
        1.6,
        "研究主軸",
        [
            "少量角落感測 + 非連網家電 + 單房間幾何配置 → 可解釋、可校正、可學習的三因子空間場估計與決策支援。",
        ],
        body_size=15,
    )
    add_card(
        slide,
        0.85,
        3.35,
        3.65,
        2.8,
        "輸入",
        [
            "房間尺寸、裝置/家具配置",
            "8 顆角落感測器",
            "室內 baseline 與外部邊界",
            "時間、設備狀態、使用者目標",
        ],
        body_size=12,
    )
    add_card(
        slide,
        4.85,
        3.35,
        3.65,
        2.8,
        "模型",
        [
            "溫度、濕度、照度各自的 nominal model",
            "power calibration",
            "trilinear residual correction",
            "optional hybrid residual",
        ],
        body_size=12,
    )
    add_card(
        slide,
        8.85,
        3.35,
        3.6,
        2.8,
        "輸出",
        [
            "任意點/區域三因子估計",
            "3D 視覺化",
            "非連網裝置影響係數",
            "反事實推薦排序與 MCP 查詢",
        ],
        body_size=12,
    )
    add_footer(slide, 4)

    # 3 background
    slide = new_slide(prs)
    add_title(slide, "研究背景與問題")
    add_card(
        slide,
        0.7,
        1.45,
        3.8,
        4.8,
        "房間場景",
        [
            "真實房間通常只有少量感測器",
            "但使用者關心的是整個空間的舒適度",
            "不是單一點位數值",
        ],
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.8,
        4.8,
        "非連網裝置",
        [
            "冷氣、窗戶、照明常無 API",
            "狀態不可直接讀取",
            "卻持續改變溫度、濕度、照度",
        ],
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        4.8,
        "核心需求",
        [
            "重建全室三因子分布",
            "學習裝置對環境的影響",
            "在 sample/cluster + 三因子目標下支援推薦",
        ],
    )
    add_footer(slide, 3)

    # 4 questions and contributions
    slide = new_slide(prs)
    add_title(slide, "研究問題與貢獻")
    add_card(
        slide,
        0.75,
        1.45,
        5.6,
        4.9,
        "研究問題",
        [
            "RQ1：8 顆角落感測器能否重建單房間三因子空間場？",
            "RQ2：能否從環境資料學習非連網裝置影響？",
            "RQ3：sample/cluster 與三因子目標下能否排序控制動作？",
            "RQ4：能否將模型封裝成 MCP 可查詢工具？",
        ],
    )
    add_card(
        slide,
        6.7,
        1.45,
        5.9,
        4.9,
        "主要貢獻",
        [
            "single-room three-factor spatial digital twin",
            "power calibration + trilinear residual correction",
            "non-networked appliance impact learning",
            "hybrid residual + Fourier denoising",
            "task-aligned public benchmark strategy",
            "MCP / Gemma / Web 可互動原型",
        ],
    )
    add_footer(slide, 4)

    # 5 literature gap
    slide = new_slide(prs)
    add_title(slide, "文獻定位、研究缺口與比較原則")
    add_card(
        slide,
        0.7,
        1.45,
        3.8,
        5.0,
        "已有研究",
        [
            "房間尺度 IEQ 實驗",
            "有限感測器場重建",
            "hybrid thermal model",
            "建築 digital twin 平台",
        ],
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.8,
        5.0,
        "常見不足",
        [
            "多只看熱環境或單雙因子",
            "少處理非連網裝置學習",
            "少將模型做成 AI 可查詢工具",
            "多半不是單房間低成本原型",
        ],
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        5.0,
        "本研究定位",
        [
            "single-room",
            "limited corner sensors",
            "temperature + humidity + illuminance",
            "control-oriented + MCP-accessible",
            "public datasets only for aligned subtasks",
        ],
    )
    add_footer(slide, 5)

    # 6 architecture
    slide = new_slide(prs)
    add_title(slide, "整體系統架構")
    add_picture(slide, ARCHITECTURE / "整體分層架構.svg", 0.8, 1.35, 6.2, 5.3)
    add_bullets(
        slide,
        7.15,
        1.5,
        5.0,
        5.2,
        [
            "入口區分為人機互動層與 AI 工具呼叫層",
            "所有請求經服務編排後進入環境數位孿生核心",
            "校正與影響學習層負責修正稀疏感測誤差",
            "輸出為空間場估測、區域估測、動作排序與 3D 視覺化",
        ],
        level0_size=17,
    )
    add_footer(slide, 6)

    # 7 execution flow
    slide = new_slide(prs)
    add_title(slide, "主要執行資料流")
    add_picture(slide, ARCHITECTURE / "主要執行資料流.svg", 0.8, 1.35, 5.9, 5.25)
    add_bullets(
        slide,
        6.95,
        1.45,
        5.25,
        5.3,
        [
            "scenario 或 direct input 先進入 service",
            "套用 indoor baseline、裝置、家具與時間設定",
            "再做場估測、感測器校正與 dashboard 輸出",
            "MCP 和 Web 都走同一條執行路徑",
        ],
        level0_size=17,
    )
    add_footer(slide, 7)

    # 8 room topology
    slide = new_slide(prs)
    add_title(slide, "房間拓樸、感測器與目標區域")
    add_picture(slide, ARCHITECTURE / "房間感測器與目標區域配置.svg", 0.8, 1.35, 6.0, 5.3)
    add_card(
        slide,
        7.05,
        1.45,
        5.15,
        4.9,
        "固定研究設定",
        [
            "房間：6 × 4 × 3 m",
            "感測器：floor/ceiling 四角，共 8 顆",
            "區域：window / center / door-side",
            "核心裝置：ac_main / window_main / light_main",
            "解析度：16 × 12 × 6",
        ],
    )
    add_footer(slide, 8)

    # 9 devices and furniture
    slide = new_slide(prs)
    add_title(slide, "模組化裝置與家具阻擋")
    add_picture(slide, ARCHITECTURE / "可模組化裝置與家具架構.svg", 0.8, 1.35, 5.9, 5.1)
    add_bullets(
        slide,
        6.95,
        1.45,
        5.25,
        5.25,
        [
            "冷氣、窗戶、燈都可視為模組化裝置",
            "家具是可開關、可移動的阻擋物件",
            "阻擋效果會依幾何位置自適應調整",
            "Web 端可新增 custom devices 與 custom furniture",
        ],
        level0_size=17,
    )
    add_footer(slide, 9)

    # 10 math model
    slide = new_slide(prs)
    add_title(slide, "數學模型：變數專屬 nominal model")
    add_card(
        slide,
        0.75,
        1.4,
        6.0,
        5.0,
        "核心形式",
        [
            "F̂ᵥ(p,t) = Nᵥ(p,t) + Cᵥ(p,t)",
            "Nᵥ：溫度、濕度、照度各自的 nominal model",
            "Cᵥ：8 顆角落感測器建立的 residual correction",
            "溫度處理熱交換與垂直分層",
            "濕度處理除濕與外氣水氣交換",
            "照度處理光源幾何、遮蔽與反射",
        ],
    )
    add_card(
        slide,
        7.0,
        1.4,
        5.2,
        5.0,
        "設計理由",
        [
            "避免只有冷氣附近變冷、全室仍近乎不變的不合理情況",
            "早期純插值與 local-only 版本都失敗過",
            "避免把同一組 bulk/local 公式硬套到三種物理量",
            "比 CFD 輕量",
            "比純插值更可解釋",
            "適合控制導向與即時服務化",
        ],
    )
    add_footer(slide, 10)

    # method rationale
    slide = new_slide(prs)
    add_title(slide, "方法選擇：為什麼不是純插值、純物理或純黑盒")
    add_card(
        slide,
        0.7,
        1.45,
        3.85,
        4.95,
        "純插值不足",
        [
            "IDW 只看距離與感測點數值",
            "不知道冷氣出風口、窗戶方向、燈具位置",
            "照度與局部熱區容易被抹平",
            "適合作 baseline，不適合作主模型",
        ],
        body_size=13,
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.85,
        4.95,
        "完整物理太重",
        [
            "CFD / ray tracing 需要更多邊界條件與計算成本",
            "一般房間很難取得精確材質、風場與設備曲線",
            "本研究目標是控制導向與即時服務化",
            "因此採 reduced-order nominal model",
        ],
        body_size=13,
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        4.95,
        "純黑盒風險",
        [
            "小資料下容易過擬合",
            "難解釋設備、邊界與空間位置的作用",
            "hybrid residual 只學剩餘誤差",
            "保留主模型的可解釋結構",
        ],
        body_size=13,
    )
    add_footer(slide, 11)

    # 11 calibration and learning
    slide = new_slide(prs)
    add_title(slide, "模型學習、推論與推薦資料流")
    add_picture(slide, ARCHITECTURE / "模型學習推論與推薦資料流.svg", 0.45, 1.1, 12.45, 5.75)
    add_footer(slide, 11)

    # 12 implementation interfaces
    slide = new_slide(prs)
    add_title(slide, "系統實作與介面")
    add_card(
        slide,
        0.7,
        1.45,
        3.85,
        4.9,
        "MCP：工具化介面",
        [
            "MCP 不做預測；核心模型負責場估計、校正與排序",
            "initialize：註冊 scenario、baseline、外部邊界、設備/家具、時間與 estimator",
            "sample point：補足非感測點、可指定 elapsed/steady state",
            "learn impacts：start/finish before-after record",
            "window direct：輸入外部資料，不走 48 組 preset",
            "rank actions：以指定座標 sample 與 T/H/L 目標評估註冊設備操作",
        ],
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.85,
        4.9,
        "Gemma / Ollama",
        [
            "Gemma 透過 bridge 做 tool calling",
            "MCP 支援來自主機與 runtime",
            "不是宣稱模型原生支援 MCP",
        ],
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        4.9,
        "Web Demo",
        [
            "可旋轉 3D 預覽",
            "時間軸與播放",
            "裝置/家具模組化調整",
            "hybrid estimator toggle",
        ],
    )
    add_footer(slide, 12)

    # 13 validation design
    slide = new_slide(prs)
    add_title(slide, "驗證設計")
    add_picture(slide, ARCHITECTURE / "驗證與實驗流程圖.svg", 0.75, 1.35, 6.0, 5.2)
    add_bullets(
        slide,
        7.0,
        1.45,
        5.2,
        5.3,
        [
            "E1-E3：truth-adjusted simulation、IDW、synthetic ablation",
            f"E4-E6：裝置影響學習、{window_summary.get('count', 0)} 組 window matrix、hybrid no-Fourier/LOO",
            f"E7：bedroom_01 {bedroom_summary['snapshot_count']} 筆快照與 pillow 位置比較",
            "E8：推薦動作 before/after intervention protocol",
            "E9：public datasets 僅作 task-aligned benchmark",
            "Web demo 與 3D 展示是呈現層，不列為量化實驗",
        ],
        level0_size=16,
    )
    add_footer(slide, 13)

    # evidence boundary
    slide = new_slide(prs)
    add_title(slide, "證據鏈與 Claim Boundary")
    add_card(
        slide,
        0.65,
        1.35,
        3.0,
        5.25,
        "Synthetic full-field",
        [
            "支援完整 3D 場誤差比較",
            "可比較 IDW、base、ablation、hybrid",
            "限制：truth 仍來自受控生成與調整，不等同長期真實場",
        ],
        body_size=12,
    )
    add_card(
        slide,
        3.85,
        1.35,
        3.0,
        5.25,
        "Real-bedroom snapshot",
        [
            "支援稀疏感測校正的真實點位檢查",
            "pillow hold-out 不參與 fitting",
            "限制：不是 dense real-room ground truth",
        ],
        body_size=12,
    )
    add_card(
        slide,
        7.05,
        1.35,
        2.85,
        5.25,
        "Public datasets",
        [
            "支援相容子任務比較",
            "SML2010 / CU-BEMS 只做 task-aligned benchmark",
            "限制：缺單房間幾何與 8 點拓樸",
        ],
        body_size=12,
    )
    add_card(
        slide,
        10.1,
        1.35,
        2.7,
        5.25,
        "Recommendation",
        [
            "目前可做反事實排序",
            "E8 定義介入驗證 protocol",
            "限制：尚未完成真實 before/after 因果驗證",
        ],
        body_size=12,
    )
    add_footer(slide, 14)

    # 14 scenarios and time/window settings
    slide = new_slide(prs)
    add_title(slide, "情境設計與輸入模式")
    add_card(
        slide,
        0.7,
        1.45,
        3.8,
        4.9,
        "標準情境",
        [
            "idle",
            "ac_only / window_only / light_only",
            "ac_window / window_light / ac_light",
            "all_active",
        ],
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.8,
        4.9,
        "窗戶模式",
        [
            "四季 × 天氣 × 時段",
            "matrix evaluation",
            "也支援 direct outdoor input",
            "可分析窗邊區與中心區差異",
        ],
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        4.9,
        "時間軸",
        [
            "所有 scenario 都有 elapsed time",
            "近似一階動態收斂",
            "Web 端可播放到 quasi-steady state",
        ],
    )
    add_footer(slide, 14)

    # 15 quantitative results
    slide = new_slide(prs)
    add_title(slide, "主要量化結果")
    add_card(
        slide,
        0.7,
        1.4,
        3.6,
        2.0,
        "平均 Field MAE",
        [
            f"Temperature: {avg_mae['temperature']}",
            f"Humidity: {avg_mae['humidity']}",
            f"Illuminance: {avg_mae['illuminance']}",
        ],
    )
    add_card(
        slide,
        0.7,
        3.7,
        3.6,
        2.4,
        "IDW 比較",
        [
            "IDW 可做 baseline",
            "但缺設備位置與方向資訊",
            "照度與局部熱區重建效果明顯較差",
        ],
    )
    add_picture(slide, FIGURES / "submission" / "field_mae_comparison.png", 4.55, 1.25, 7.7, 3.25)
    add_card(
        slide,
        4.6,
        4.75,
        7.7,
        1.35,
        "真實臥室校正檢查",
        [
            f"Raw pillow MAE: {metric_triplet(bedroom_aggregate['raw_pillow_mae'])}",
            f"Corrected pillow MAE: {metric_triplet(bedroom_aggregate['estimated_pillow_mae'])}",
            "推薦動作以實測 penalty 下降作為後續驗證指標",
        ],
    )
    add_footer(slide, 15)

    # real-room and recommendation status
    slide = new_slide(prs)
    add_title(slide, "真實臥室快照與推薦驗證狀態")
    add_card(
        slide,
        0.7,
        1.45,
        5.75,
        4.95,
        "E7：real-bedroom sparse calibration",
        [
            f"Snapshot count: {bedroom_summary['snapshot_count']}",
            "pillow 參考點不參與 8 角點 residual fitting",
            f"Raw pillow MAE: {metric_triplet(bedroom_aggregate['raw_pillow_mae'])}",
            f"Corrected pillow MAE: {metric_triplet(bedroom_aggregate['estimated_pillow_mae'])}",
            "可主張：稀疏校正在此真實快照設定下能改善 held-out pillow 點",
            "不可主張：已完成全室 dense real-room truth 驗證",
        ],
        body_size=12,
    )
    add_card(
        slide,
        6.8,
        1.45,
        5.75,
        4.95,
        "E8：recommendation validation protocol",
        [
            "目前 rank actions 是模型反事實重跑後的 penalty reduction 排序",
            "不可解讀為已完成實際控制有效性的因果驗證",
            "正式驗證應執行 before/after intervention",
            "比較 predicted improvement、actual penalty reduction 與 action rank consistency",
            "這一點是未來工作，也是目前 claim boundary 的重點",
        ],
        body_size=12,
    )
    add_footer(slide, 16)

    # 16 qualitative visual results
    slide = new_slide(prs)
    add_title(slide, "3D 視覺化結果")
    add_picture(slide, FIGURES / "all_active_temperature_3d.svg", 0.65, 1.5, 4.0, 4.8)
    add_picture(slide, FIGURES / "window_only_illuminance_3d.svg", 4.68, 1.5, 4.0, 4.8)
    add_picture(slide, FIGURES / "ac_only_temperature_3d.svg", 8.7, 1.5, 4.0, 4.8)
    add_footer(slide, 16)

    # 17 hybrid result
    slide = new_slide(prs)
    add_title(slide, "Hybrid Residual Neural Network 結果")
    default_hybrid = submission_summary["default_holdout_hybrid"]
    no_fourier = submission_summary["no_fourier_holdout_hybrid"]
    loo = submission_summary["leave_one_scenario_out"]
    add_card(
        slide,
        0.75,
        1.45,
        5.2,
        4.9,
        "Robustness checks",
        [
            f"Default samples: {default_hybrid['dataset']['train_samples']} / {default_hybrid['dataset']['test_samples']}",
            f"Default MAE: {metric_triplet(default_hybrid['hybrid_test_field_mae'])}",
            f"No-Fourier MAE: {metric_triplet(no_fourier['hybrid_test_field_mae'])}",
            f"LOO avg MAE: {metric_triplet(loo['average_hybrid_field_mae'])}",
            f"LOO reduction: T {loo['average_field_mae_reduction_percent']['temperature']:.2f}%, H {loo['average_field_mae_reduction_percent']['humidity']:.2f}%, L {loo['average_field_mae_reduction_percent']['illuminance']:.2f}%",
        ],
    )
    add_picture(slide, FIGURES / "submission" / "field_mae_comparison.png", 6.15, 1.35, 6.1, 3.0)
    add_bullets(
        slide,
        6.25,
        4.55,
        5.8,
        1.7,
        [
            "no-Fourier 對照顯示照度改善不是頻域處理造成",
            "LOO 降低單一 held-out split 過度樂觀的風險",
            "LOO 結果仍限標準情境 family，不等同任意房間泛化",
            "真實臥室快照已驗證 calibration，推薦有效性仍需介入實驗",
        ],
        level0_size=15,
    )
    add_footer(slide, 17)

    # 18 SML2010 public task breakdown
    slide = new_slide(prs)
    add_title(slide, "公開資料任務拆解：SML2010")
    add_picture(slide, PUBLIC_BENCHMARK_FIGURES / "sml2010_task_breakdown.svg", 0.75, 1.3, 7.1, 4.0)
    add_card(
        slide,
        8.15,
        1.4,
        4.35,
        4.85,
        "S1 / S2 / S3 判讀",
        [
            "S1：純照度短視窗，persistence 最強，是劣勢",
            "S2：長視窗溫度有優勢，但濕度有尺度對齊問題",
            "S3：事件 delta response 是主要優勢",
            "SML2010 24 任務：12 lowest MAE、15 勝 LR、14 勝 persistence",
            "不能宣稱公開資料等同 full 3D 場驗證",
        ],
    )
    add_footer(slide, 18)

    # 19 CU-BEMS public task breakdown
    slide = new_slide(prs)
    add_title(slide, "公開資料任務拆解：CU-BEMS")
    add_picture(slide, PUBLIC_BENCHMARK_FIGURES / "cu_bems_task_breakdown.svg", 0.75, 1.3, 7.1, 4.0)
    add_card(
        slide,
        8.15,
        1.4,
        4.35,
        4.85,
        "C1 / C2 / C3 判讀",
        [
            "C1：AC 溫濕度可補強 LR，但不勝 persistence",
            "C2：商辦照度與單房間假設差距大，是劣勢",
            "C3：compound event 可穩定勝過 LR",
            "CU-BEMS 12 任務：9 勝 LR、0 勝 persistence",
            "結果是 zone-level 外部壓力測試，不是 8 點房間拓樸驗證",
        ],
    )
    add_footer(slide, 19)

    # 20 conclusion and future
    slide = new_slide(prs)
    add_title(slide, "結論、限制與未來工作")
    add_card(
        slide,
        0.7,
        1.45,
        3.85,
        4.9,
        "結論",
        [
            "單房間三因子數位孿生原型已可運作",
            "可估場、可校正、可學習、可推薦",
            "bedroom_01 快照顯示校正後 pillow 點 MAE 明顯下降",
            "資料比較需依任務層級切分",
            "hybrid residual 泛化仍需更多房間與 dense ground truth",
            "MCP 與 Web 展示已完成",
        ],
    )
    add_card(
        slide,
        4.75,
        1.45,
        3.85,
        4.9,
        "限制",
        [
            "已有小型真實臥室快照，但仍缺長期 dense field",
            "LOO hybrid 目前只支持標準情境 family 內殘差可學習",
            "不是 CFD 等級模型",
            "公開資料集缺乏 full-field ground truth",
            "推薦動作尚未完成真實介入式因果驗證",
        ],
    )
    add_card(
        slide,
        8.8,
        1.45,
        3.8,
        4.9,
        "未來工作",
        [
            "擴大 ESP32 長期真實資料",
            "加入 CO2 / PM2.5",
            "發展 multi-zone / partition model",
            "把 CU-BEMS / SML2010 / ASHRAE 納入 task-aligned benchmark",
            "執行推薦動作 before/after 介入驗證",
            "朝閉環控制與遠端 MCP 延伸",
        ],
    )
    add_footer(slide, 20)

    # formula section guide
    slide = new_slide(prs)
    add_title(slide, "公式與指標整理")
    add_card(
        slide,
        0.7,
        1.4,
        3.85,
        4.95,
        "場模型",
        [
            "三因子場與查詢點",
            "總估計式",
            "indoor baseline",
            "高度正規化",
            "設備 activation 與 envelope",
        ],
        body_size=13,
    )
    add_card(
        slide,
        4.75,
        1.4,
        3.85,
        4.95,
        "三因子公式",
        [
            "溫度：全室熱響應 + 局部熱場",
            "濕度：除濕 + 外氣水氣交換",
            "照度：直射、環境光、一次漫反射",
            "重點是三種物理量不共用同一套 nominal model",
        ],
        body_size=13,
    )
    add_card(
        slide,
        8.8,
        1.4,
        3.8,
        4.95,
        "校正與評估",
        [
            "8 參數 trilinear correction",
            "可表示空間與平滑 residual 誤差界",
            "非連網裝置影響學習",
            "hybrid residual、MAE/RMSE/correlation、IDW、推薦排序",
        ],
        body_size=13,
    )
    add_footer(slide, 21)
    add_formula_walkthrough(prs, 21, compact=False)
    renumber_footers(prs)
    return prs


def build_outline() -> str:
    slides = [
        ("封面", ["題目、姓名、雙指導教授、研究定位"]),
        ("研究問題與動機", ["非連網裝置無法直接回報狀態", "有限感測器下仍需估計全室環境", "早期純插值與 local-only 模型都不合理"]),
        ("系統架構", ["入口分成使用者互動層與 AI 工具呼叫層", "服務編排、主模型與 residual 修正的分工"]),
        ("房間拓樸、感測器與目標區域", ["8 顆角落感測器", "三個主要區域與三個核心裝置"]),
        ("數學模型", ["變數專屬 nominal model", "trilinear correction", "裝置與家具模組化", "溫度、濕度、照度分別使用不同公式"]),
        ("模型學習、推論與推薦資料流", ["學習端：raw records → 對齊 → scenario state → labels → coefficients/checkpoint", "推論端：runtime input → nominal field → correction / hybrid → point or zone prediction", "推薦端：sample / cluster + T/H/L 目標 → 反事實重跑 → penalty reduction 排序"]),
        ("系統實作與介面", ["MCP 是工具化介面，不是預測模型本身", "initialize：設定 scenario、室內 baseline、外部邊界、設備/家具、預設時間與 estimator", "sample point：查指定座標在特定時間或穩定態的溫濕照度", "learn impacts：start/finish before-after record", "window direct / rank actions：輸入外部窗戶資料；rank actions 需指定 sample 與 T/H/L 目標", "Gemma bridge 與 Web demo 分別負責 AI tool calling 與人機展示"]),
        ("驗證流程與比較原則", ["E1-E3：synthetic full-field、IDW baseline、ablation", "E4：非連網裝置影響學習與推薦排序", "E5：48 組窗戶矩陣與 direct input", "E6：hybrid residual no-Fourier 與 LOO cross-validation", "E7：bedroom_01 7 天真實快照與 pillow hold-out", "E8 protocol、E9 public task-aligned benchmark；demo 不是量化實驗"]),
        ("主要結果", ["平均 field MAE", "IDW / Base / LOO Hybrid 誤差比較", "真實臥室 pillow MAE 比較", "推薦排序目前為 counterfactual simulation", "3D 視覺化案例"]),
        ("Hybrid Residual 結果", ["default held-out、no-Fourier、LOO MAE", "train/test sample count", "研究定位不是黑盒替代", "LOO 結果限標準情境 family"]),
        ("公開資料任務拆解", ["SML2010：S1 純照度劣勢、S2 長視窗溫度部分優勢、S3 事件 delta 主要優勢", "CU-BEMS：C1/C3 勝 linear regression 但不勝 persistence，C2 照度劣勢", "明確說明 public benchmark 不是 full 3D 場驗證"]),
        ("研究貢獻與資料策略", ["三因子、有限感測器、非連網裝置、服務化", "canonical synthetic benchmark + real-bedroom snapshots + task-aligned public datasets", "明確列出每種資料支援的 claim boundary"]),
        ("結論與未來工作", ["長期真實資料、dense real-room ground truth、更多因子、multi-zone、推薦動作介入驗證、閉環控制"]),
    ]
    slides.extend(
        (
            title,
            [
                left_title,
                right_title,
            ],
        )
        for title, left_title, _, right_title, _ in FORMULA_WALKTHROUGH
    )
    lines = ["# 論文報告投影片大綱", ""]
    for index, (title, bullets) in enumerate(slides, start=1):
        lines.append(f"## Slide {index}: {title}")
        lines.extend([f"- {item}" for item in bullets])
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def build_outline_30min() -> str:
    slides = [
        ("封面", ["題目、姓名、雙指導教授、研究定位"]),
        ("報告流程", ["背景、文獻、方法、實作、驗證、結論、公式與指標整理"]),
        ("研究主軸與輸入輸出", ["研究主軸：少量角落感測 + 非連網家電 + 單房間幾何配置 → 三因子空間場估計與決策支援", "輸入：房間、8 點感測、baseline、外部邊界、時間與設備狀態", "模型：三因子 nominal model、power calibration、trilinear correction、hybrid residual", "輸出：任意點/區域估計、3D 視覺化、影響學習、推薦排序與 MCP 查詢"]),
        ("研究背景與問題", ["非連網裝置造成空間影響但無法直接讀取", "有限感測器仍需估全室環境"]),
        ("研究問題與貢獻", ["RQ1-RQ4、主要技術貢獻、task-aligned benchmark 策略"]),
        ("文獻定位、研究缺口與比較原則", ["IEQ 實驗、場重建、hybrid model、digital twin 平台之差異", "公開資料集只比較相容子任務"]),
        ("整體系統架構", ["人機互動層與 AI 工具呼叫層共用服務編排入口"]),
        ("主要執行資料流", ["runtime request 到 dashboard / MCP response 的流程"]),
        ("房間拓樸、感測器與目標區域", ["8 顆角落感測器與三個區域"]),
        ("模組化裝置與家具阻擋", ["裝置模組化、家具自適應阻擋"]),
        ("數學模型", ["變數專屬 nominal model + residual correction", "早期純插值與 local-only 模型失敗後的調整", "避免把同一套公式套用到溫度、濕度、照度"]),
        ("方法選擇：為什麼不是純插值、純物理或純黑盒", ["IDW 適合作 baseline 但缺設備與方向資訊", "完整 CFD/ray tracing 對低成本即時服務太重", "hybrid residual 只學剩餘誤差，不取代可解釋主模型"]),
        ("模型學習、推論與推薦資料流", ["學習資料流：raw data → 對齊 → scenario state → labels → coefficients/checkpoint", "推論資料流：runtime input → nominal field → correction/hybrid → 溫濕照度", "推薦資料流：sample / cluster + T/H/L 目標 → 反事實重跑 → penalty reduction 排序"]),
        ("系統實作與介面", ["MCP 是工具化介面，不是預測模型本身", "initialize：設定 scenario、baseline、外部邊界、設備/家具、時間與 estimator", "sample point：註冊環境後查指定座標三因子估計", "learn impacts：以 before/after observations 建立可學習資料", "window direct / rank actions：直接輸入窗戶外部資料；rank actions 需指定 sample 與 T/H/L 目標", "Gemma/Ollama 透過 bridge 呼叫 tools；Web demo 負責人機互動展示"]),
        ("驗證設計", ["E1-E3：truth-adjusted simulation、IDW、synthetic ablation", "E4-E6：裝置影響學習、window matrix、hybrid no-Fourier/LOO", "E7：bedroom_01 7 天真實快照與 pillow 位置比較", "E8：推薦動作 before/after intervention protocol", "E9：public datasets 僅作 task-aligned benchmark", "Web demo 與 3D 展示是呈現層，不列為量化實驗"]),
        ("證據鏈與 Claim Boundary", ["Synthetic full-field 支援完整 3D 場比較，但不等同長期真實場", "Real-bedroom snapshot 支援稀疏校正的 held-out 點位檢查，但不是 dense truth", "Public datasets 僅支援相容子任務，不是單房間 8 點拓樸驗證", "Recommendation 目前是反事實排序，仍需 before/after 介入驗證"]),
        ("情境設計與輸入模式", ["8 組 scenario、48 組窗戶矩陣、direct input、timeline"]),
        ("主要量化結果", ["平均 MAE、IDW/Base/LOO Hybrid 誤差圖", "真實臥室 raw vs corrected pillow MAE", "推薦有效性以 actual comfort-penalty reduction 驗證", "實驗 E1-E7 與 E9 已有數值輸出；E8 僅為介入 protocol"]),
        ("真實臥室快照與推薦驗證狀態", ["E7：pillow hold-out 不參與 8 角點 residual fitting，呈現 raw vs corrected MAE", "E8：rank actions 目前是模型反事實排序，需實測介入驗證因果效果"]),
        ("3D 視覺化結果", ["溫度與照度熱區案例"]),
        ("Hybrid Residual 結果", ["default held-out、no-Fourier、LOO robustness checks", "train/test sample count 與 synthetic benchmark 限制", "LOO 結果限標準情境 family", "真實快照作為 sparse calibration 驗證"]),
        ("公開資料任務拆解：SML2010", ["S1：純照度短視窗是劣勢", "S2：長視窗溫度有優勢但濕度有尺度對齊問題", "S3：事件 delta response 是主要優勢"]),
        ("公開資料任務拆解：CU-BEMS", ["C1：AC 溫濕度可補強 linear regression", "C2：商辦照度與單房間假設差距大", "C3：compound event 可勝 linear regression 但不勝 persistence"]),
        ("結論、限制與未來工作", ["目前完成度、真實快照限制、hybrid 泛化限制、推薦動作尚需介入驗證、task-aligned benchmark 與後續方向"]),
        ("公式與指標整理", ["場模型：三因子場、總估計式、baseline、activation、envelope", "三因子公式：溫度、濕度、照度分別說明", "校正與評估：8 點三線性校正、影響學習、hybrid residual、metrics、IDW、推薦排序"]),
    ]
    slides.extend(
        (
            title,
            [
                left_title,
                right_title,
            ],
        )
        for title, left_title, _, right_title, _ in FORMULA_WALKTHROUGH
    )
    lines = ["# 論文報告投影片大綱（30min 版）", ""]
    for index, (title, bullets) in enumerate(slides, start=1):
        lines.append(f"## Slide {index}: {title}")
        lines.extend([f"- {item}" for item in bullets])
        lines.append("")
    return "\n".join(lines).strip() + "\n"


SPEAKER_NOTE_GLOSSARY: List[Tuple[str, Tuple[str, ...], str]] = [
    ("單房間", ("單房間", "single-room"), "本研究限定在單一矩形房間，不處理多房間或整棟建築的氣流與能量交換。"),
    ("非連網家電/裝置", ("非連網家電", "非連網裝置", "非連網"), "沒有穩定 API 或遙測資料可讀取狀態的冷氣、窗戶、照明等設備。"),
    ("稀疏感測", ("稀疏感測", "少量感測", "sparse"), "感測點數量少於完整空間場需求，需靠模型與校正推估未量測位置。"),
    ("角落感測器", ("角落感測器", "8 顆", "8 點", "8個角點", "8 個角點"), "配置在房間地面四角與天花板四角的 8 個感測點，用於建立 sparse observation 與 residual correction。"),
    ("空間數位孿生", ("空間數位孿生", "數位孿生", "digital twin"), "以房間幾何、裝置、感測器與模型維持一個可查詢的室內環境狀態估計。"),
    ("三因子", ("三因子", "temperature", "humidity", "illuminance", "溫度、濕度、照度", "溫濕照度"), "本研究同時估計溫度、相對濕度與照度三種室內環境量。"),
    ("空間場", ("空間場", "field", "場估計", "場重建", "field MAE"), "不是單一平均值，而是在房間 3D 座標中任意位置可查詢的環境量分布。"),
    ("控制導向", ("控制導向", "control-oriented"), "模型重點在支援查詢、比較與推薦排序，而不是取代高精度物理模擬器。"),
    ("CFD", ("CFD",), "Computational Fluid Dynamics，計算流體力學；可模擬細緻氣流，但邊界條件與計算成本高。"),
    ("Ray tracing", ("ray tracing", "光線追蹤"), "依光線路徑追蹤照明傳播的精密光學方法；本研究只採輕量照度幾何與一次漫反射近似。"),
    ("API", ("API",), "Application Programming Interface，讓系統讀取或控制設備狀態的程式介面。"),
    ("遙測", ("遙測", "telemetry"), "設備主動回報狀態或感測資料；非連網裝置通常缺少這類資料。"),
    ("Sample", ("sample point", "sample ", "sample/cluster", "sample 或 zone"), "指定房間中的查詢點，用來取得該座標的三因子估計。"),
    ("Zone", ("zone", "區域", "目標區域"), "房間內的目標區域，用於彙整多個點的平均狀態或舒適度評估。"),
    ("Baseline", ("baseline", "Indoor baseline", "室內 baseline"), "泛指比較或模型參考基準；在模型脈絡中常指未加入設備作用前的室內溫濕照度基準。"),
    ("外部邊界", ("外部邊界", "外部環境邊界", "outdoor"), "室外溫度、濕度、日照等會透過窗戶或邊界條件影響室內的輸入。"),
    ("Nominal model", ("nominal model", "N_v", "Nᵥ"), "主模型的可解釋估計部分，描述設備、邊界與空間結構造成的主要趨勢。"),
    ("Residual", ("residual", "剩餘誤差", "殘差"), "觀測或 truth 與模型預測之間的差，用於校正或第二層學習。"),
    ("Residual correction", ("residual correction", "校正場", "感測器校正"), "利用感測器 residual 修正 nominal model，使估計更貼近觀測。"),
    ("Trilinear correction", ("trilinear", "三線性", "三線性補間", "三線性校正"), "用 X/Y/Z 三個座標方向的一階補間，由 8 個角點 residual 推估室內 residual 場。"),
    ("Power calibration", ("power calibration", "power scale", "P_ac", "P_win", "P_light"), "依觀測差異調整設備影響強度，避免裝置作用尺度只依預設值決定。"),
    ("Hybrid residual", ("hybrid residual", "F_hybrid", "hybrid"), "在可解釋 base estimator 後面再加一個資料驅動 residual 模型，不直接取代主模型。"),
    ("MCP", ("MCP", "Model Context Protocol"), "Model Context Protocol，是讓 LLM application 以標準化方式連接外部資料與工具的 open protocol；本研究用它封裝數位孿生工具。"),
    ("MCP host/client/server", ("MCP host", "MCP client", "MCP server", "client-server", "server"), "MCP 採 client-server 概念；host/client 是使用工具的 AI 應用端，server 則暴露工具、資源或 prompt。"),
    ("MCP Tools", ("MCP tools", "tools/list", "tools/call", "initialize_environment", "sample_point", "learn_impacts", "rank_actions"), "MCP server 可暴露可執行工具；client 可列出工具並以結構化 arguments 呼叫。"),
    ("MCP Resources", ("resources/list", "Resources", "resources"), "MCP resources 是 server 提供給 client 的上下文資料，例如檔案、資料庫 schema 或應用資料。"),
    ("MCP Prompts", ("prompts/list", "Prompts", "prompt templates"), "MCP prompts 是 server 提供的結構化 prompt template，可由 client 取得並填入參數。"),
    ("JSON-RPC", ("JSON-RPC", "jsonrpc"), "MCP 使用 JSON-RPC 2.0 編碼 request、response 與 notification。"),
    ("stdio transport", ("stdio", "stdin", "stdout", "standard input", "standard output"), "MCP 的本地 transport 之一；client 啟動 server subprocess，透過 stdin/stdout 傳送 UTF-8 JSON-RPC 訊息。"),
    ("Streamable HTTP", ("Streamable HTTP", "HTTP MCP", "遠端 HTTP MCP"), "MCP 的另一種標準 transport，適合遠端或網路化部署。"),
    ("Protocol version", ("protocolVersion", "2024-11-05", "2025-06-18"), "MCP 會在 initialize 階段協商 protocolVersion；本研究本地 server 目前回傳既有版本並實作 tools workflow。"),
    ("Web demo", ("Web demo", "dashboard"), "人機互動展示介面，用於查看 3D 場、時間軸、設備狀態與查詢結果。"),
    ("Gemma/Ollama bridge", ("Gemma", "Ollama", "bridge"), "讓本地語言模型透過工具呼叫流程存取模型服務的橋接層。"),
    ("Tool calling", ("tool calling", "AI 工具呼叫", "AI agent"), "語言模型不是直接計算答案，而是呼叫外部工具取得模型查詢或操作結果。"),
    ("服務編排", ("服務編排", "service", "service orchestration"), "把 scenario、模型估計、校正、推薦與輸出流程串接起來的中介層。"),
    ("Estimator", ("estimator",), "實際負責產生場估計的模型物件，可切換 base、corrected 或 hybrid 版本。"),
    ("Scenario", ("scenario", "標準情境", "情境"), "一組房間、設備、外部邊界與時間設定，用於模擬或驗證。"),
    ("Direct input", ("direct input", "window direct"), "不使用預設矩陣情境，直接輸入外部溫濕度、日照與開窗比例。"),
    ("Window matrix", ("window matrix", "窗戶矩陣"), "依季節、天氣、時段等組合建立的窗戶外部邊界情境集合。"),
    ("Quasi-steady state", ("quasi-steady", "穩定態"), "近似達到穩定但非嚴格物理穩態的狀態，用於簡化時間響應解讀。"),
    ("Activation", ("activation", "Aⱼ", "A_j"), "設備啟動強度隨時間接近穩態的函數。"),
    ("Influence envelope", ("influence envelope", "Eⱼ", "E_j", "envelope"), "設備對某位置的空間作用權重，通常含時間強度、距離、方向與遮蔽。"),
    ("Distance decay", ("距離衰減", "Rⱼ", "R_j"), "距離設備越遠，局部作用越弱的權重函數。"),
    ("Directionality", ("方向性", "Dⱼ", "D_j"), "冷氣出風方向、窗戶日照方向或光源方向造成的非均向影響。"),
    ("Visibility/obstruction", ("可見性", "遮蔽", "阻擋", "Vⱼ", "V_j", "obstruction"), "家具或幾何遮擋造成設備影響變弱的因素。"),
    ("IDW", ("IDW",), "Inverse Distance Weighting，反距離加權插值；只使用距離與感測值，不含設備物理先驗。"),
    ("黑盒模型", ("黑盒", "純黑盒"), "主要依資料學習輸入輸出關係、但內部物理意義較不明確的模型。"),
    ("IEQ", ("IEQ",), "Indoor Environmental Quality，室內環境品質，通常涵蓋熱舒適、空氣品質、照明等因素。"),
    ("Hybrid thermal model", ("hybrid thermal model",), "結合物理結構與資料驅動方法的熱環境模型。"),
    ("Task-aligned benchmark", ("task-aligned benchmark", "相容子任務", "task-aligned"), "只比較公開資料集中與本研究觀測型態相容的子任務，不宣稱完整場驗證。"),
    ("SML2010", ("SML2010",), "公開智慧建築資料集；本研究用於 two-point boundary-response 類任務。"),
    ("CU-BEMS", ("CU-BEMS",), "商辦建築能源管理資料集；本研究用於 zone-level device-response 類任務。"),
    ("Public dataset", ("Public datasets", "公開資料", "公開資料集"), "外部公開資料來源；本研究只用於相容任務壓力測試，不作完整 3D truth。"),
    ("Synthetic full-field", ("Synthetic full-field", "truth-adjusted", "受控完整場", "synthetic"), "可取得完整場 truth 的受控驗證資料，用於完整 3D 場誤差比較。"),
    ("Ablation", ("ablation", "消融"), "移除或替換模型元件後比較性能，用來檢查各元件貢獻。"),
    ("LOO", ("LOO", "leave-one-scenario-out"), "Leave-one-scenario-out，輪流留下一個情境作測試，檢查模型是否只對單一切分有效。"),
    ("No-Fourier", ("No-Fourier", "Fourier"), "去除或比較 Fourier 相關處理的對照設定，用於確認改善不是單一頻域技巧造成。"),
    ("Real-bedroom snapshot", ("Real-bedroom snapshot", "bedroom_01", "真實臥室", "快照"), "真實臥室中的稀疏量測快照，用於檢查校正對未參與 fitting 點位的改善。"),
    ("Pillow hold-out", ("pillow", "hold-out"), "將 pillow 位置作為未參與校正 fitting 的參考點，用於測試非感測點估計效果。"),
    ("Dense ground truth", ("dense", "dense truth", "dense real-room ground truth", "完整 3D 場真值"), "房間內大量點位的真實環境場資料，是更嚴格但較難取得的驗證基準。"),
    ("MAE", ("MAE",), "Mean Absolute Error，平均絕對誤差；數值越低代表平均偏差越小。"),
    ("RMSE", ("RMSE",), "Root Mean Squared Error，均方根誤差；比 MAE 更放大尖峰或離群誤差。"),
    ("Correlation", ("Correlation", "相關係數"), "衡量預測與真值趨勢方向一致性的指標。"),
    ("Persistence", ("persistence",), "直接沿用上一時步值作預測的時間序列 baseline，在高慣性短視窗資料中通常很強。"),
    ("Linear regression", ("linear regression", "LR"), "線性回歸 baseline，用線性權重將輸入特徵映射到目標值。"),
    ("Structured prior", ("structured prior",), "模型內建的設備、邊界與物理結構先驗。"),
    ("Facade event delta", ("facade event delta", "event delta", "delta response"), "檢查外部邊界或事件造成的變化量，而非只預測下一時步絕對值。"),
    ("Device-response benchmark", ("device-response", "zone-level device-response"), "以設備用電或啟動訊號對 zone 環境變化的響應作為比較任務。"),
    ("Zone-level", ("zone-level",), "以建築區域平均值為資料粒度，不含房間內細緻 3D 幾何。"),
    ("Counterfactual simulation", ("counterfactual", "反事實"), "假設某候選動作發生後重新估計結果，用來比較預期改善。"),
    ("Comfort penalty", ("comfort penalty", "Penalty", "舒適度"), "偏離目標溫濕照度時的懲罰值，推薦排序用它衡量改善幅度。"),
    ("Before/after intervention", ("before/after", "intervention", "介入驗證"), "實際採取動作前後量測環境變化，用於驗證推薦是否有因果改善效果。"),
    ("ESP32", ("ESP32",), "低成本微控制器平台，可用於後續長期真實感測資料蒐集。"),
    ("CO2", ("CO2", "CO₂"), "二氧化碳濃度，可作為未來室內空氣品質因子。"),
    ("PM2.5", ("PM2.5",), "細懸浮微粒濃度，可作為未來室內空氣品質因子。"),
    ("Multi-zone model", ("multi-zone", "partition model"), "將房間或建築切成多個區域處理交換與隔間效應的模型。"),
    ("閉環控制", ("閉環控制", "closed-loop"), "模型輸出進一步驅動控制動作，並用後續感測結果回饋修正決策。"),
    ("座標系", ("座標系", "原點", "p = (x,y,z)", "p=(x,y,z)"), "用 x/y/z 公尺座標描述房間內位置；本研究原點在地面西南角。"),
    ("T(p,t)", ("T(p,t)", "T："), "位置 p、時間 t 的溫度場值。"),
    ("H(p,t)", ("H(p,t)", "H："), "位置 p、時間 t 的相對濕度場值。"),
    ("L(p,t)", ("L(p,t)", "L："), "位置 p、時間 t 的照度場值。"),
    ("b₀", ("b₀", "T₀", "H₀", "L₀"), "室內基準狀態，包含基準溫度、相對濕度與照度。"),
    ("ζ", ("ζ",), "高度正規化座標，用於描述查詢點相對於房間高度的位置。"),
    ("τⱼ", ("τⱼ", "τ_j"), "設備時間響應常數，控制 activation 接近穩態的速度。"),
    ("B 項", ("B_T", "B_H", "B_ac", "B_win", "B_light", "全室項"), "表示全室平均狀態偏移的 bulk/global effect。"),
    ("S 項", ("S_T", "S_H", "S_ac", "S_win", "S_light", "局部項"), "表示設備附近、窗邊或光源附近的局部空間差異。"),
    ("clip[0,100]", ("clip[0,100]",), "把相對濕度限制在 0% 到 100% 的合理物理範圍內。"),
    ("max{0}", ("max{0", "max{0,"), "把照度限制為非負值，避免模型輸出不合理的負照度。"),
    ("直射光", ("直射光", "direct source", "L_winᵈ", "L_lightᵈ"), "直接由窗戶或燈具到達查詢點的照度貢獻。"),
    ("環境光", ("環境光", "ambient"), "非單一路徑直射、較均勻分布的背景照度貢獻。"),
    ("一次漫反射", ("一次漫反射", "Iʳ", "reflect", "漫反射"), "只計算一次表面反射對照度的回填效果，是輕量近似而非完整 radiosity。"),
    ("反射率 ρ", ("ρ_s", "ρ"), "表面反射率，描述牆面、地板或家具把入射光反射出去的比例。"),
    ("三線性函數空間", ("𝒱", "span{1", "XYZ"), "由 1、X、Y、Z 與交互項組成的 8 維 residual 表示空間。"),
    ("角點 residual", ("rᵛ", "r^", "rᵛ_", "p_{abc}"), "角落感測器觀測值與 nominal model 預測值的差。"),
    ("補間權重 ℓ", ("ℓ₀", "ℓ₁"), "三線性補間中每個角點 residual 對內部點的權重函數。"),
    ("誤差上界", ("誤差上界", "M_xx", "M_yy", "M_zz"), "用 residual 二階曲率限制三線性補間與真實 residual 的最大偏差。"),
    ("特徵向量 φᵢ", ("φᵢ", "feature vector"), "模型訓練時輸入的情境特徵，例如座標、時間、baseline、外部條件與設備作用。"),
    ("標籤 yᵢ", ("yᵢ", "label", "labels"), "監督式學習中的目標值；本研究常用 true field 與 base estimator 的差作 residual label。"),
    ("損失函數 ℒ", ("ℒ", "loss"), "訓練 neural network 時要最小化的目標函數。"),
    ("正則化 λ", ("λ", "regularization", "正則化"), "限制模型參數大小以降低過擬合的項。"),
    ("IDW 權重", ("w_s", "q：距離權重"), "IDW 中由距離決定的感測器權重，距離越近權重越高。"),
    ("Score(a)", ("Score(a)", "Score"), "候選動作 a 的推薦分數，通常由採取前後 comfort penalty 的下降量決定。"),
]


def glossary_notes_for_slide(title: str, paragraphs: Sequence[str]) -> List[Tuple[str, str]]:
    text = "\n".join([title, *paragraphs]).lower()
    notes: List[Tuple[str, str]] = []
    for label, aliases, explanation in SPEAKER_NOTE_GLOSSARY:
        if any(alias.lower() in text for alias in aliases):
            notes.append((label, explanation))
    return notes


def build_speaker_notes_30min() -> str:
    validation_summary = read_json(DATA / "validation_summary.json")
    submission_summary = read_json(DATA / "submission_readiness_summary.json")
    window_summary = read_json(DATA / "window_matrix_summary.json")
    bedroom_summary = read_json(DATA / "bedroom_01_weekly" / "weekly_simulation_summary.json")
    avg_mae = average_field_mae(validation_summary)
    bedroom_aggregate = bedroom_summary["aggregate"]
    default_hybrid = submission_summary["default_holdout_hybrid"]
    no_fourier = submission_summary["no_fourier_holdout_hybrid"]
    loo = submission_summary["leave_one_scenario_out"]

    slides: List[Tuple[str, List[str]]] = [
        (
            "封面",
            [
                "各位老師好，我是林昀佑。今天報告的題目是「單房間非連網家電環境影響學習之稀疏感測空間數位孿生原型」。這個題目聚焦在一般房間中常見但不一定連網的冷氣、窗戶與照明，以及少量感測器下如何估計完整室內環境分布。",
                "整體研究不是要做完整 CFD 或精密光學模擬，而是建立一個控制導向、可解釋、可校正、也能被工具介面查詢的三因子空間數位孿生原型。",
            ],
        ),
        (
            "報告流程",
            [
                "接下來會先從研究背景與問題開始，說明為什麼非連網家電與有限感測器會造成空間感知困難。",
                "再來會說明文獻定位、系統架構、數學模型、感測器校正與影響學習，然後進入系統實作、驗證設計與實驗結果。",
                "最後整理結論、限制與未來工作；後半段的公式與指標整理可以用來補充每個模型元件的細節。",
            ],
        ),
        (
            "研究主軸與輸入輸出",
            [
                "這頁先把整個研究壓縮成輸入、模型與輸出三個部分。輸入端包含房間幾何、8 顆角落感測器、室內 baseline、外部環境邊界、時間與設備狀態。",
                "模型端的核心是三因子 nominal model，再加上 power calibration、trilinear residual correction，以及 optional 的 hybrid residual 修正。",
                "輸出端則包含任意點或區域的溫度、濕度、照度估計，3D 視覺化，非連網裝置影響係數，以及後續可以提供 MCP 或 Web demo 使用的反事實推薦排序。",
            ],
        ),
        (
            "研究背景與問題",
            [
                "一般智慧居家或智慧建築需要知道室內環境狀態，才能支援舒適度評估、能源管理與設備控制。但實際房間中，冷氣、窗戶和照明常常沒有可讀取的 API 或遙測資料。",
                "另一個限制是感測器數量。使用者關心的是整個房間不同位置的舒適狀態，但實際上通常只會放少數幾顆感測器。",
                "因此本研究的問題是：在非連網裝置狀態不完整、感測器稀疏的情況下，如何估計整個房間的環境分布，並學習這些裝置對環境的影響。",
            ],
        ),
        (
            "研究問題與貢獻",
            [
                "本研究可以拆成四個研究問題。第一，8 顆角落感測器能不能支援單房間三因子空間場估計。第二，能不能從環境變化中學習非連網裝置影響。",
                "第三，當使用者指定 sample 或 zone 以及溫濕照度目標時，能不能排序可能的控制動作。第四，這個模型能不能封裝成 Web 與 MCP 可查詢服務。",
                "主要貢獻是把三因子 nominal model、8 點 residual correction、非連網裝置影響學習、hybrid residual，以及 task-aligned public benchmark 放在同一個可執行原型中。",
            ],
        ),
        (
            "文獻定位、研究缺口與比較原則",
            [
                "相關研究大致可以分成 IEQ 實驗、有限感測器場重建、hybrid thermal model 與 digital twin 平台。但很多研究只處理單一或雙因子，或依賴較完整的設備遙測。",
                "本研究的差異在於把單房間、低成本角落感測、非連網裝置，以及溫度、濕度、照度三因子放在同一個控制導向模型中。",
                "公開資料集的比較採取 task-aligned 原則，因為 SML2010 和 CU-BEMS 沒有本研究需要的單房間幾何、8 點拓樸和 dense 3D 場真值，所以只能比較相容子任務。",
            ],
        ),
        (
            "整體系統架構",
            [
                "整體架構分為前端互動、AI 工具呼叫、服務編排、數位孿生核心、校正學習與視覺化輸出幾層。",
                "Web demo 與 MCP tools 都不直接做模型推論，而是呼叫同一個服務編排入口，由核心模型負責場估計、校正、學習與推薦排序。",
                "這樣設計的好處是展示介面、AI tool calling 與後端模型可以分工，未來也比較容易替換 estimator 或新增裝置與環境因子。",
            ],
        ),
        (
            "主要執行資料流",
            [
                "執行時，系統先取得 scenario 或 direct input，包含房間狀態、baseline、外部邊界、家具和設備設定。",
                "接著服務層把這些資料交給 estimator，先建立 nominal field，再套用感測器校正或 hybrid residual，最後輸出 dashboard、point sample、zone summary 或 MCP response。",
                "所以 Web 和 MCP 只是不同入口，核心資料流是共用的，避免展示結果和工具查詢結果不一致。",
            ],
        ),
        (
            "房間拓樸、感測器與目標區域",
            [
                "本研究使用單一矩形房間作為主要研究場景，尺寸為 6 m × 4 m × 3 m。座標系使用公尺，原點在房間地面西南角。",
                "感測器放在地面四角與天花板四角，共 8 顆。這 8 點不是直接量到全室，而是提供 sparse observation，用來對 nominal model 的 residual 做三線性補間。",
                "目標區域分成窗邊、中心與門側等 zone，方便後續做區域平均、舒適度評估和推薦排序。",
            ],
        ),
        (
            "模組化裝置與家具阻擋",
            [
                "冷氣、窗戶與燈具都被視為模組化裝置，每個裝置都有位置、方向、作用尺度與啟動狀態。這讓模型可以支援新增或移動裝置。",
                "家具則被視為空間中的阻擋物，會影響冷氣、窗戶日照或燈具光源的局部作用權重。",
                "這個設計的目的不是做非常精細的流場或光線追蹤，而是用低成本幾何資訊修正單純距離衰減太粗略的問題。",
            ],
        ),
        (
            "數學模型",
            [
                "核心估計式是 F_hat_v(p,t)=N_v(p,t)+C_v(p,t)。N_v 是變數專屬 nominal model，C_v 是由角落感測 residual 建立的校正場。",
                "溫度模型處理熱交換、熱源和垂直分層；濕度模型處理除濕與外氣水氣交換；照度模型處理光源幾何、遮蔽與一次漫反射。",
                "這裡最重要的是三個環境變數不共用同一套物理公式。它們共用座標、裝置框架與校正流程，但 nominal model 根據物理意義分開設計。",
            ],
        ),
        (
            "方法選擇：為什麼不是純插值、純物理或純黑盒",
            [
                "純插值，例如 IDW，只知道感測器位置與距離，不知道冷氣出風、窗戶日照或燈具位置，因此在局部熱區和照度場會比較吃虧。",
                "完整物理模擬像 CFD 或 ray tracing，需要大量邊界條件、材料與計算成本，不符合低成本房間原型與即時查詢的需求。",
                "純黑盒模型在資料量有限時也容易過擬合，而且不容易解釋設備與空間結構的作用。因此本研究採用可解釋 base model，再用 residual correction 和 hybrid residual 補足誤差。",
            ],
        ),
        (
            "模型學習、推論與推薦資料流",
            [
                "學習端會把 raw records 對齊成 scenario state，產生訓練 labels，再更新裝置影響係數或 hybrid checkpoint。",
                "推論端從 runtime input 開始，先建立 nominal field，再做 correction 或 hybrid 修正，最後輸出 point 或 zone 的三因子估計。",
                "推薦端不是直接控制設備，而是把候選動作做反事實重跑，計算採取動作前後的 comfort penalty reduction，再依改善幅度排序。",
            ],
        ),
        (
            "系統實作與介面",
            [
                "MCP 的全名是 Model Context Protocol。它不是我的預測模型，也不是一個新的神經網路架構，而是一個讓 LLM application 用標準化方式連接外部資料與工具的 open protocol。",
                "如果老師問 std 或 standard，我會回答：MCP 本身是標準化的 protocol；官方規格用 JSON-RPC 2.0 表示 request、response 與 notification。它的標準 transport 包含 stdio 和 Streamable HTTP。",
                "stdio 是 standard input/output 的意思，適合本機工具。client 會啟動 MCP server subprocess，server 從 stdin 讀 newline-delimited JSON-RPC message，再把 response 寫到 stdout；stderr 只用於 log。",
                "在我的系統裡，數位孿生核心服務被包成本地 MCP server，主要暴露 tools/list 與 tools/call。工具包含 initialize_environment、sample_point、learn_impacts、run_window_direct 和 rank_actions。",
                "initialize 負責註冊 scenario、baseline、外部邊界、設備、家具、時間與 estimator。sample_point 查詢指定座標的溫濕照度估計；rank_actions 則在給定 sample 與三因子目標後做反事實排序。",
                "所以本研究對 MCP 的定位是系統整合與工具化封裝：證明這個數位孿生模型可以被 AI client 操作。我的研究貢獻不是提出新的 MCP protocol，也不是宣稱模型權重原生支援 MCP。",
                "Web demo 負責人機互動展示，Gemma/Ollama bridge 負責把自然語言轉成 tool calling。兩者底層都呼叫同一個模型服務，因此結果可以保持一致。",
            ],
        ),
        (
            "驗證設計",
            [
                "驗證採分層設計。E1 到 E3 是受控完整場重建、IDW baseline 與 ablation，主要檢查完整 3D 場估計。",
                f"E4 到 E6 包含非連網裝置影響學習、{window_summary.get('count', 0)} 組窗戶矩陣，以及 hybrid no-Fourier 和 leave-one-scenario-out 檢查。",
                f"E7 使用 bedroom_01 的 {bedroom_summary['snapshot_count']} 筆真實快照做 pillow hold-out 檢查。E8 是推薦動作介入驗證 protocol，E9 是公開資料集 task-aligned benchmark。",
            ],
        ),
        (
            "證據鏈與 Claim Boundary",
            [
                "Synthetic full-field 支援完整 3D 場誤差比較，因此可以用來比較 base model、IDW、ablation 與 hybrid residual。",
                "Real-bedroom snapshot 支援真實稀疏校正檢查，尤其是 pillow hold-out 點，但它不是 dense real-room ground truth。",
                "Public datasets 只支援相容子任務，不能被解讀為單房間 8 點拓樸或完整 3D 場驗證。Recommendation 目前也只是反事實排序，尚需真實 before/after intervention 證明因果改善。",
            ],
        ),
        (
            "情境設計與輸入模式",
            [
                "標準情境包含 idle、單裝置啟動、雙裝置組合與 all_active，用來觀察不同裝置組合對三因子場的影響。",
                "窗戶部分除了 48 組季節、天氣與時段矩陣，也支援 direct input，讓使用者直接指定外部溫度、濕度、日照與開窗比例。",
                "所有 scenario 都有時間軸設定，設備影響採一階收斂近似，因此可以看啟動後逐步接近 quasi-steady state 的過程。",
            ],
        ),
        (
            "主要量化結果",
            [
                f"8 組標準情境中，base model 平均 field MAE 為 temperature {avg_mae['temperature']:.4f}、humidity {avg_mae['humidity']:.4f}、illuminance {avg_mae['illuminance']:.4f}。",
                "圖中比較 IDW、base model 與 leave-one-scenario-out hybrid residual。IDW 因缺少裝置位置、方向與物理先驗，在照度與局部場上特別不利。",
                f"真實臥室 pillow 點的 raw MAE 為 {metric_triplet(bedroom_aggregate['raw_pillow_mae'])}，校正後 MAE 為 {metric_triplet(bedroom_aggregate['estimated_pillow_mae'])}，顯示稀疏校正在此設定下有明顯改善。",
            ],
        ),
        (
            "真實臥室快照與推薦驗證狀態",
            [
                "E7 的重點是 pillow 參考點沒有參與 8 個角點 residual fitting，因此它可以用來檢查校正場是否改善非感測點估計。",
                f"結果上，校正後 pillow MAE 從 raw 的 {metric_triplet(bedroom_aggregate['raw_pillow_mae'])} 降到 {metric_triplet(bedroom_aggregate['estimated_pillow_mae'])}。",
                "E8 目前是推薦動作驗證 protocol。也就是說，本研究已能做模型反事實排序，但實際控制行為是否真的改善舒適度，仍需要 before/after intervention 來驗證。",
            ],
        ),
        (
            "3D 視覺化結果",
            [
                "這頁展示不同情境下的 3D 場分布，例如 all_active 的溫度場、window_only 的照度場，以及 ac_only 的溫度場。",
                "3D 圖的功能是幫助理解空間分布，不直接作為新的量化實驗。量化結果仍以前面提到的 field MAE、baseline comparison 和 hold-out 檢查為主。",
                "可以看到裝置位置與作用方向會造成局部差異，這也是為什麼模型不能只用單一平均值或純距離插值處理。",
            ],
        ),
        (
            "Hybrid Residual 結果",
            [
                f"Default held-out 的 train/test samples 為 {default_hybrid['dataset']['train_samples']} / {default_hybrid['dataset']['test_samples']}，hybrid test MAE 為 {metric_triplet(default_hybrid['hybrid_test_field_mae'])}。",
                f"No-Fourier 對照的 MAE 為 {metric_triplet(no_fourier['hybrid_test_field_mae'])}，LOO 平均 hybrid MAE 為 {metric_triplet(loo['average_hybrid_field_mae'])}。",
                "這些結果表示標準情境 family 內的 residual 有可學習性，但不能直接擴大解讀為任意房間、任意裝置配置都能泛化。",
            ],
        ),
        (
            "公開資料任務拆解：SML2010",
            [
                "SML2010 被映射成 two-point boundary-response benchmark。它適合檢查外氣、日照與室內兩點響應，但沒有完整 3D 場真值。",
                "S1 純照度短視窗是主要劣勢，因為 persistence 在短時間照度高度自相關時很強。S2 長視窗溫度有部分優勢，但濕度有尺度對齊問題。",
                "S3 facade event delta 是主要優勢，因為事件後變化方向和長視窗響應更能受益於 structured prior。",
            ],
        ),
        (
            "公開資料任務拆解：CU-BEMS",
            [
                "CU-BEMS 被映射成商辦 zone-level device-response benchmark。它有 AC power 和 lighting power 等欄位，但不是本研究的單房間 8 點拓樸。",
                "C1 中 AC 溫濕度可補強 linear regression，但不勝 persistence。C2 商辦照度與單房間光學假設差距大，是明確劣勢。",
                "C3 compound event 可勝 linear regression，但仍不勝 persistence。這表示本研究特徵對事件讀出有幫助，但不能宣稱在商辦時序任務全面勝出。",
            ],
        ),
        (
            "結論、限制與未來工作",
            [
                "本研究完成一個單房間三因子空間數位孿生原型，能在少量角落感測器下估計溫度、濕度與照度分布，並支援非連網裝置影響學習。",
                "限制方面，目前仍缺長期 dense real-room ground truth，hybrid residual 的泛化也主要限於標準情境 family。",
                "未來工作包括擴大 ESP32 長期資料、加入 CO2/PM2.5、發展 multi-zone model、執行推薦動作介入驗證，以及往閉環控制與遠端 MCP 延伸。",
            ],
        ),
        (
            "公式與指標整理",
            [
                "後半段整理公式與指標。第一組是場模型，包括三因子場、總估計式、baseline、activation 與 influence envelope。",
                "第二組是三因子 nominal model，分別說明溫度、濕度與照度為什麼要採用不同的物理近似。",
                "第三組是校正與評估，包括 8 點三線性 residual correction、非連網裝置影響學習、hybrid residual、MAE/RMSE/correlation、IDW baseline 與推薦排序。",
            ],
        ),
    ]

    for title, left_title, left_lines, right_title, right_lines in FORMULA_WALKTHROUGH:
        left_text = "；".join(left_lines[:3])
        right_text = "；".join(right_lines[:3])
        slides.append(
            (
                title,
                [
                    f"這頁說明「{left_title}」。可以先從公式或定義開始，指出它在整體模型中負責哪一部分。",
                    f"左側重點包含：{left_text}。",
                    f"接著說明「{right_title}」。這一部分通常用來補上模型設計理由、限制或可主張範圍。",
                    f"右側重點包含：{right_text}。",
                ],
            )
        )

    lines = [
        "# 30 分鐘論文簡報逐頁講稿",
        "",
        "本檔是 `thesis_presentation_zh_30min.pptx` 的講稿，不放入投影片畫面。投影片維持正式內容；這份 Markdown 用於練習口頭說明與答辯準備。",
        "",
    ]
    for index, (title, paragraphs) in enumerate(slides, start=1):
        lines.append(f"## Slide {index}: {title}")
        for paragraph in paragraphs:
            lines.append("")
            lines.append(paragraph)
        glossary_notes = glossary_notes_for_slide(title, paragraphs)
        if glossary_notes:
            lines.append("")
            lines.append("### 名詞註釋")
            for label, explanation in glossary_notes:
                lines.append(f"- **{label}**：{explanation}")
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def main() -> None:
    PAPERS.mkdir(parents=True, exist_ok=True)
    THESIS_PAPERS.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(PRESENTATION_PATH)
    prs.save(STORED_PRESENTATION_PATH)
    prs_long = build_presentation_30min()
    prs_long.save(LONG_PRESENTATION_PATH)
    prs_long.save(STORED_LONG_PRESENTATION_PATH)
    OUTLINE_PATH.write_text(build_outline(), encoding="utf-8")
    LONG_OUTLINE_PATH.write_text(build_outline_30min(), encoding="utf-8")
    LONG_SPEAKER_NOTES_PATH.write_text(build_speaker_notes_30min(), encoding="utf-8")
    print(f"Wrote {PRESENTATION_PATH}")
    print(f"Wrote {STORED_PRESENTATION_PATH}")
    print(f"Wrote {LONG_PRESENTATION_PATH}")
    print(f"Wrote {STORED_LONG_PRESENTATION_PATH}")
    print(f"Wrote {OUTLINE_PATH}")
    print(f"Wrote {LONG_OUTLINE_PATH}")
    print(f"Wrote {LONG_SPEAKER_NOTES_PATH}")


if __name__ == "__main__":
    main()
