#!/usr/bin/env python3
"""Build the agent-centered thesis defense deck.

Usage:
    pip install python-pptx
    python3 scripts/build_thesis_presentation_agent_zh.py outputs/papers/thesis_presentation_agent_zh.pptx
"""

from __future__ import annotations

import sys
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

W, H = 13.333, 7.5
FONT = "Noto Sans CJK TC"
C = dict(bg="F7F8FA", text="111827", muted="4B5563", line="E5E7EB", white="FFFFFF",
         blue="1D4ED8", blue2="DBEAFE", green="047857", green2="D1FAE5",
         amber="B45309", amber2="FEF3C7", purple="6D28D9", purple2="EDE9FE",
         rose="BE123C", rose2="FFE4E6", slate="0F172A", slate2="E2E8F0", slate_text="CBD5E1")

def I(v): return Inches(v)
def R(h):
    h = h.strip("#")
    return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:6],16))

def fmt(shape, text, size=14, bold=False, color=C["text"], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, margin=.05):
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = I(margin)
    tf.vertical_anchor = anchor
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text, p.alignment, p.space_after = line, align, Pt(4)
        p.font.name, p.font.size, p.font.bold, p.font.color.rgb = FONT, Pt(size), bold, R(color)
    return shape

def txt(slide, text, x, y, w, h, size=14, bold=False, color=C["text"], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    return fmt(slide.shapes.add_textbox(I(x), I(y), I(w), I(h)), text, size, bold, color, align, anchor)

def box(slide, text, x, y, w, h, fill=C["white"], line=C["line"], size=13.5, bold=True, dash=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
    s.fill.solid(); s.fill.fore_color.rgb = R(fill)
    s.line.color.rgb = R(line); s.line.width = Pt(1)
    if dash: s.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return fmt(s, text, size, bold, C["text"], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, .06)

def bullets(slide, items, x, y, w, h, size=14.5):
    return txt(slide, "\n".join("• " + s for s in items), x, y, w, h, size=size)

def line(slide, x1, y1, x2, y2, color="334155", width=1.3):
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    s.line.color.rgb = R(color); s.line.width = Pt(width)
    return s

def table(slide, rows, x, y, widths, rh=.58, size=11):
    for r, row in enumerate(rows):
        cx = x
        for c, val in enumerate(row):
            s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(cx), I(y+r*rh), I(widths[c]), I(rh))
            s.fill.solid(); s.fill.fore_color.rgb = R(C["slate2"] if r == 0 else C["white"])
            s.line.color.rgb = R(C["line"]); s.line.width = Pt(.7)
            fmt(s, val, size + (.5 if r == 0 else 0), r == 0, C["text"], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, .04)
            cx += widths[c]

def background(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = R(C["bg"])
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0), I(0), I(W), I(H))
    s.fill.solid(); s.fill.fore_color.rgb = R(C["bg"]); s.line.color.rgb = R(C["bg"])

def header(slide, title, tag, n):
    txt(slide, title, .62, .32, 8.8, .42, size=20, bold=True)
    if tag: txt(slide, tag, 9.55, .39, 3.1, .28, size=10.5, color=C["muted"], align=PP_ALIGN.RIGHT)
    line(slide, .62, .93, 12.67, .93, color=C["line"], width=1)
    txt(slide, str(n).zfill(2), 12.28, 7.08, .42, .22, size=9, color=C["muted"], align=PP_ALIGN.RIGHT)

def new_slide(prs, title, tag, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    background(s); header(s, title, tag, n)
    return s

def cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); background(s)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0), I(0), I(5.25), I(6.95))
    panel.fill.solid(); panel.fill.fore_color.rgb = R(C["slate"]); panel.line.color.rgb = R(C["slate"])
    box(s, "論文口試報告｜30 min", .72, .56, 2.05, .34, fill=C["blue2"], line=C["blue2"], size=10)
    txt(s, "單房間非連網家電環境影響學習之\n稀疏感測空間數位孿生原型", .72, 1.35, 4.15, 1.45, 26, True, C["white"])
    txt(s, "A Sparse-Sensing Spatial Digital Twin for Learning Environmental Impacts of Non-Networked Appliances in a Single Room", .72, 3.05, 4.05, .95, 13.5, False, C["slate_text"])
    txt(s, "林昀佑｜指導教授：易昶霈教授、沈慧宇副教授", .72, 6.43, 4.15, .35, 11.5, False, "E5E7EB")
    box(s, "核心定位", 6.25, 1.0, 1.2, .42, fill=C["purple2"], line=C["purple2"], size=11)
    txt(s, "不是單純 MCP / Web demo，\n而是 Agent-ready indoor spatial digital twin。", 6.25, 1.65, 5.7, .85, 22, True)
    bullets(s, ["少量感測器推估房間內三因子空間場", "學習冷氣、窗戶、燈光等非連網設備影響", "將數位孿生核心封裝為 AI Agent 可呼叫工具"], 6.3, 3.0, 5.8, 1.35, 15.5)
    box(s, "可解釋｜可校正｜可工具化", 6.3, 5.25, 5.15, .62, size=19)
    txt(s, "01", 12.28, 7.08, .42, .22, 9, False, C["muted"], PP_ALIGN.RIGHT)

def cards(slide, items, y=1.45, card_w=2.55):
    for i, item in enumerate(items):
        title, desc, fill, edge = item
        x = .8 + i * 3.12
        box(slide, title, x, y, card_w, 1.0, fill=fill, line=edge, size=16)
        txt(slide, desc, x+.1, y+1.27, card_w-.2, .95, 13.5, align=PP_ALIGN.CENTER)

def build():
    prs = Presentation(); prs.slide_width = I(W); prs.slide_height = I(H)
    props = prs.core_properties
    props.author = "林昀佑"; props.title = "單房間非連網家電環境影響學習之稀疏感測空間數位孿生原型"
    cover(prs)

    s = new_slide(prs, "研究問題：普通房間沒有完整環境場", "Problem setting", 2)
    box(s, "現實住宅", .8, 1.35, 2.2, .58, fill=C["slate2"], line="94A3B8"); bullets(s,["設備多半沒有 API","感測器數量有限","不同位置體感不同"],.95,2.18,3,1.05,13.5)
    box(s,"研究問題",4.85,1.25,3.65,.62,fill=C["blue2"],line=C["blue"],size=18); txt(s,"在稀疏感測與非連網家電條件下，如何估計單房間內溫度、相對濕度與照度的空間分布，並推論設備對環境的影響？",4.35,2.2,4.6,1.35,18,True,align=PP_ALIGN.CENTER)
    box(s,"數位孿生需求",10,1.35,2.2,.58,fill=C["green2"],line=C["green"]); bullets(s,["可查詢特定座標","可模擬設備狀態","可供 Agent 決策輔助"],9.85,2.18,2.8,1.05,13.2)
    line(s,3.7,2.35,4.35,2.35,C["blue"]); line(s,9.1,2.35,9.75,2.35,C["green"]); box(s,"本研究處理的是「空間環境狀態重建」與「設備影響學習」，不是一般智慧家電控制平台。",1.55,5.35,10.25,.7,size=17)

    s = new_slide(prs,"研究缺口：插值知道距離，但不知道設備語意","Research gap",3)
    table(s, [["方法類型","優點","限制"],["一般感測器監控","容易實作、能記錄時間序列","只能觀察測點，無法形成完整空間場"],["IDW / 空間插值","簡單、可作 baseline","只考慮距離，不理解冷氣、窗戶、燈光"],["純黑箱模型","可能提高預測精度","缺乏設備與幾何語意，難以解釋"],["本研究","結合幾何、設備、稀疏感測校正","限制在單房間 prototype 與分層驗證"]],.75,1.35,[2.15,4.4,5.25],.68)
    box(s,"關鍵差異：本研究不是把感測值直接插值，而是先建立可解釋環境模型，再用感測器殘差修正。",1.2,5.75,10.7,.62,fill=C["blue2"],line=C["blue"],size=16)

    s = new_slide(prs,"總體貢獻：四個層次","Contributions",4)
    cards(s, [("1\n三因子空間數位孿生","溫度、相對濕度、照度\n同時估計房間內座標狀態",C["blue2"],C["blue"]),("2\n非連網設備影響學習","用 before/after 觀測\n學習冷氣、窗戶、燈光影響",C["green2"],C["green"]),("3\n稀疏感測校正","power calibration\ntrilinear residual correction",C["amber2"],C["amber"]),("4\nAgent 工具介面","查詢、模擬、影響學習\n候選動作排序",C["purple2"],C["purple"])])
    box(s,"主張邊界：本研究證明的是單房間稀疏感測條件下的 prototype 與分層驗證，不宣稱已完成多房間泛化或真實閉迴路控制。",1,5.15,11.25,.78,line=C["rose"],size=16)

    s = new_slide(prs,"系統架構：從房間語意到 Agent 工具","System architecture",5)
    xs=[.65,3,5.35,7.7,10.05]; labels=["Room schema\ngeometry / devices","Sparse sensing\ncorner sensors","Digital twin core\nthree-factor fields","Calibration / learning\nresidual + impact","Agent tool layer\nquery / simulate / rank"]
    fills=[C["slate2"],C["blue2"],C["green2"],C["amber2"],C["purple2"]]; edges=["94A3B8",C["blue"],C["green"],C["amber"],C["purple"]]
    for x,l,f,e in zip(xs,labels,fills,edges): box(s,l,x,2.05,1.95,1.05,fill=f,line=e,size=12.5)
    for x in [2.6,4.95,7.3,9.65]: line(s,x,2.58,x+.4,2.58)
    bullets(s,["MCP 是其中一種工具協定實作；論文主體應放在 Agent 可呼叫的 digital twin core","Web demo 是展示層；不應取代方法與驗證層敘事","模型、校正、工具化三者要分層描述"],1.05,4.55,11.2,1.2,14.5)

    s = new_slide(prs,"三因子模型：變數專屬 reduced-order nominal model","Variable-specific modeling",6)
    table(s, [["因子","主要來源","模型語意"],["溫度","冷氣、窗戶、燈光熱效應","熱交換、設備距離衰減、空間分布"],["相對濕度","冷氣除濕、開窗換氣、外部濕度","水氣交換與設備狀態影響"],["照度","窗戶自然光、燈光、遮蔽與反射","幾何光照、距離衰減、家具遮蔽"]],1,1.5,[1.6,4.25,5.35],.72)
    box(s,"方法定位：不是三個因子套同一個插值公式，而是依照變數特性建立 nominal model。",1.25,5.15,10.6,.68,fill=C["green2"],line=C["green"],size=16.5)

    s = new_slide(prs,"稀疏感測校正：從少量觀測補正整體空間場","Sparse calibration",7)
    for l,x,w,f,e in [("Nominal prediction",.95,2.35,C["slate2"],"94A3B8"),("Sensor residuals",4.1,2.35,C["blue2"],C["blue"]),("Trilinear correction",7.25,2.35,C["green2"],C["green"]),("Corrected field",10.35,2.1,C["amber2"],C["amber"])]: box(s,l,x,1.7,w,.75,fill=f,line=e)
    for x,c in [(3.3,C["blue"]),(6.45,C["green"]),(9.6,C["amber"])]: line(s,x,2.08,x+.8,2.08,c)
    bullets(s,["active-device power calibration：根據設備附近誤差調整設備影響強度","trilinear residual correction：用角落感測器殘差修正房間內部估計","校正目標是吸收真實觀測差異，而不是取代原本的物理／幾何模型"],1.15,4.15,10.9,1.25,15)

    s = new_slide(prs,"非連網家電影響學習：用環境變化反推設備效果","Impact learning",8)
    for l,x,f,e in [("Before\n設備作用前",1,C["white"],C["line"]),("After\n設備作用後",4,C["white"],C["line"]),("Delta field\n影響場",7,C["blue2"],C["blue"]),("Impact model\n設備參數",10,C["green2"],C["green"])]: box(s,l,x,1.65,2,.85,fill=f,line=e)
    for x in [3,6,9]: line(s,x,2.08,x+1,2.08)
    table(s, [["設備","學習／模擬的主要影響"],["冷氣","降溫、除濕、距離與方向性影響"],["窗戶","外部溫濕度、日照、開啟比例影響"],["燈光","照度提升與局部熱效應"]],2,4.2,[2.15,7.1],.58)

    s = new_slide(prs,"Hybrid residual neural network：黑箱只學剩餘誤差","Residual learning",9)
    for l,x,w,f,e in [("可解釋主模型",1,2.4,C["green2"],C["green"]),("+ 感測器校正",4,2.4,C["blue2"],C["blue"]),("+ 殘差神經網路",7,2.4,C["purple2"],C["purple"]),("= 最終估計場",10,2.25,C["amber2"],C["amber"])]: box(s,l,x,1.6,w,.72,fill=f,line=e,size=17)
    for x in [3.4,6.4,9.4]: line(s,x,1.96,x+.6,1.96)
    bullets(s,["神經網路不是取代數位孿生模型，而是學 nominal model 無法解釋的 residual","保留設備、幾何、因子語意，並以資料驅動方式補償誤差","口試時要避免讓審查者以為整個方法只是 black-box regression"],1.15,4,10.8,1.3,15)

    s = new_slide(prs,"Agent tool interface：讓模型被查詢、模擬與排序","Agent integration",10)
    for k,(name,desc) in enumerate([("initialize_environment","建立房間、設備、感測器與外部條件"),("sample_point","查詢指定座標的三因子狀態"),("learn_impacts","由 before/after 資料學習設備影響"),("run_window_direct","模擬外部天氣與開窗比例"),("rank_actions","根據目標環境排序候選動作")]):
        y=1.25+k*.78; box(s,name,1,y,3.1,.5,fill=C["purple2"],line=C["purple"],size=11.5); txt(s,desc,4.35,y+.08,6.9,.32,13.5)
    box(s,"推薦表述：MCP 是實作通道；Agent tool interface 才是論文中要強調的系統整合貢獻。",1.2,5.55,10.6,.68,line=C["purple"],size=16)

    s = new_slide(prs,"展示層：Web demo 的正確定位","Visualization layer",11)
    for l,x,w,f,e in [("Digital twin core",1,2.5,C["green2"],C["green"]),("API / tool layer",4.2,2.5,C["purple2"],C["purple"]),("Web demo",7.4,2.5,C["blue2"],C["blue"]),("Defense explanation",10.4,1.9,C["amber2"],C["amber"])]: box(s,l,x,1.8,w,.72,fill=f,line=e,size=13)
    for x in [3.5,6.7,9.9]: line(s,x,2.16,x+.7,2.16)
    bullets(s,["Web demo 用於展示 3D 空間點、設備狀態切換與三因子分布","它不是主要方法，也不是主要驗證來源","口試時應用它輔助說明模型輸入、輸出與系統流程"],1.2,4.1,10.8,1.2,15)

    tables = [
      ("驗證設計：分層驗證，不混淆證據層級","Evaluation design", [["層級","驗證目的","不能過度宣稱"],["受控模擬","檢查模型是否能重建合理空間場","不能等同真實房間完整驗證"],["IDW baseline","比較設備語意模型與純距離插值","不能只看單一情境"],["Ablation","檢查各模組是否有貢獻","不能說所有指標必然改善"],["真實臥室 snapshot","驗證 sparse calibration 對保留點有幫助","不能宣稱完整 3D ground truth"],["Public benchmark","補充 task-aligned 外部比較","不能說完全驗證本系統"]], [.65,1.2,[2,5.05,5],.55]),
      ("結果二：消融分析檢查每個模組的角色","Ablation study", [["移除模組","觀察重點"],["without reflection / geometry detail","照度與遮蔽相關誤差變化"],["without power calibration","設備附近誤差是否放大"],["without trilinear correction","角落感測器殘差是否無法傳遞到內部空間"],["without hybrid residual NN","是否缺少資料驅動補償能力"]], [1.1,1.35,[3.9,7.2],.65]),
      ("限制與風險：提前講清楚比較安全","Limitations", [["類別","限制","對應說法"],["資料限制","真實資料為單房間 snapshot，缺完整 dense-field ground truth","定位為 prototype 與保留點驗證"],["方法限制","reduced-order 模型是工程近似，hybrid residual 需防 overfitting","保留可解釋主模型並明確標註誤差來源"],["系統限制","Agent/MCP 是工具化層，尚非真實閉迴路控制","稱為 counterfactual ranking 與 decision support"]], [.8,1.45,[1.8,5.45,4.5],.78]),
    ]
    s = new_slide(prs,tables[0][0],tables[0][1],12); table(s,tables[0][2],*tables[0][3])

    s = new_slide(prs,"結果一：受控模擬與 IDW baseline","Controlled simulation",13)
    for title,bs,x,w,bw,f,e in [("IDW baseline",["只看感測點距離","缺少冷氣／窗戶／燈光位置語意","設備影響明顯時容易失真"],1,2.5,3.1,C["rose2"],C["rose"]),("本研究模型",["使用設備影響函數","保留房間幾何與因子差異","再用感測器殘差修正"],5.35,2.5,3.1,C["blue2"],C["blue"]),("結論",["模型價值來自設備語意","不是單純增加參數","適合 sparse sensing 場景"],9.55,2.3,2.85,C["green2"],C["green"])]: box(s,title,x,1.45,w,.72,fill=f,line=e,size=18); bullets(s,bs,x+.05,2.45,bw,1.1,13.2)
    box(s,"報告時應說明：IDW 是合理 baseline，但不是能處理非連網設備影響的充分模型。",1.2,5.6,10.6,.62,size=16)

    s = new_slide(prs,tables[1][0],tables[1][1],14); table(s,tables[1][2],*tables[1][3]); box(s,"Ablation 的用法：證明系統不是功能堆疊，而是每一層都對特定誤差來源有對應作用。",1.3,5.55,10.4,.62,fill=C["amber2"],line=C["amber"],size=16)

    s = new_slide(prs,"結果三：真實臥室 snapshot 的定位","Real-bedroom snapshot",15)
    for title,bs,x,f,e in [("可以主張",["真實觀測可被納入校正流程","保留點估計可用來檢查校正效果","可證明 prototype 與真實場域銜接"],1,C["green2"],C["green"]),("不要主張",["已完整驗證 3D dense field","已泛化到所有房型","已完成閉迴路控制實驗"],5,C["rose2"],C["rose"]),("論文寫法",["稱為 snapshot validation","明確標註 claim boundary","搭配模擬與 benchmark 補足"],9,C["blue2"],C["blue"])]: box(s,title,x,1.45,2.2,.6,fill=f,line=e); bullets(s,bs,x+.05,2.3,3.35,1.3,13.5)

    s = new_slide(prs,"結果四：公開資料集作為 task-aligned benchmark","External benchmarks",16)
    for title,desc,x,f,e in [("SML2010","用於補充室內環境預測相關任務比較",1,C["blue2"],C["blue"]),("CU-BEMS","用於補充建築能源與環境感測任務比較",5.4,C["amber2"],C["amber"]),("正確定位","沒有完整房間幾何\n沒有 8 角落感測配置\n不是 3D dense-field ground truth",9.55,C["rose2"],C["rose"])]: box(s,title,x,1.6,2.45,.7,fill=f,line=e,size=18); txt(s,desc,x-.05,2.65,2.75,.95,13,align=PP_ALIGN.CENTER)
    box(s,"結論：公開資料集能補充外部任務相容性，但不能替代本研究的空間場驗證。",1.25,5.4,10.6,.62,size=16)

    s = new_slide(prs,"擬新增延伸：自由空間與模組化 estimator","Proposed extension",17)
    box(s,"目前主線\nPhysics spine",.95,2.1,2.2,.85,fill=C["green2"],line=C["green"],size=15)
    for l,x,y,w in [("自由空間拓樸\nΩ_room / Ω_occ / Ω_free",3.95,1.45,2.7),("節點語義\nV_geom / V_obs / V_target",3.95,2.75,2.7),("模組化 estimator\nIDW / 2D / 3D / Cell-IDW",7.45,2.1,2.8),("Confidence / provenance\nblocked CV / intervention",10.85,2.1,1.85)]: box(s,l,x,y,w,.85,line=C["purple"],dash=True,size=12)
    for a in [(3.15,2.52,3.95,1.88),(3.15,2.52,3.95,3.18),(6.65,2.18,7.45,2.52),(6.65,3.18,7.45,2.52),(10.25,2.52,10.85,2.52)]: line(s,*a,color=C["purple"])
    bullets(s,["虛線外框表示 proposed extension，不放進已完成驗證主張","可放在 future work 或系統擴充章節","用來表達系統不是純樹狀，而是多模組互相關聯"],1.1,4.85,10.9,1,14)

    s = new_slide(prs,tables[2][0],tables[2][1],18); table(s,tables[2][2],*tables[2][3]); box(s,"限制不是弱點，而是 claim boundary。只要證據層次與主張層次對齊，論文會更可信。",1.25,5.7,10.55,.62,size=16)

    s = new_slide(prs,"口試敘事主線：從問題到貢獻","Defense storyline",19)
    labels=["普通房間缺完整環境場","建立三因子可解釋模型","用稀疏感測校正","學習非連網設備影響","提供 Agent 工具化使用"]
    for k,l in enumerate(labels):
        x=.75+k*2.45; y=1.7+(1.15 if k==4 else 0); box(s,f"{k+1}\n{l}",x,y,2.05,.82,fill=C["purple2"] if k==4 else C["blue2"],line=C["purple"] if k==4 else C["blue"],size=12.5)
        if k<4: line(s,x+2.05,2.1,x+2.45,3.25 if k==3 else 2.1)
    for k,l in enumerate(["可解釋","可校正","可工具化"]): box(s,l,[1.2,3.3,5.4][k],5.1,[1.8,1.8,2][k],.55,line=[C["green"],C["blue"],C["purple"]][k],size=18)
    txt(s,"不建議一直講「我做了 MCP」；建議說「我把 digital twin core 做成 Agent 可以使用的 tool interface」。",8.1,4.82,4.2,.92,15,True)

    s = new_slide(prs,"結論：本研究的主要貢獻","Conclusion",20)
    for k,(num,body,col) in enumerate([("1","提出單房間稀疏感測三因子空間數位孿生原型。",C["blue"]),("2","以變數專屬 reduced-order nominal model 保留溫度、濕度、照度的物理與幾何語意。",C["green"]),("3","結合 power calibration、trilinear residual correction 與 hybrid residual NN 改善估計。",C["purple"]),("4","把模型封裝為 Agent 可呼叫工具介面，支援查詢、模擬、影響學習與反事實排序。",C["amber"])]):
        y=1.48+k*1.05; box(s,num,1,y,.55,.55,fill=col,line=col,size=18); txt(s,body,1.85,y+.08,10.4,.4,17,True)
    box(s,"最終定位：不是一般插值，也不是單純 Web/MCP demo，而是 Agent-ready indoor spatial digital twin。",1.35,6.1,10.7,.65,size=18)
    return prs

def main():
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("outputs/papers/thesis_presentation_agent_zh.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs = build(); prs.save(out)
    print(f"Wrote {out}"); print(f"Slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
