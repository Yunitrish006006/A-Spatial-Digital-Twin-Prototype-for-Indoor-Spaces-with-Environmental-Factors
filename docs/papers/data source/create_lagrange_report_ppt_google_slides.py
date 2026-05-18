"""
Revised PowerPoint generator for the paper:
Galindo, Ike, Liu (2022), “Error-constant estimation under the maximum norm
for linear Lagrange interpolation”.

Outputs:
  - /mnt/data/lagrange_error_constant_report_revised.pptx
  - /mnt/data/lagrange_error_constant_speaker_notes_revised.txt

Design goals:
  - large readable Chinese typography
  - no text overflow/out-of-bounds elements
  - professional one-hour professor-facing report
  - top-down explanation with demonstration calculation
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterable, List, Tuple
import math

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

OUT_DIR = Path('/mnt/data')
PPTX_PATH = OUT_DIR / 'lagrange_error_constant_report_google_slides.pptx'
NOTES_PATH = OUT_DIR / 'lagrange_error_constant_speaker_notes_google_slides.txt'

# --- page and theme ---
W, H = Inches(13.333), Inches(7.5)
FONT_CJK = 'Noto Sans TC'
FONT_EN = 'Arial'
FONT_MATH = 'Arial'

NAVY = RGBColor(20, 35, 62)
BLUE = RGBColor(42, 91, 149)
TEAL = RGBColor(45, 134, 137)
ORANGE = RGBColor(213, 122, 68)
GREEN = RGBColor(74, 137, 85)
RED = RGBColor(174, 76, 73)
GRAY = RGBColor(93, 103, 118)
LIGHT = RGBColor(246, 248, 251)
MID = RGBColor(225, 232, 242)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(31, 38, 50)
LIGHT_BLUE = RGBColor(229, 239, 252)
LIGHT_TEAL = RGBColor(226, 245, 244)
LIGHT_ORANGE = RGBColor(252, 239, 229)
LIGHT_GREEN = RGBColor(232, 245, 235)
LIGHT_RED = RGBColor(251, 233, 232)



def inch(v: float):
    return Inches(v)


def apply_text(p, text='', size=20, color=DARK, bold=False, font=FONT_CJK, align=None):
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    if align is not None:
        p.alignment = align


def set_tf(tf, margin=0.12, valign=MSO_ANCHOR.TOP):
    tf.margin_left = inch(margin)
    tf.margin_right = inch(margin)
    tf.margin_top = inch(margin)
    tf.margin_bottom = inch(margin)
    tf.word_wrap = True
    tf.vertical_anchor = valign


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()
    return bg


def add_title(slide, title: str, section: str = '', no: int | None = None):
    if section:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(0), inch(0), inch(13.333), inch(0.22))
        bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
        bar.line.fill.background()
        st = slide.shapes.add_textbox(inch(0.68), inch(0.30), inch(11.0), inch(0.28))
        apply_text(st.text_frame.paragraphs[0], section, 11, BLUE, True)
    tx = slide.shapes.add_textbox(inch(0.68), inch(0.56), inch(11.6), inch(0.52))
    apply_text(tx.text_frame.paragraphs[0], title, 25, NAVY, True)
    if no is not None:
        sn = slide.shapes.add_textbox(inch(12.35), inch(0.60), inch(0.45), inch(0.25))
        apply_text(sn.text_frame.paragraphs[0], f'{no:02d}', 10, GRAY, False, FONT_EN, PP_ALIGN.RIGHT)
    return tx


def add_footer(slide, no: int):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(0.68), inch(7.05), inch(12.0), inch(0.012))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(210, 218, 230)
    line.line.fill.background()
    foot = slide.shapes.add_textbox(inch(0.68), inch(7.13), inch(10.4), inch(0.22))
    apply_text(foot.text_frame.paragraphs[0], 'Galindo, Ike & Liu (2022)｜Maximum-norm error constant for linear Lagrange interpolation', 7.2, GRAY, False, FONT_EN)
    s = slide.shapes.add_textbox(inch(12.25), inch(7.12), inch(0.45), inch(0.22))
    apply_text(s.text_frame.paragraphs[0], str(no), 8.5, GRAY, False, FONT_EN, PP_ALIGN.RIGHT)


def add_box(slide, x, y, w, h, fill=WHITE, line=MID, radius=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 inch(x), inch(y), inch(w), inch(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    return shp


def add_text(slide, x, y, w, h, text, size=20, color=DARK, bold=False, fill=None, align=None, valign=MSO_ANCHOR.TOP):
    shp = add_box(slide, x, y, w, h, fill=fill if fill else WHITE, line=MID, radius=True)
    tf = shp.text_frame; tf.clear(); set_tf(tf, 0.16, valign)
    p = tf.paragraphs[0]
    apply_text(p, text, size, color, bold, FONT_CJK, align)
    return shp


def add_label(slide, x, y, w, h, text, color=BLUE, size=17):
    # Use light fills + dark text so rendered slides remain readable even if viewed as images.
    tint_map = {
        (BLUE.rgb if hasattr(BLUE, 'rgb') else None): LIGHT_BLUE,
    }
    if color == BLUE or color == NAVY:
        fill = LIGHT_BLUE; txt = NAVY; line = BLUE
    elif color == TEAL:
        fill = LIGHT_TEAL; txt = NAVY; line = TEAL
    elif color == ORANGE:
        fill = LIGHT_ORANGE; txt = NAVY; line = ORANGE
    elif color == GREEN:
        fill = LIGHT_GREEN; txt = NAVY; line = GREEN
    elif color == RED:
        fill = LIGHT_RED; txt = NAVY; line = RED
    else:
        fill = RGBColor(239, 244, 250); txt = NAVY; line = color
    shp = add_box(slide, x, y, w, h, fill=fill, line=line, radius=True)
    tf = shp.text_frame; tf.clear(); set_tf(tf, 0.07, MSO_ANCHOR.MIDDLE)
    apply_text(tf.paragraphs[0], text, size, txt, True, FONT_CJK, PP_ALIGN.CENTER)
    return shp


def add_heading_body(slide, x, y, w, h, heading, body_lines: Iterable[str], accent=BLUE, body_size=17):
    shp = add_box(slide, x, y, w, h, WHITE, MID, True)
    tf = shp.text_frame; tf.clear(); set_tf(tf, 0.20, MSO_ANCHOR.TOP)
    apply_text(tf.paragraphs[0], heading, 18, accent, True)
    for line in body_lines:
        p = tf.add_paragraph()
        apply_text(p, '• ' + line, body_size, DARK, False)
        p.space_after = Pt(4)
        p.line_spacing = 1.08
    return shp


def add_formula(slide, x, y, w, h, formula, caption='', size=20):
    shp = add_box(slide, x, y, w, h, RGBColor(239, 244, 250), RGBColor(195, 208, 225), True)
    tf = shp.text_frame; tf.clear(); set_tf(tf, 0.12, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    apply_text(p, formula, size, NAVY, True, FONT_MATH, PP_ALIGN.CENTER)
    if caption:
        p2 = tf.add_paragraph()
        apply_text(p2, caption, 9.5, GRAY, False, FONT_CJK, PP_ALIGN.CENTER)
    return shp


def add_triangle(slide, x, y, w, h, label='K', grid=False):
    ax, ay, aw, ah = inch(x), inch(y), inch(w), inch(h)
    p1 = (ax + aw*0.14, ay + ah*0.80)
    p2 = (ax + aw*0.86, ay + ah*0.80)
    p3 = (ax + aw*0.24, ay + ah*0.18)
    if grid:
        # thin internal triangulation lines behind the main triangle
        for t in [0.25, 0.5, 0.75]:
            a = (p1[0] + (p3[0]-p1[0])*t, p1[1] + (p3[1]-p1[1])*t)
            b = (p2[0] + (p3[0]-p2[0])*t, p2[1] + (p3[1]-p2[1])*t)
            ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, a[0], a[1], b[0], b[1])
            ln.line.color.rgb = RGBColor(170, 185, 205); ln.line.width = Pt(0.7)
        for t in [0.25, 0.5, 0.75]:
            a = (p1[0] + (p2[0]-p1[0])*t, p1[1])
            b = (p3[0] + (p2[0]-p3[0])*t, p3[1] + (p2[1]-p3[1])*t)
            ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, a[0], a[1], b[0], b[1])
            ln.line.color.rgb = RGBColor(170, 185, 205); ln.line.width = Pt(0.7)
    for a, b in [(p1,p2),(p2,p3),(p3,p1)]:
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, a[0], a[1], b[0], b[1])
        ln.line.color.rgb = BLUE; ln.line.width = Pt(2.2)
    for i, pt in enumerate([p1,p2,p3], 1):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, pt[0]-inch(0.055), pt[1]-inch(0.055), inch(0.11), inch(0.11))
        dot.fill.solid(); dot.fill.fore_color.rgb = ORANGE; dot.line.fill.background()
        lab = slide.shapes.add_textbox(pt[0]-inch(0.10), pt[1]+inch(0.08), inch(0.35), inch(0.20))
        apply_text(lab.text_frame.paragraphs[0], f'p{i}', 8.5, GRAY, False, FONT_EN)
    tt = slide.shapes.add_textbox(ax + aw*0.48, ay + ah*0.46, inch(1.1), inch(0.28))
    apply_text(tt.text_frame.paragraphs[0], label, 16, NAVY, True, FONT_EN)


def add_arrow(slide, x1, y1, x2, y2, color=GRAY, width=1.4):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    ln.line.end_arrowhead = True
    return ln


def add_flow(slide, items: List[Tuple[str, str]], x0=0.78, y=2.0, gap=0.25, box_w=2.25, box_h=1.0):
    # connectors first so they stay behind boxes
    for i in range(len(items)-1):
        x1 = x0 + i*(box_w+gap) + box_w
        x2 = x0 + (i+1)*(box_w+gap)
        add_arrow(slide, x1+0.02, y+box_h/2, x2-0.08, y+box_h/2, GRAY, 1.2)
    for i, (head, body) in enumerate(items):
        x = x0 + i*(box_w+gap)
        shp = add_box(slide, x, y, box_w, box_h, WHITE, MID, True)
        tf = shp.text_frame; tf.clear(); set_tf(tf, 0.09, MSO_ANCHOR.MIDDLE)
        apply_text(tf.paragraphs[0], head, 14.5, BLUE, True, FONT_CJK, PP_ALIGN.CENTER)
        p2 = tf.add_paragraph()
        apply_text(p2, body, 10.5, DARK, False, FONT_CJK, PP_ALIGN.CENTER)


def add_table(slide, x, y, w, h):
    data = [
        ['最大角 θ', 'Cᴸ 下界', 'λh,B', 'Cᴸ 上界'],
        ['π/6  (30°)', '0.31511', '9.8925', '0.31799'],
        ['π/4  (45°)', '0.26777', '13.574', '0.27146'],
        ['π/3  (60°)', '0.25209', '15.457', '0.25439'],
        ['π/2  (90°)', '0.40432', '5.7812', '0.41596'],
    ]
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, inch(x), inch(y), inch(w), inch(h))
    tbl = tbl_shape.table
    widths = [1.9, 1.35, 1.35, 1.35]
    for c, ww in enumerate(widths):
        tbl.columns[c].width = inch(ww)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = data[r][c]
            cell.margin_left = inch(0.04); cell.margin_right = inch(0.04)
            cell.margin_top = inch(0.04); cell.margin_bottom = inch(0.04)
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BLUE if r == 0 else WHITE
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT_CJK if c == 0 else FONT_EN
                    run.font.size = Pt(12.5 if r > 0 else 11.5)
                    run.font.bold = r == 0
                    run.font.color.rgb = NAVY if r == 0 else DARK
    return tbl_shape


def add_chart(slide, x, y, w, h):
    data = CategoryChartData()
    data.categories = ['π/6', 'π/4', 'π/3', 'π/2']
    data.add_series('下界', (0.31511, 0.26777, 0.25209, 0.40432))
    data.add_series('上界', (0.31799, 0.27146, 0.25439, 0.41596))
    frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, inch(x), inch(y), inch(w), inch(h), data)
    ch = frame.chart
    ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.value_axis.minimum_scale = 0.20
    ch.value_axis.maximum_scale = 0.45
    ch.category_axis.tick_labels.font.size = Pt(10)
    ch.value_axis.tick_labels.font.size = Pt(10)
    ch.chart_title.text_frame.text = 'Table 2：誤差常數上下界'
    ch.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    return frame


def warnIfSlideElementsOutOfBounds(prs: Presentation) -> List[str]:
    warnings = []
    for si, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            try:
                if shp.left < 0 or shp.top < 0 or shp.left + shp.width > prs.slide_width or shp.top + shp.height > prs.slide_height:
                    warnings.append(f'Slide {si}: out-of-bounds shape={shp.name}')
            except Exception:
                pass
    return warnings


def warnIfSlideHasOverlaps(prs: Presentation) -> List[str]:
    warnings = []
    for si, slide in enumerate(prs.slides, 1):
        boxes = []
        for shp in slide.shapes:
            try:
                if shp.shape_type == 9:  # connector
                    continue
                if shp.width > prs.slide_width * 0.95 and shp.height > prs.slide_height * 0.90:  # background
                    continue
                if shp.height < Inches(0.05):
                    continue
                # ignore text labels on diagrams: intentionally close to nodes
                if shp.width < Inches(0.5) and shp.height < Inches(0.35):
                    continue
                boxes.append((shp.left, shp.top, shp.left + shp.width, shp.top + shp.height, shp.name))
            except Exception:
                continue
        for i in range(len(boxes)):
            l1, t1, r1, b1, n1 = boxes[i]
            a1 = max(1, (r1-l1)*(b1-t1))
            for j in range(i+1, len(boxes)):
                l2, t2, r2, b2, n2 = boxes[j]
                iw = min(r1, r2) - max(l1, l2)
                ih = min(b1, b2) - max(t1, t2)
                if iw <= 0 or ih <= 0:
                    continue
                a2 = max(1, (r2-l2)*(b2-t2))
                ratio = (iw*ih) / min(a1, a2)
                if ratio > 0.40:
                    warnings.append(f'Slide {si}: overlap {n1} / {n2} ({ratio:.2f})')
    return warnings


def build_deck():
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    blank = prs.slide_layouts[6]
    notes: List[str] = []

    # 1 title
    slide = prs.slides.add_slide(blank); add_bg(slide)
    add_label(slide, 0.72, 0.58, 2.7, 0.40, '論文報告｜一小時版', BLUE, 15)
    tb = slide.shapes.add_textbox(inch(0.72), inch(1.28), inch(8.6), inch(1.35))
    apply_text(tb.text_frame.paragraphs[0], '最大範數下線性 Lagrange 插值\n誤差常數估計', 33, NAVY, True)
    sub = slide.shapes.add_textbox(inch(0.76), inch(2.95), inch(8.5), inch(0.42))
    apply_text(sub.text_frame.paragraphs[0], 'Error-constant estimation under the maximum norm for linear Lagrange interpolation', 14, GRAY, False, FONT_EN)
    add_triangle(slide, 9.25, 1.05, 3.1, 2.65, 'K')
    add_label(slide, 0.95, 4.42, 3.1, 0.78, '核心結果\n0.40432 ≤ Cᴸ ≤ 0.41596', TEAL, 17)
    add_label(slide, 4.45, 4.42, 3.0, 0.78, '方法\nFEM + Bernstein', BLUE, 17)
    add_label(slide, 7.95, 4.42, 3.0, 0.78, '時間\n≤ 60 分鐘', ORANGE, 17)
    src = slide.shapes.add_textbox(inch(0.75), inch(6.17), inch(11.7), inch(0.45))
    apply_text(src.text_frame.paragraphs[0], 'Galindo, Ike & Liu, Journal of Inequalities and Applications, 2022:109', 11.5, GRAY, False, FONT_EN)
    add_footer(slide, 1)
    notes.append('開場：本報告介紹 Galindo、Ike 與 Liu 在 2022 年提出的最大範數下線性 Lagrange 插值誤差常數估計。先講結論：對單位等腰直角三角形，Cᴸ 被夾在 0.40432 與 0.41596 之間。')

    # 2 top-down map
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '報告架構：先答案，再拆方法', 'Top-down roadmap', 2)
    add_text(slide, 0.85, 1.35, 4.15, 4.65, '主線問題\n\n1. 這篇論文要保證什麼誤差？\n\n2. 為什麼 L∞ 最大範數特別難？\n\n3. FEM、Fujino–Morley 與 Bernstein 如何合作？\n\n4. 數值結果如何解讀？', 18, DARK, False)
    add_flow(slide, [('背景', '10 分'), ('定義', '8 分'), ('方法', '20 分'), ('計算', '12 分'), ('結論', '8 分')], x0=5.35, y=1.55, box_w=1.25, box_h=0.78, gap=0.20)
    add_text(slide, 5.2, 3.10, 6.9, 1.25, '今天只需抓住一個核心：\n把「無限多函數中的最壞誤差」轉成「有限維矩陣最佳化」。', 21, NAVY, True, RGBColor(237, 245, 252), PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_text(slide, 5.2, 4.78, 6.9, 0.90, '報告節奏：主講約 58 分鐘，最後保留 Q&A。', 18, BLUE, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_footer(slide, 2)
    notes.append('說明報告採 top-down 結構：先回答論文核心答案，再拆開方法。聽眾不需要一開始就進入所有公式，只要先知道研究在建立一個嚴格的最大誤差保證。')

    # 3 intuition
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '直觀問題：用平面近似曲面，最大誤差是多少？', 'Motivation', 3)
    add_triangle(slide, 0.88, 1.45, 3.7, 3.35, 'triangle K', grid=True)
    add_heading_body(slide, 4.90, 1.30, 3.45, 3.05, '線性 Lagrange 插值', ['在三角形三個頂點取函數值', '用一個線性函數通過三點', '比較原函數 u 與 Πᴸu 的差'], TEAL, 16.5)
    add_heading_body(slide, 8.75, 1.30, 3.45, 3.05, '最大範數 L∞', ['不是平均誤差', '看整個三角形上的最壞點', '適合做保守誤差保證'], ORANGE, 16.5)
    add_formula(slide, 1.55, 5.35, 10.3, 0.86, '||u − Πᴸu||∞,K  ≤  Cᴸ(K) |u|₂,K', 'Cᴸ(K) 就是本文要估計的誤差常數', 21)
    add_footer(slide, 3)
    notes.append('用幾何直覺說明：線性 Lagrange 插值就是用一個平面通過三角形三個頂點。最大範數問的是所有點中最壞點的誤差，因此比平均誤差更保守，也更適合嚴格估計。')

    # 4 contribution
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '論文貢獻：把難解的最壞誤差變成可算的嚴格界', 'Contribution', 4)
    add_heading_body(slide, 0.80, 1.30, 3.75, 4.45, '問題層', ['估計 Cᴸ(K)', '誤差量：L∞ 最大範數', '控制量：H² 半範數，反映曲率'], BLUE, 16.5)
    add_heading_body(slide, 4.82, 1.30, 3.75, 4.45, '方法層', ['FEM 離散化函數空間', 'Fujino–Morley 提供正交分解', 'Bernstein 凸包性質處理最大值'], TEAL, 16.5)
    add_heading_body(slide, 8.84, 1.30, 3.75, 4.45, '結果層', ['多種三角形的上下界', '等腰直角三角形：0.40432 ≤ Cᴸ ≤ 0.41596', '相對誤差小，可支援 FEM 誤差控制'], ORANGE, 16.5)
    add_footer(slide, 4)
    notes.append('這張總結貢獻：論文不是只做數值實驗，而是建立一套可嚴格估計誤差常數的方法。教授聽的重點是定義、方法與上下界三件事。')

    # 5 formulation
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '數學定義：Cᴸ(K) 是最壞比例', 'Problem definition', 5)
    add_formula(slide, 0.95, 1.30, 5.55, 1.00, 'Cᴸ(K) = sup  ||u − Πᴸu||∞,K / |u|₂,K', '在所有 u∈H²(K) 中取最壞情況', 18.5)
    add_formula(slide, 6.90, 1.30, 5.55, 1.00, 'λ(K) = inf  |u|²₂,K / ||u||²∞,K', '等價形式：Cᴸ(K)=1/√λ(K)', 18.5)
    add_heading_body(slide, 0.95, 3.00, 5.55, 2.42, '為什麼轉成 λ？', ['Cᴸ 是誤差常數；λ 是能量比例', 'λ 越大，Cᴸ = 1/√λ 越小', '估計 λ 的下界，可推得 Cᴸ 的上界'], GREEN, 17)
    add_heading_body(slide, 6.90, 3.00, 5.55, 2.42, '困難在哪裡？', ['最大值可能在三角形任意位置', 'L∞ 約束不是簡單線性條件', '直接在 H²(K) 上求解不可行'], RED, 17)
    add_footer(slide, 5)
    notes.append('這張正式定義 Cᴸ(K)。它是最大誤差除以曲率尺度的最壞比例。作者再引入 λ，因為 Cᴸ 等於 λ 的倒平方根；只要把 λ 夾住，就能得到 Cᴸ 的界。')

    # 6 1D warmup
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '暖身示範：一維案例讓誤差常數變直觀', 'Warm-up calculation', 6)
    add_text(slide, 0.85, 1.35, 3.75, 3.75, '取 u(x)=x²，x∈[0,1]\n\n端點值：u(0)=0, u(1)=1\n\n線性插值：Πᴸu(x)=x\n\n誤差：u−Πᴸu=x²−x', 18.2, DARK, False)
    add_formula(slide, 5.15, 1.42, 3.30, 0.74, 'max |x² − x| = 1/4', '', 18)
    add_formula(slide, 5.15, 2.45, 3.30, 0.74, '||u″||∞ = 2', '', 18)
    add_formula(slide, 5.15, 3.48, 3.30, 0.74, '(1/4) / 2 = 1/8', '', 18)
    add_heading_body(slide, 8.95, 1.35, 3.30, 3.75, '這個例子說明', ['常數是「最壞比例」', '二維三角形是同一想法的高階版本', '本文改用 H² 半範數控制 L∞ 誤差'], TEAL, 16.5)
    add_footer(slide, 6)
    notes.append('一維暖身：u=x² 的線性插值是 x，誤差 x²−x 在中點最大，值是 1/4。二階導數最大值是 2，因此比例是 1/8。這說明誤差常數的意義。')

    # 7 FEM
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '第一個工具：FEM 將函數問題離散成矩陣問題', 'Method 1: FEM', 7)
    add_flow(slide, [('原始問題', 'u∈H²(K)'), ('剖分 K', '小三角形'), ('FM 空間', 'VhFM'), ('矩陣 A', 'xᵀAx'), ('可計算', '有限維')], x0=0.75, y=1.65, box_w=2.20, box_h=1.00, gap=0.20)
    add_heading_body(slide, 0.95, 3.55, 5.50, 2.45, 'Fujino–Morley 空間的角色', ['用分片二次多項式近似誤差函數', '利用插值條件與正交分解', '把 H² 能量寫成矩陣二次型 xᵀAx'], BLUE, 17)
    add_heading_body(slide, 6.90, 3.55, 5.50, 2.45, '直觀理解', ['原本要在所有曲面中找最壞者', '現在改在有限個基底係數中搜尋', '網格細化後上下界會更接近'], ORANGE, 17)
    add_footer(slide, 7)
    notes.append('FEM 的作用是把無限維函數空間轉為有限維係數向量。Fujino–Morley 空間的正交性讓 H² 能量能以 xᵀAx 的形式表示，這是後續矩陣最佳化的基礎。')

    # 8 Bernstein
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '第二個工具：Bernstein 凸包性質處理最大範數', 'Method 2: Bernstein', 8)
    add_formula(slide, 0.95, 1.25, 5.30, 0.82, 'p = Σ dᵢⱼₖ Jᵢⱼₖ⁽ⁿ⁾', '', 19)
    add_formula(slide, 0.95, 2.34, 5.30, 0.82, 'Jᵢⱼₖ ≥ 0,    ΣJᵢⱼₖ = 1', '', 19)
    add_formula(slide, 0.95, 3.43, 5.30, 0.82, '||p||∞,K ≤ max |dᵢⱼₖ|', '', 19)
    add_triangle(slide, 7.20, 1.20, 4.25, 3.10, 'control points', grid=True)
    add_text(slide, 6.95, 4.70, 5.20, 1.25, '直觀：函數值可視為控制點係數的加權平均。\n所以「找函數最大值」可改成「檢查控制點係數」。', 18.2, NAVY, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_footer(slide, 8)
    notes.append('Bernstein 多項式的關鍵在於基底非負且總和為一，因此多項式值落在控制點係數形成的範圍內。這使 L∞ 最大值限制能被控制點最大值界住。')

    # 9 optimization
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '演算法核心：由 λ 的下界推得 Cᴸ 的上界', 'Algorithm', 9)
    add_flow(slide, [('建立 VhFM', '基底 φᵢ'), ('組裝 A,B', '能量與控制點'), ('放寬限制', '||Bx||∞ ≥ 1'), ('計算 λh,B', '矩陣公式'), ('轉回 Cᴸ', '1/√λ')], x0=0.75, y=1.45, box_w=2.20, box_h=1.00, gap=0.20)
    add_formula(slide, 1.05, 3.30, 5.25, 0.85, 'λh,B = min xᵀAx,  s.t. ||Bx||∞ ≥ 1', '', 17.5)
    add_formula(slide, 6.95, 3.30, 5.20, 0.85, 'D = BA⁻¹Bᵀ,   λh,B = 1 / max diag(D)', '', 17.5)
    add_text(slide, 1.05, 4.80, 11.10, 0.90, '方向要抓清楚：λ 的下界越大，Cᴸ = 1/√λ 的上界越小，誤差保證就越銳利。', 20, NAVY, True, RGBColor(237, 245, 252), PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_footer(slide, 9)
    notes.append('演算法步驟：先建立 Fujino–Morley 空間與矩陣 A，再用 B 表示 Bernstein 控制點。透過放寬的 L∞ 限制計算 λh,B，最後由 Cᴸ=1/√λ 得到誤差常數上界。')

    # 10 calculation
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '示範計算：等腰直角三角形的上界', 'Worked example', 10)
    add_text(slide, 0.90, 1.30, 4.15, 3.25, '已知條件\n\n三角形：K₁,π/2\n單位等腰直角三角形\n\nTable 2 給出：λh,B = 5.7812\n\n關係式：Cᴸ = 1 / √λ', 18, DARK, False)
    add_formula(slide, 5.42, 1.30, 6.55, 0.78, 'Cᴸ ≤ 1 / √5.7812', '', 19)
    add_formula(slide, 5.42, 2.30, 6.55, 0.78, '√5.7812 ≈ 2.4044', '', 19)
    add_formula(slide, 5.42, 3.30, 6.55, 0.78, '1 / 2.4044 ≈ 0.41596', '', 19)
    add_label(slide, 1.20, 5.20, 3.15, 0.72, '下界\n0.40432', GREEN, 17)
    add_label(slide, 5.10, 5.20, 3.15, 0.72, '真值區間\n0.40432–0.41596', BLUE, 17)
    add_label(slide, 9.00, 5.20, 3.15, 0.72, '上界\n0.41596', ORANGE, 17)
    add_footer(slide, 10)
    notes.append('示範計算：Table 2 對等腰直角三角形給出 λh,B=5.7812。因為 Cᴸ=1/√λ，所以上界是 1/√5.7812。根號約為 2.4044，倒數約為 0.41596。')

    # 11 numerical results
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '數值結果：不同角度三角形的上下界', 'Numerical results', 11)
    add_table(slide, 0.75, 1.25, 5.95, 3.58)
    add_chart(slide, 7.05, 1.15, 5.50, 3.88)
    add_text(slide, 0.95, 5.45, 11.30, 0.86, '觀察：θ=π/3 附近常數較小；θ=π/2 的常數明顯較大。三角形形狀會直接影響最大範數誤差保證。', 18.2, NAVY, True, WHITE, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_footer(slide, 11)
    notes.append('這張呈現 Table 2 的主要數字。上下界越接近，代表估計越銳利。π/3 附近常數較小，而 π/2 顯著較大，說明元素形狀對最大範數誤差有直接影響。')

    # 12 shape dependence
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '形狀影響：三角形退化時常數可能變大', 'Interpretation', 12)
    add_triangle(slide, 0.85, 1.30, 3.5, 2.8, 'regular', grid=True)
    add_triangle(slide, 4.90, 1.30, 3.5, 2.8, 'right', grid=True)
    # Degenerate-looking thin triangle
    ax, ay, aw, ah = inch(9.0), inch(1.35), inch(3.1), inch(2.55)
    p1 = (ax + aw*0.05, ay + ah*0.82); p2 = (ax + aw*0.94, ay + ah*0.82); p3 = (ax + aw*0.18, ay + ah*0.66)
    for a, b in [(p1,p2),(p2,p3),(p3,p1)]:
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, a[0], a[1], b[0], b[1]); ln.line.color.rgb = RED; ln.line.width = Pt(2.2)
    lab = slide.shapes.add_textbox(inch(9.60), inch(3.70), inch(2.2), inch(0.28)); apply_text(lab.text_frame.paragraphs[0], 'degenerate', 15, RED, True, FONT_EN, PP_ALIGN.CENTER)
    add_heading_body(slide, 0.90, 4.45, 3.55, 1.45, '形狀正則', ['最小角不太小，常數有界'], BLUE, 16)
    add_heading_body(slide, 4.88, 4.45, 3.55, 1.45, '直角案例', ['可取得銳利上下界'], TEAL, 16)
    add_heading_body(slide, 8.86, 4.45, 3.55, 1.45, '退化案例', ['角度或邊長退化時，Cᴸ 可能趨近無窮'], RED, 16)
    add_footer(slide, 12)
    notes.append('這張解釋形狀的重要性。論文指出對形狀正則的三角形族，誤差常數是有界的；但若三角形逐漸退化成線段，常數可能趨近無窮。')

    # 13 result meaning
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '結果解讀：這些數字在 FEM 中代表什麼？', 'Meaning', 13)
    add_heading_body(slide, 0.85, 1.30, 3.75, 4.45, '對數值分析', ['提供最大範數下的明確誤差常數', '可支援逐點誤差估計', '比只知道收斂階數更具體'], BLUE, 16.5)
    add_heading_body(slide, 4.83, 1.30, 3.75, 4.45, '對網格設計', ['三角形形狀會影響常數', '形狀正則性仍是核心條件', '角度與長寬比會影響最壞誤差'], TEAL, 16.5)
    add_heading_body(slide, 8.81, 1.30, 3.75, 4.45, '對方法論', ['FEM + Bernstein 是可移植策略', '適合處理最大範數限制', '可延伸到高階插值或 PDE 問題'], ORANGE, 16.5)
    add_footer(slide, 13)
    notes.append('這張把數值結果轉成 FEM 應用意義。若要嚴格控制有限元素解的逐點誤差，只知道收斂階數不夠，還需要明確的常數。')

    # 14 conclusion
    slide = prs.slides.add_slide(blank); add_bg(slide); add_title(slide, '結論：三句話總結這篇論文', 'Conclusion', 14)
    add_text(slide, 1.05, 1.35, 11.20, 0.86, '1｜研究問題：估計三角形上線性 Lagrange 插值的 L∞ 誤差常數 Cᴸ(K)。', 19.2, NAVY, True, WHITE, None, MSO_ANCHOR.MIDDLE)
    add_text(slide, 1.05, 2.58, 11.20, 0.86, '2｜技術核心：用 Fujino–Morley FEM 離散化，再用 Bernstein 凸包性質處理最大範數約束。', 19.2, NAVY, True, WHITE, None, MSO_ANCHOR.MIDDLE)
    add_text(slide, 1.05, 3.81, 11.20, 0.86, '3｜代表結果：單位等腰直角三角形滿足 0.40432 ≤ Cᴸ(K) ≤ 0.41596。', 19.2, NAVY, True, WHITE, None, MSO_ANCHOR.MIDDLE)
    add_text(slide, 1.05, 5.35, 11.20, 0.72, 'Takeaway：這是一個把「最大誤差保證」做成可計算、可驗證數值界的研究。', 19.2, BLUE, True, RGBColor(237,245,252), PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    add_footer(slide, 14)
    notes.append('總結三句話：問題是 Cᴸ(K)，方法是 Fujino–Morley FEM 加 Bernstein 凸包，結果是等腰直角三角形的嚴格上下界。')

    # 15 Q&A
    slide = prs.slides.add_slide(blank); add_bg(slide)
    tb = slide.shapes.add_textbox(inch(0.82), inch(1.15), inch(11.2), inch(0.95))
    apply_text(tb.text_frame.paragraphs[0], '謝謝聆聽', 42, NAVY, True)
    sub = slide.shapes.add_textbox(inch(0.86), inch(2.28), inch(10.6), inch(0.55))
    apply_text(sub.text_frame.paragraphs[0], 'Q&A：λ 與 Cᴸ 的關係、Bernstein 放寬的銳利度、FEM 誤差估計應用。', 19, GRAY)
    add_label(slide, 1.05, 4.05, 3.25, 0.78, '關鍵公式\nCᴸ = 1/√λ', BLUE, 17)
    add_label(slide, 4.95, 4.05, 3.25, 0.78, '關鍵工具\nBernstein 凸包', TEAL, 17)
    add_triangle(slide, 9.00, 3.30, 3.15, 2.45, 'Cᴸ(K)')
    add_footer(slide, 15)
    notes.append('結尾提醒可討論的問題：λ 與 Cᴸ 的倒平方根關係，Bernstein 放寬是否銳利，以及這個常數如何進一步用在 FEM 的最大範數誤差估計。')

    return prs, notes


def main():
    prs, notes = build_deck()
    prs.save(PPTX_PATH)
    NOTES_PATH.write_text('\n\n'.join(f'Slide {i+1}\n{t}' for i, t in enumerate(notes)), encoding='utf-8')
    check = Presentation(PPTX_PATH)
    oob = warnIfSlideElementsOutOfBounds(check)
    overlaps = warnIfSlideHasOverlaps(check)
    if oob:
        print('OUT_OF_BOUNDS_WARNINGS')
        print('\n'.join(oob))
    if overlaps:
        print('OVERLAP_WARNINGS')
        print('\n'.join(overlaps[:40]))
    print(f'Created: {PPTX_PATH}')
    print(f'Created: {NOTES_PATH}')


if __name__ == '__main__':
    main()
